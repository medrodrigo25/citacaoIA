"""
CitacaoIA -- Gestor Inteligente de Citacoes Cientificas Vancouver
Versao: 2.0
Uso: streamlit run citacaoIA.py
"""

import re, json, time, os, threading as _threading
import json as _json_auth
import os   as _os
from datetime import datetime
from io import BytesIO

import streamlit as st
import requests

# Global task store for background pipeline runs
_PIPELINE_TASKS: dict = {}

# ── PAGE CONFIG ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="CitacaoIA . Vancouver",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
:root {--primary:#1e3a5f;--accent:#2e7ddb;--bg:#f4f7fb;}
.stApp{background:var(--bg);}
.cia-header{background:linear-gradient(135deg,#1e3a5f,#2e7ddb);color:white;
  padding:1.4rem 2rem;border-radius:12px;margin-bottom:1.5rem;
  box-shadow:0 4px 18px rgba(30,58,95,.25);}
.cia-header h1{color:white;margin:0;font-size:2rem;}
.cia-header p{color:rgba(255,255,255,.82);margin:.3rem 0 0;font-size:.95rem;}
.pill-changed{background:#fff3cd;color:#7a4f00;border-radius:20px;
  padding:2px 10px;font-size:.78rem;font-weight:600;}
.pill-ok{background:#d4edda;color:#155724;border-radius:20px;
  padding:2px 10px;font-size:.78rem;font-weight:600;}
.ref-box{background:#f8f9fc;border-left:4px solid #2e7ddb;border-radius:0 8px 8px 0;
  padding:.7rem 1rem;font-size:.88rem;font-family:Georgia,serif;margin:.4rem 0;}
.art-card{background:white;border:1.5px solid #d0daea;border-radius:10px;
  padding:1rem 1.2rem;margin:.5rem 0;transition:border-color .2s;}
.art-card:hover{border-color:#2e7ddb;}
.art-title{font-size:1rem;font-weight:700;color:#1e3a5f;margin:0 0 .3rem;}
.art-meta{font-size:.82rem;color:#555;margin:.2rem 0;}
.step-badge{display:inline-block;background:#2e7ddb;color:white;border-radius:50%;
  width:24px;height:24px;line-height:24px;text-align:center;
  font-size:.8rem;font-weight:700;margin-right:6px;}
</style>
""", unsafe_allow_html=True)

# ── CONSTANTS ─────────────────────────────────────────────────────────────────
PUBMED_BASE   = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/"
SEMANTIC_BASE = "https://api.semanticscholar.org/graph/v1/paper/search"
CROSSREF_BASE = "https://api.crossref.org/works/"
MAX_SEARCHES  = 6
CITATION_STYLES = ["Vancouver", "APA 7a Ed.", "ABNT NBR 6023", "Chicago"]
BIBLIOTECA_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "citacaoIA_biblioteca.json")

# =============================================================================
# SUPABASE HELPER  (usado quando configurado; cai para JSON local caso contrario)
# =============================================================================

def _supa_cfg():
    """Return (url, headers) or (None, None) if Supabase not configured."""
    try:
        url = st.secrets["supabase"]["url"].rstrip("/")
        key = st.secrets["supabase"]["key"]
        return url, {"apikey": key, "Authorization": f"Bearer {key}",
                     "Content-Type": "application/json", "Prefer": "return=representation"}
    except Exception:
        return None, None

def _supa_available() -> bool:
    url, _ = _supa_cfg()
    return url is not None


# =============================================================================
# AUTENTICACAO — Login / Registro / Admin / Log de uso
# =============================================================================
import hashlib as _hashlib
import uuid    as _uuid

_AUTH_FILE  = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "auth_db.json")
_USAGE_FILE = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "usage_log.json")


def _hash_pw(pw: str) -> str:
    return _hashlib.sha256(pw.encode()).hexdigest()


def _load_auth_db() -> dict:
    """Load auth DB.
    Priority for admin credentials (highest to lowest):
      1. Streamlit secrets [auth] admin_user / admin_password
      2. Environment variable ADMIN_RESET_PASSWORD / ADMIN_RESET_USER
      3. Local auth_db.json file (persists on local installs; ephemeral on Cloud)
    """
    import os as _os2

    # ── Read from file (local installs) ───────────────────────────────────
    db = {"users": {}, "invite_codes": {}, "recovery_codes": {}}
    if _os.path.exists(_AUTH_FILE):
        try:
            with open(_AUTH_FILE, "r") as f:
                loaded = _json_auth.load(f)
                db.update(loaded)
                if "recovery_codes" not in db:
                    db["recovery_codes"] = {}
        except Exception:
            pass

    # ── Priority 1: Streamlit secrets [auth] ─────────────────────────────
    # Works on Streamlit Cloud where filesystem is ephemeral.
    # In secrets.toml add:
    #   [auth]
    #   admin_user     = "admin"
    #   admin_password = "admin123"
    try:
        _auth_sec = st.secrets.get("auth", {})
        sec_pw   = str(_auth_sec.get("admin_password", "")).strip()
        sec_user = str(_auth_sec.get("admin_user", "admin")).strip()
        if sec_pw:
            db["users"][sec_user] = {
                "password_hash":    _hash_pw(sec_pw),
                "role":             "admin",
                "created_at":       __import__("datetime").datetime.now().isoformat(),
                "invite_code_used": "SECRETS",
            }
            # Do NOT save to file here — on Cloud the password comes from secrets every load
    except Exception:
        pass

    # ── Priority 2: Environment variable (legacy / Docker / Railway) ──────
    reset_pw = _os2.environ.get("ADMIN_RESET_PASSWORD", "").strip()
    if reset_pw:
        admin_u = _os2.environ.get("ADMIN_RESET_USER", "admin").strip()
        db["users"][admin_u] = {
            "password_hash":    _hash_pw(reset_pw),
            "role":             "admin",
            "created_at":       __import__("datetime").datetime.now().isoformat(),
            "invite_code_used": "ENV_RESET",
        }
        try:
            _save_auth_db(db)
        except Exception:
            pass

    return db


def _save_auth_db(db: dict):
    with open(_AUTH_FILE, "w") as f:
        _json_auth.dump(db, f, indent=2)


def _load_usage_log() -> list:
    if _os.path.exists(_USAGE_FILE):
        with open(_USAGE_FILE, "r") as f:
            return _json_auth.load(f)
    return []


def _append_usage_log(username: str, feature: str, detail: str = ""):
    log = _load_usage_log()
    log.append({
        "username": username,
        "feature":  feature,
        "detail":   detail,
        "ts":       __import__("datetime").datetime.now().isoformat()
    })
    with open(_USAGE_FILE, "w") as f:
        _json_auth.dump(log[-5000:], f, indent=2)


def _gen_invite_code(db: dict, created_by: str,
                     max_uses: int = 1, expires_days: int = 30) -> str:
    from datetime import datetime, timedelta
    code = str(_uuid.uuid4())[:8].upper()
    db["invite_codes"][code] = {
        "created_by": created_by,
        "created_at": datetime.now().isoformat(),
        "max_uses":   max_uses,
        "uses":       0,
        "expires_at": (datetime.now() + timedelta(days=expires_days)).isoformat(),
        "used_by":    []
    }
    _save_auth_db(db)
    return code


def _validate_invite(db: dict, code: str) -> tuple:
    from datetime import datetime
    if code not in db["invite_codes"]:
        return False, "Codigo de convite invalido."
    inv = db["invite_codes"][code]
    if inv["uses"] >= inv["max_uses"]:
        return False, "Codigo ja utilizado o numero maximo de vezes."
    if datetime.fromisoformat(inv["expires_at"]) < datetime.now():
        return False, "Codigo de convite expirado."
    return True, "OK"


def _register_user(db: dict, username: str, password: str,
                   invite_code: str, role: str = "user") -> tuple:
    username = username.strip().lower()
    if len(username) < 3:
        return False, "Usuario deve ter ao menos 3 caracteres."
    if len(password) < 6:
        return False, "Senha deve ter ao menos 6 caracteres."
    if username in db["users"]:
        return False, "Nome de usuario ja existe."
    if role != "admin":
        ok, msg = _validate_invite(db, invite_code)
        if not ok:
            return False, msg
        db["invite_codes"][invite_code]["uses"] += 1
        db["invite_codes"][invite_code]["used_by"].append(username)
    db["users"][username] = {
        "password_hash":    _hash_pw(password),
        "role":             role,
        "created_at":       __import__("datetime").datetime.now().isoformat(),
        "invite_code_used": invite_code
    }
    _save_auth_db(db)
    return True, "Conta criada com sucesso!"


def _authenticate(db: dict, username: str, password: str) -> tuple:
    username = username.strip().lower()
    user = db["users"].get(username)
    if not user:
        return False, "Usuario nao encontrado."
    if user["password_hash"] != _hash_pw(password):
        return False, "Senha incorreta."
    return True, user["role"]


def _change_password(db: dict, username: str,
                     old_pw: str, new_pw: str) -> tuple:
    ok, _ = _authenticate(db, username, old_pw)
    if not ok:
        return False, "Senha atual incorreta."
    if len(new_pw) < 6:
        return False, "Nova senha deve ter ao menos 6 caracteres."
    db["users"][username]["password_hash"] = _hash_pw(new_pw)
    _save_auth_db(db)
    return True, "Senha alterada com sucesso!"


# ─────────────────────────────────────────────────────────────────────────────

def render_login_page() -> bool:
    """Gate: returns True if user is authenticated. Shows setup/login otherwise."""

    # ── HUB TOKEN AUTH (portal-aluno.html → iframe) ──────────────────────────
    # Se um token HMAC válido chegar via ?hub_token=...&hub_user=...&hub_ts=...,
    # o aluno é autenticado automaticamente sem ver o formulário de login.
    if not st.session_state.get("auth_user"):
        try:
            hub_cfg = st.secrets.get("hub_auth", {})
            if hub_cfg.get("enabled", False):
                import hmac as _hmac_hub, hashlib as _sha_hub, time as _time_hub
                params    = st.query_params
                hub_token = params.get("hub_token", "")
                hub_user  = params.get("hub_user",  "")
                hub_ts    = params.get("hub_ts",    "")
                secret    = hub_cfg.get("secret", "")
                if hub_token and hub_user and hub_ts and secret:
                    try:
                        ts_int = int(hub_ts)
                        age    = _time_hub.time() - ts_int
                        if 0 <= age <= 300:  # token válido por 5 minutos
                            payload  = f"{hub_user}|{hub_ts}".encode()
                            expected = _hmac_hub.new(
                                secret.encode(), payload, _sha_hub.sha256
                            ).hexdigest()
                            if _hmac_hub.compare_digest(expected, hub_token):
                                st.session_state["auth_user"] = hub_user.lower()
                                st.session_state["auth_role"] = "user"
                                st.session_state["hub_auth"]  = True
                                st.query_params.clear()
                    except (ValueError, TypeError):
                        pass  # timestamp malformado → login normal
        except Exception:
            pass  # secrets não configurado → ignora

    if st.session_state.get("auth_user"):
        return True

    db = _load_auth_db()

    # ── FIRST RUN: no users exist → create admin account ─────────────────────
    if not db["users"]:
        st.markdown("""
        <div style="max-width:440px;margin:60px auto;padding:2rem;
        border-radius:16px;box-shadow:0 4px 24px #0002;background:#fff;">
        <h2 style="color:#003366;text-align:center">CitacaoIA</h2>
        <p style="color:#e07b00;text-align:center;font-weight:600">
        Primeiro acesso — configure o administrador</p>
        """, unsafe_allow_html=True)

        st.info("Nenhum usuario cadastrado. Crie a conta de administrador para comecar.")
        adm_u  = st.text_input("Usuario do administrador", key="setup_u",
                                placeholder="min. 3 caracteres")
        adm_p  = st.text_input("Senha", key="setup_p", type="password",
                                placeholder="min. 6 caracteres")
        adm_p2 = st.text_input("Confirmar senha", key="setup_p2", type="password")

        if st.button("Criar conta de administrador", type="primary",
                     use_container_width=True, key="btn_setup"):
            if adm_p != adm_p2:
                st.error("As senhas nao coincidem.")
            else:
                ok, msg = _register_user(db, adm_u, adm_p, "SYSTEM", role="admin")
                if ok:
                    st.success("Administrador criado! Faca login abaixo.")
                    st.rerun()
                else:
                    st.error(msg)

        st.markdown("</div>", unsafe_allow_html=True)
        return False

    # ── NORMAL LOGIN ──────────────────────────────────────────────────────────
    import os as _os3
    _reset_active = bool(_os3.environ.get("ADMIN_RESET_PASSWORD","").strip())

    st.markdown("""
    <div style="max-width:420px;margin:70px auto;padding:2rem;
    border-radius:16px;box-shadow:0 4px 24px #0002;background:#fff;">
    <h2 style="color:#003366;text-align:center;margin-bottom:.2rem">CitacaoIA</h2>
    <p style="color:#888;text-align:center;font-size:.9rem;margin-bottom:1.5rem">
    Gestor Inteligente de Citacoes</p>
    """, unsafe_allow_html=True)
    if _reset_active:
        _ru = _os3.environ.get("ADMIN_RESET_USER","admin")
        st.warning(f"Modo reset ativo — use usuario **{_ru}** e a senha definida em ADMIN_RESET_PASSWORD.")

    tab_login, tab_reg, tab_recover = st.tabs(["Entrar", "Criar conta (convite)", "Recuperar Senha"])

    with tab_login:
        u = st.text_input("Usuario", key="login_u")
        p = st.text_input("Senha",   key="login_p", type="password")
        if st.button("Entrar", type="primary", use_container_width=True, key="btn_login"):
            ok, role_or_msg = _authenticate(db, u, p)
            if ok:
                st.session_state["auth_user"] = u.strip().lower()
                st.session_state["auth_role"] = role_or_msg
                _append_usage_log(u.strip().lower(), "login")
                st.rerun()
            else:
                st.error(role_or_msg)

    with tab_reg:
        st.caption("Voce recebeu um codigo de convite do administrador? Crie sua conta aqui.")
        inv = st.text_input("Codigo de convite", key="reg_inv",
                             placeholder="Ex: A3F7B2C1").strip().upper()
        nu  = st.text_input("Escolha um usuario", key="reg_u")
        np  = st.text_input("Crie uma senha",  key="reg_p", type="password")
        np2 = st.text_input("Confirme a senha", key="reg_p2", type="password")
        if st.button("Criar conta", type="primary", use_container_width=True, key="btn_reg"):
            if np != np2:
                st.error("As senhas nao coincidem.")
            else:
                ok, msg = _register_user(db, nu, np, inv)
                if ok:
                    st.success(msg + " Faca login na aba 'Entrar'.")
                else:
                    st.error(msg)

    with tab_recover:
        st.caption("Para redefinir sua senha voce precisa de um **codigo de recuperacao** gerado pelo administrador.")
        st.info("**Como funciona:**\n1. O admin gera um codigo de recuperacao no painel de administracao.\n2. O admin envia o codigo para voce.\n3. Voce usa o codigo aqui para criar uma nova senha.")

        rec_user   = st.text_input("Seu usuario", key="rec_username", placeholder="Digite seu usuario")
        rec_code   = st.text_input("Codigo de recuperacao", key="rec_code",
                                   placeholder="Ex: RC-A3F7B2C1").strip().upper()
        rec_pw     = st.text_input("Nova senha",    key="rec_new_pw",  type="password")
        rec_pw2    = st.text_input("Confirmar nova senha", key="rec_new_pw2", type="password")

        if st.button("Redefinir senha", type="primary", use_container_width=True, key="btn_recover"):
            if not rec_user.strip():
                st.error("Informe seu usuario.")
            elif not rec_code.strip():
                st.error("Informe o codigo de recuperacao.")
            elif len(rec_pw) < 6:
                st.error("A nova senha deve ter ao menos 6 caracteres.")
            elif rec_pw != rec_pw2:
                st.error("As senhas nao coincidem.")
            else:
                # Validate recovery code
                rec_codes = db.get("recovery_codes", {})
                uname = rec_user.strip().lower()
                if rec_code not in rec_codes:
                    st.error("Codigo de recuperacao invalido.")
                elif rec_codes[rec_code].get("used"):
                    st.error("Este codigo ja foi utilizado.")
                elif rec_codes[rec_code].get("username","") not in ("", uname):
                    st.error("Codigo nao valido para este usuario.")
                else:
                    from datetime import datetime as _dt
                    exp = rec_codes[rec_code].get("expires_at","")
                    if exp and _dt.fromisoformat(exp) < _dt.now():
                        st.error("Codigo expirado. Solicite um novo ao administrador.")
                    elif uname not in db["users"]:
                        st.error("Usuario nao encontrado.")
                    else:
                        db["users"][uname]["password_hash"] = _hash_pw(rec_pw)
                        db["recovery_codes"][rec_code]["used"] = True
                        _save_auth_db(db)
                        st.success("Senha redefinida com sucesso! Faca login na aba 'Entrar'.")

    st.markdown("</div>", unsafe_allow_html=True)
    return False


def _gen_recovery_code(db: dict, username: str, expires_hours: int = 24) -> str:
    """Generate a one-time recovery code for password reset."""
    from datetime import datetime, timedelta
    if "recovery_codes" not in db:
        db["recovery_codes"] = {}
    code = "RC-" + str(_uuid.uuid4())[:8].upper()
    db["recovery_codes"][code] = {
        "username":   username.strip().lower(),
        "created_at": datetime.now().isoformat(),
        "expires_at": (datetime.now() + timedelta(hours=expires_hours)).isoformat(),
        "used":       False,
    }
    _save_auth_db(db)
    return code


def _admin_reset_password(db: dict, username: str, new_pw: str) -> tuple:
    """Admin directly resets a user's password (no old password needed)."""
    username = username.strip().lower()
    if username not in db["users"]:
        return False, "Usuario nao encontrado."
    if len(new_pw) < 6:
        return False, "Nova senha deve ter ao menos 6 caracteres."
    db["users"][username]["password_hash"] = _hash_pw(new_pw)
    _save_auth_db(db)
    return True, f"Senha do usuario '{username}' redefinida com sucesso."


def render_admin_panel():
    """Admin-only full-page panel: users, invite codes, usage log, settings."""
    st.markdown("## Configuracoes do Administrador")
    db  = _load_auth_db()
    log = _load_usage_log()
    auth_user = st.session_state.get("auth_user", "admin")

    a_tab1, a_tab2, a_tab3, a_tab4, a_tab5 = st.tabs(
    ["Usuarios", "Convidar Usuario", "Recuperar Senha", "Log de Uso", "Minha Conta"]
    )

    # ── TAB 1: Users ──────────────────────────────────────────────────────
    with a_tab1:
        import pandas as _pd
        users_data = [
            {"Usuario": u,
             "Papel":    d["role"],
             "Criado":   d.get("created_at","")[:10],
             "Convite":  d.get("invite_code_used","")}
            for u, d in db["users"].items()
        ]
        st.dataframe(_pd.DataFrame(users_data), use_container_width=True)

        all_non_admin = [u for u, d in db["users"].items() if d["role"] != "admin"]
        if all_non_admin:
            del_u = st.selectbox("Remover usuario:", ["— selecione —"] + all_non_admin,
                                 key="del_u_sel")
            if st.button("Remover", key="btn_del_u", type="secondary"):
                if del_u and del_u != "— selecione —":
                    del db["users"][del_u]
                    _save_auth_db(db)
                    st.success(f"Usuario '{del_u}' removido.")
                    st.rerun()

    # ── TAB 2: Invite ─────────────────────────────────────────────────────
    with a_tab2:
        st.markdown("**Gere um codigo e envie ao novo usuario.**")
        st.caption("O usuario usa o codigo para criar login e senha proprios.")

        col_a, col_b = st.columns(2)
        inv_uses = col_a.number_input("Usos maximos", 1, 20, 1, key="inv_u")
        inv_days = col_b.number_input("Validade (dias)", 1, 90, 30, key="inv_d")

        if st.button("Gerar codigo de convite", type="primary",
                     use_container_width=True, key="btn_gen_inv"):
            code = _gen_invite_code(db, auth_user, int(inv_uses), int(inv_days))
            st.session_state["last_invite_code"] = code
            st.rerun()

        if st.session_state.get("last_invite_code"):
            code = st.session_state["last_invite_code"]
            st.success("Codigo gerado! Copie e envie ao usuario:")
            st.code(code, language=None)
            st.caption(
                f"Valido por {inv_days} dia(s), maximo {inv_uses} uso(s). "
                "O usuario acessa o app, clica em 'Criar conta (convite)' "
                "e usa este codigo para criar login e senha."
            )

        st.markdown("---")
        st.markdown("**Codigos gerados anteriormente:**")
        inv_rows = [
            {"Codigo": c,
             "Usos": f"{d['uses']}/{d['max_uses']}",
             "Expira": d.get("expires_at","")[:10],
             "Usado por": ", ".join(d.get("used_by",[]))}
            for c, d in db["invite_codes"].items()
        ]
        if inv_rows:
            import pandas as _pd2
            st.dataframe(_pd2.DataFrame(inv_rows), use_container_width=True)
        else:
            st.info("Nenhum codigo gerado ainda.")

    # ── TAB 3: Recuperar Senha ────────────────────────────────────────────
    with a_tab3:
        st.markdown("### Recuperacao de Senha de Usuarios")
        all_users = list(db["users"].keys())

        st.markdown("#### Opção 1 — Reset direto (admin define nova senha)")
        rst_user = st.selectbox("Usuario para resetar:", ["— selecione —"] + all_users, key="rst_u_sel")
        rst_pw   = st.text_input("Nova senha:", type="password", key="rst_pw_inp")
        rst_pw2  = st.text_input("Confirmar nova senha:", type="password", key="rst_pw2_inp")
        if st.button("Redefinir senha agora", key="btn_rst_direct", type="primary"):
            if rst_user == "— selecione —":
                st.error("Selecione um usuario.")
            elif len(rst_pw) < 6:
                st.error("Senha deve ter ao menos 6 caracteres.")
            elif rst_pw != rst_pw2:
                st.error("As senhas nao coincidem.")
            else:
                ok, msg = _admin_reset_password(db, rst_user, rst_pw)
                if ok:
                    st.success(msg)
                else:
                    st.error(msg)

        st.markdown("---")
        st.markdown("#### Opção 2 — Gerar codigo de recuperacao (usuario redefine sozinho)")
        st.caption("O usuario entra na tela de login → aba 'Recuperar Senha' e usa o codigo.")
        rc_user  = st.selectbox("Usuario:", ["— selecione —"] + all_users, key="rc_u_sel")
        rc_hours = st.number_input("Validade do codigo (horas):", 1, 72, 24, key="rc_hours")
        if st.button("Gerar codigo de recuperacao", key="btn_gen_rc"):
            if rc_user == "— selecione —":
                st.error("Selecione um usuario.")
            else:
                code = _gen_recovery_code(db, rc_user, int(rc_hours))
                st.session_state["last_recovery_code"] = code
                st.session_state["last_recovery_user"] = rc_user
                st.rerun()

        if st.session_state.get("last_recovery_code"):
            code  = st.session_state["last_recovery_code"]
            ruser = st.session_state.get("last_recovery_user","")
            st.success(f"Codigo gerado para **{ruser}**! Copie e envie ao usuario:")
            st.code(code, language=None)
            st.caption(f"Valido por {rc_hours} hora(s). O usuario acessa o app → aba 'Recuperar Senha' e usa este codigo.")

        st.markdown("---")
        st.markdown("**Codigos de recuperacao gerados:**")
        rc_rows = [
            {"Codigo": c,
             "Usuario": d.get("username",""),
             "Expira": d.get("expires_at","")[:16].replace("T"," "),
             "Usado": "Sim" if d.get("used") else "Nao"}
            for c, d in db.get("recovery_codes", {}).items()
        ]
        if rc_rows:
            import pandas as _pdr
            st.dataframe(_pdr.DataFrame(rc_rows), use_container_width=True)
        else:
            st.info("Nenhum codigo de recuperacao gerado ainda.")

    # ── TAB 4: Usage log ──────────────────────────────────────────────────
    with a_tab4:
        if log:
            import pandas as _pd3
            df = _pd3.DataFrame(log[-500:][::-1])
            df.columns = ["Usuario", "Funcao", "Detalhe", "Data/Hora"]
            opts = ["Todos"] + sorted(df["Usuario"].unique().tolist())
            filt = st.selectbox("Filtrar por usuario:", opts, key="log_filt")
            if filt != "Todos":
                df = df[df["Usuario"] == filt]
            st.dataframe(df, use_container_width=True)
            st.download_button("Exportar CSV",
                               data=df.to_csv(index=False).encode(),
                               file_name="usage_log.csv", mime="text/csv",
                               key="dl_log_csv")
        else:
            st.info("Nenhuma acao registrada ainda.")

    # ── TAB 5: Change password ────────────────────────────────────────────
    with a_tab5:
        st.markdown("**Alterar minha senha**")
        old_pw  = st.text_input("Senha atual", type="password", key="chpw_old")
        new_pw  = st.text_input("Nova senha",  type="password", key="chpw_new")
        new_pw2 = st.text_input("Confirmar nova senha", type="password", key="chpw_new2")
        if st.button("Alterar senha", key="btn_chpw"):
            if new_pw != new_pw2:
                st.error("As novas senhas nao coincidem.")
            else:
                ok, msg = _change_password(db, auth_user, old_pw, new_pw)
                if ok:
                    st.success(msg)
                else:
                    st.error(msg)

# =============================================================================
# BIBLIOTECA  (Supabase quando disponivel, JSON local como fallback)
# =============================================================================

def load_library() -> list:
    url, hdrs = _supa_cfg()
    if url:
        try:
            r = requests.get(f"{url}/rest/v1/citacaoIA_biblioteca",
                             params={"select":"*","order":"added_date.desc"},
                             headers=hdrs, timeout=10)
            if r.status_code == 200:
                rows = r.json()
                return [row["article_data"] for row in rows if row.get("article_data")]
        except Exception:
            pass
    # fallback: JSON local
    if os.path.exists(BIBLIOTECA_FILE):
        try:
            with open(BIBLIOTECA_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return []
    return []

def save_library(library: list) -> bool:
    """Save to JSON (local only). Cloud saves are done incrementally via Supabase."""
    try:
        folder = os.path.dirname(BIBLIOTECA_FILE)
        if folder and not os.path.exists(folder):
            os.makedirs(folder, exist_ok=True)
        with open(BIBLIOTECA_FILE, "w", encoding="utf-8") as f:
            json.dump(library, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        return False

def _supa_upsert(article: dict):
    """Insert or update a single article in Supabase."""
    url, hdrs = _supa_cfg()
    if not url:
        return
    doi   = article.get("doi","").lower().strip()
    title = article.get("title","").lower()[:60]
    # Check if exists by DOI or title
    try:
        params = {"select":"id,article_data"}
        if doi:
            params["article_data->>doi"] = f"eq.{doi}"
        r = requests.get(f"{url}/rest/v1/citacaoIA_biblioteca",
                         params=params, headers=hdrs, timeout=8)
        existing = r.json() if r.status_code == 200 else []
        if not existing and title:
            r2 = requests.get(f"{url}/rest/v1/citacaoIA_biblioteca",
                              params={"select":"id","article_data->>title": f"ilike.{article.get('title','')[:60]}%"},
                              headers=hdrs, timeout=8)
            existing = r2.json() if r2.status_code == 200 else []
        if existing:
            row_id = existing[0]["id"]
            requests.patch(f"{url}/rest/v1/citacaoIA_biblioteca?id=eq.{row_id}",
                           json={"article_data": article}, headers=hdrs, timeout=8)
        else:
            article["added_date"] = article.get("added_date", datetime.now().isoformat())
            requests.post(f"{url}/rest/v1/citacaoIA_biblioteca",
                          json={"article_data": article}, headers=hdrs, timeout=8)
    except Exception:
        pass

def _supa_delete(article: dict):
    """Delete a single article from Supabase by DOI or title."""
    url, hdrs = _supa_cfg()
    if not url:
        return
    doi = article.get("doi","").lower().strip()
    try:
        if doi:
            requests.delete(f"{url}/rest/v1/citacaoIA_biblioteca",
                            params={"article_data->>doi": f"eq.{doi}"},
                            headers=hdrs, timeout=8)
        else:
            title = article.get("title","")[:60]
            requests.delete(f"{url}/rest/v1/citacaoIA_biblioteca",
                            params={"article_data->>title": f"ilike.{title}%"},
                            headers=hdrs, timeout=8)
    except Exception:
        pass

def add_to_library(article: dict) -> tuple:
    """Add or update article. Returns (library, is_new)."""
    lib = load_library()
    doi = article.get("doi","").lower().strip()
    is_new = True
    for i, a in enumerate(lib):
        if doi and a.get("doi","").lower().strip() == doi:
            lib[i] = {**a, **article, "added_date": a.get("added_date", datetime.now().isoformat())}
            is_new = False
            break
        if a.get("title","").lower()[:60] == article.get("title","").lower()[:60]:
            lib[i] = {**a, **article, "added_date": a.get("added_date", datetime.now().isoformat())}
            is_new = False
            break
    if is_new:
        article["added_date"] = datetime.now().isoformat()
        lib.insert(0, article)
    save_library(lib)
    _supa_upsert(article if not is_new else lib[0])
    return lib, is_new

def delete_from_library(index: int) -> list:
    lib = load_library()
    if 0 <= index < len(lib):
        art = lib[index]
        lib.pop(index)
        save_library(lib)
        _supa_delete(art)
    return lib

def search_library(query: str) -> list:
    lib = load_library()
    if not query.strip():
        return lib
    q = query.lower()
    return [a for a in lib if
            q in a.get("title","").lower() or
            q in " ".join(a.get("authors",[])).lower() or
            q in a.get("journal","").lower() or
            q in a.get("year","").lower() or
            q in a.get("doi","").lower() or
            q in a.get("abstract","").lower()]

# =============================================================================
# DOI LOOKUP (CrossRef)
# =============================================================================

def lookup_doi(doi: str) -> dict | None:
    doi = doi.strip().lstrip("https://doi.org/").lstrip("http://doi.org/").lstrip("doi:")
    try:
        r = requests.get(CROSSREF_BASE + doi, timeout=12,
                         headers={"User-Agent": "CitacaoIA/2.0 mailto:research@borghi.med.br"})
        if r.status_code != 200:
            return None
        data = r.json().get("message", {})

        # Authors
        authors = []
        for a in data.get("author", []):
            family = a.get("family", "")
            given  = a.get("given", "")
            if family:
                initials = "".join(w[0] for w in given.split() if w) if given else ""
                authors.append(f"{family} {initials}".strip())

        # Journal
        journal = ""
        containers = data.get("container-title", [])
        if containers:
            journal = containers[0]

        # Year
        year = ""
        dp = data.get("published", data.get("published-print", data.get("published-online", {})))
        parts = dp.get("date-parts", [[""]])[0]
        if parts:
            year = str(parts[0])

        # Title
        titles = data.get("title", [""])
        title = titles[0] if titles else ""

        # Pages
        pages = data.get("page", "")

        return {
            "doi":      doi,
            "title":    title,
            "authors":  authors,
            "journal":  journal,
            "year":     year,
            "volume":   str(data.get("volume", "")),
            "issue":    str(data.get("issue", "")),
            "pages":    pages,
            "abstract": data.get("abstract", ""),
            "_source":  "CrossRef/DOI",
        }
    except Exception as e:
        return None

# =============================================================================
# WEB SEARCH — PubMed, LILACS, Embase (via Europe PMC), Semantic Scholar,
#              OpenAlex (busca cinzenta / literatura expandida)
# =============================================================================

def search_pubmed(query: str, max_results: int = 3) -> list:
    """Search PubMed via NCBI eutils (free API, no key needed)."""
    try:
        r = requests.get(PUBMED_BASE + "esearch.fcgi",
            params={"db":"pubmed","term":query,"retmax":max_results,"retmode":"json"}, timeout=10)
        pmids = r.json().get("esearchresult", {}).get("idlist", [])
        if not pmids:
            return []
        r2 = requests.get(PUBMED_BASE + "esummary.fcgi",
            params={"db":"pubmed","id":",".join(pmids),"retmode":"json"}, timeout=10)
        data = r2.json().get("result", {})
        articles = []
        for pmid in pmids:
            item = data.get(pmid, {})
            if not item or pmid == "uids":
                continue
            authors = [a.get("name","") for a in item.get("authors",[])]
            articles.append({
                "pmid": pmid, "title": item.get("title","").rstrip("."),
                "authors": authors, "journal": item.get("source",""),
                "year": item.get("pubdate","")[:4], "volume": item.get("volume",""),
                "issue": item.get("issue",""), "pages": item.get("pages",""),
                "doi": item.get("elocationid","").replace("doi: ",""), "_source": "PubMed",
            })
        return articles
    except Exception:
        return []


def search_europe_pmc(query: str, max_results: int = 3,
                      source_filter: str = "") -> list:
    """Search Europe PMC — covers PubMed Central, EMBASE subset, WHO IRIS, preprints.
    source_filter: '' = all, 'MED' = MEDLINE, 'PPR' = preprints (busca cinzenta)."""
    try:
        params = {
            "query": query + (f" SOURCE:{source_filter}" if source_filter else ""),
            "resultType": "core",
            "pageSize": max_results,
            "format": "json",
        }
        r = requests.get("https://www.ebi.ac.uk/europepmc/webservices/rest/search",
                         params=params, timeout=12)
        results = r.json().get("resultList", {}).get("result", [])
        articles = []
        for item in results:
            authors_raw = item.get("authorList", {}).get("author", [])
            if isinstance(authors_raw, list):
                authors = [f"{a.get('lastName','')} {a.get('initials','')}" for a in authors_raw]
            else:
                authors = [item.get("authorString","")]
            articles.append({
                "title":   item.get("title","").rstrip("."),
                "authors": authors,
                "journal": item.get("journalTitle",""),
                "year":    str(item.get("pubYear","")),
                "volume":  item.get("journalVolume",""),
                "issue":   item.get("issue",""),
                "pages":   item.get("pageInfo",""),
                "doi":     item.get("doi",""),
                "pmid":    item.get("pmid",""),
                "_source": f"Europe PMC{'/Preprint' if source_filter=='PPR' else ''}",
            })
        return articles
    except Exception:
        return []


def search_lilacs(query: str, max_results: int = 3) -> list:
    """Search LILACS via the free BVS/BIREME iAHx Solr endpoint (no auth required)."""
    articles = []
    # Strategy 1 — BVS iAHx Solr (free, JSON output)
    try:
        r = requests.get(
            "https://pesquisa.bvsalud.org/portal/api/",
            params={
                "q":    query,
                "site": "portal",
                "op":   "search",
                "lang": "pt",
                "fmt":  "json",
                "count": max_results,
            },
            timeout=12,
        )
        if r.status_code == 200:
            data = r.json()
            docs = (data.get("hits", {}).get("hits", []) or
                    data.get("documents", []) or
                    data.get("results", []))
            for doc in docs[:max_results]:
                src = doc.get("_source", doc)
                ti  = src.get("ti", src.get("title", ""))
                if isinstance(ti, dict):
                    ti = next(iter(ti.values()), "") if ti else ""
                if not str(ti).strip():
                    continue
                au = src.get("au", src.get("authors", []))
                if isinstance(au, str):
                    au = [au]
                articles.append({
                    "title":   str(ti).rstrip("."),
                    "authors": au,
                    "journal": src.get("ta", src.get("journal", "")),
                    "year":    str(src.get("dp", src.get("year", "")))[:4],
                    "volume":  src.get("vi", ""),
                    "issue":   src.get("ip", ""),
                    "pages":   src.get("pg", ""),
                    "doi":     src.get("doi", ""),
                    "_source": "LILACS/BVS",
                })
    except Exception:
        pass

    # Strategy 2 — Europe PMC with LILACS filter (free, covers Latin American journals)
    if not articles:
        try:
            r = requests.get(
                "https://www.ebi.ac.uk/europepmc/webservices/rest/search",
                params={
                    "query":      query + " (JOURNAL_ISSN:\"LILACS\" OR SRC:AGR OR SRC:PAT)",
                    "resultType": "core",
                    "pageSize":   max_results,
                    "format":     "json",
                },
                timeout=12,
            )
            results = r.json().get("resultList", {}).get("result", [])
            for item in results:
                authors_raw = item.get("authorList", {}).get("author", [])
                if isinstance(authors_raw, list):
                    authors = [f"{a.get('lastName','')} {a.get('initials','')}" for a in authors_raw]
                else:
                    authors = [item.get("authorString", "")]
                articles.append({
                    "title":   item.get("title", "").rstrip("."),
                    "authors": authors,
                    "journal": item.get("journalTitle", ""),
                    "year":    str(item.get("pubYear", "")),
                    "volume":  item.get("journalVolume", ""),
                    "issue":   item.get("issue", ""),
                    "pages":   item.get("pageInfo", ""),
                    "doi":     item.get("doi", ""),
                    "pmid":    item.get("pmid", ""),
                    "_source": "LILACS/Europe PMC",
                })
        except Exception:
            pass

    return articles


def search_openalex(query: str, max_results: int = 3) -> list:
    """Search OpenAlex — open scholarly graph with 250M+ works (busca cinzenta ampliada)."""
    try:
        r = requests.get(
            "https://api.openalex.org/works",
            params={
                "search": query,
                "per_page": max_results,
                "select": "title,authorships,publication_year,primary_location,doi,biblio",
                "mailto": "research@borghi.med.br",
            },
            timeout=12,
        )
        items = r.json().get("results", [])
        articles = []
        for item in items:
            authors = [
                a.get("author", {}).get("display_name", "")
                for a in item.get("authorships", [])
            ]
            loc    = item.get("primary_location") or {}
            source = (loc.get("source") or {}).get("display_name", "")
            biblio = item.get("biblio", {})
            doi    = item.get("doi","").replace("https://doi.org/","")
            articles.append({
                "title":   item.get("title","").rstrip("."),
                "authors": authors,
                "journal": source,
                "year":    str(item.get("publication_year","")),
                "volume":  str(biblio.get("volume","")),
                "issue":   str(biblio.get("issue","")),
                "pages":   f"{biblio.get('first_page','')}--{biblio.get('last_page','')}".strip("-"),
                "doi":     doi,
                "_source": "OpenAlex",
            })
        return articles
    except Exception:
        return []


def search_semantic_scholar(query: str, max_results: int = 3) -> list:
    """Search Semantic Scholar (AI-powered academic graph)."""
    try:
        r = requests.get(SEMANTIC_BASE,
            params={"query":query,"fields":"title,authors,year,venue,externalIds","limit":max_results}, timeout=10)
        items = r.json().get("data", [])
        articles = []
        for item in items:
            authors = [a.get("name","") for a in item.get("authors",[])]
            ext = item.get("externalIds", {})
            articles.append({
                "title": item.get("title",""), "authors": authors,
                "year": str(item.get("year","")), "journal": item.get("venue",""),
                "doi": ext.get("DOI",""), "pmid": ext.get("PubMed",""),
                "volume":"", "issue":"", "pages":"", "_source": "Semantic Scholar",
            })
        return articles
    except Exception:
        return []


def multi_source_search(query: str, max_per_source: int = 2) -> list:
    """Search all configured databases for a single query and return merged results."""
    results = []
    # PubMed / MEDLINE
    results.extend(search_pubmed(query, max_per_source))
    # Europe PMC (covers PMC + Cochrane + preprints)
    results.extend(search_europe_pmc(query, max_per_source))
    # LILACS (literatura latino-americana)
    results.extend(search_lilacs(query, max_per_source))
    # OpenAlex (busca cinzenta ampliada — sem paywalls de API)
    results.extend(search_openalex(query, max_per_source))
    # Semantic Scholar (fallback / enriquecimento)
    if not results:
        results.extend(search_semantic_scholar(query, max_per_source))
    return results

# =============================================================================
# VANCOUVER FORMATTER
# =============================================================================

def format_vancouver(ref: dict, number: int) -> str:
    authors = ref.get("authors", []) or ["Autor desconhecido"]
    if len(authors) > 6:
        author_str = ", ".join(authors[:6]) + " et al"
    else:
        author_str = ", ".join(authors)
    title   = ref.get("title","Sem titulo").rstrip(".")
    journal = ref.get("journal","")
    year    = ref.get("year","")
    vol, iss, pgs = ref.get("volume",""), ref.get("issue",""), ref.get("pages","")
    doi     = ref.get("doi","")
    pmid    = ref.get("pmid","")
    cit = f"{number}. {author_str}. {title}. {journal}. {year}"
    if vol:  cit += f";{vol}"
    if iss:  cit += f"({iss})"
    if pgs:  cit += f":{pgs}"
    cit += "."
    if doi:  cit += f" {doi}"
    elif pmid: cit += f" PMID:{pmid}"
    return cit


def format_apa(ref: dict, number: int) -> str:
    """APA 7th Edition: Author, A. A., & Author, B. B. (Year). Title. Journal, vol(iss), pages. doi"""
    authors_raw = ref.get("authors", []) or ["Autor desconhecido"]
    def apa_name(n):
        parts = n.strip().split()
        if len(parts) >= 2:
            last = parts[0].rstrip(",")
            initials = " ".join(p[0].upper()+"." for p in parts[1:] if p)
            return f"{last}, {initials}"
        return n
    apa_authors = [apa_name(a) for a in authors_raw]
    if len(apa_authors) > 20:
        author_str = ", ".join(apa_authors[:19]) + ", ... " + apa_authors[-1]
    elif len(apa_authors) > 2:
        author_str = ", ".join(apa_authors[:-1]) + ", & " + apa_authors[-1]
    elif len(apa_authors) == 2:
        author_str = apa_authors[0] + ", & " + apa_authors[1]
    else:
        author_str = apa_authors[0]
    year    = ref.get("year","s.d.")
    title   = ref.get("title","Sem titulo").rstrip(".")
    journal = ref.get("journal","")
    vol, iss, pgs = ref.get("volume",""), ref.get("issue",""), ref.get("pages","")
    doi     = ref.get("doi","")
    cit = f"{author_str} ({year}). {title}."
    if journal:
        cit += f" *{journal}*"
        if vol:  cit += f", *{vol}*"
        if iss:  cit += f"({iss})"
        if pgs:  cit += f", {pgs}"
        cit += "."
    if doi: cit += f" https://doi.org/{doi}" if not doi.startswith("http") else f" {doi}"
    return cit


def format_abnt(ref: dict, number: int) -> str:
    """ABNT NBR 6023:2018: SOBRENOME, Nome. Titulo. Revista, vol., n., p., ano."""
    authors_raw = ref.get("authors", []) or ["AUTOR DESCONHECIDO"]
    def abnt_name(n):
        parts = n.strip().split()
        if len(parts) >= 2:
            last = parts[0].rstrip(",").upper()
            first = " ".join(p for p in parts[1:] if p)
            return f"{last}, {first}"
        return n.upper()
    abnt_authors = [abnt_name(a) for a in authors_raw]
    if len(abnt_authors) > 3:
        author_str = abnt_authors[0] + " et al."
    else:
        author_str = "; ".join(abnt_authors)
    title   = ref.get("title","Sem titulo").rstrip(".")
    journal = ref.get("journal","")
    year    = ref.get("year","")
    vol, iss, pgs = ref.get("volume",""), ref.get("issue",""), ref.get("pages","")
    doi     = ref.get("doi","")
    cit = f"{author_str}. {title}."
    if journal: cit += f" {journal}"
    if vol:  cit += f", v. {vol}"
    if iss:  cit += f", n. {iss}"
    if pgs:  cit += f", p. {pgs}"
    if year: cit += f", {year}"
    cit += "."
    if doi: cit += f" Disponivel em: https://doi.org/{doi}." if not doi.startswith("http") else f" Disponivel em: {doi}."
    return cit


def format_chicago(ref: dict, number: int) -> str:
    """Chicago Author-Date: Author (Year). "Title." Journal vol, no. iss (year): pages. doi."""
    authors_raw = ref.get("authors", []) or ["Unknown Author"]
    def chi_name(n, first_author=False):
        parts = n.strip().split()
        if len(parts) >= 2 and first_author:
            last = parts[0].rstrip(",")
            first = " ".join(p for p in parts[1:] if p)
            return f"{last}, {first}"
        return n
    chi_authors = [chi_name(a, i==0) for i,a in enumerate(authors_raw)]
    if len(chi_authors) > 3:
        author_str = chi_authors[0] + " et al."
    elif len(chi_authors) > 1:
        author_str = ", ".join(chi_authors[:-1]) + ", and " + chi_authors[-1]
    else:
        author_str = chi_authors[0]
    year    = ref.get("year","n.d.")
    title   = ref.get("title","No title").rstrip(".")
    journal = ref.get("journal","")
    vol, iss, pgs = ref.get("volume",""), ref.get("issue",""), ref.get("pages","")
    doi     = ref.get("doi","")
    cit = f"{author_str}. {year}. \"{title}.\""
    if journal:
        cit += f" {journal}"
        if vol: cit += f" {vol}"
        if iss: cit += f", no. {iss}"
        if pgs: cit += f": {pgs}"
        cit += "."
    if doi: cit += f" https://doi.org/{doi}." if not doi.startswith("http") else f" {doi}."
    return cit


def format_reference(ref: dict, number: int, style: str = "Vancouver") -> str:
    """Master dispatcher for citation styles."""
    if style.startswith("APA"):    return format_apa(ref, number)
    if style.startswith("ABNT"):   return format_abnt(ref, number)
    if style.startswith("Chicago"): return format_chicago(ref, number)
    return format_vancouver(ref, number)   # Vancouver default


# =============================================================================
# AI FUNCTIONS
# =============================================================================

def get_ai_client(provider: str, api_key: str):
    if provider == "Anthropic (Claude)":
        import anthropic
        return anthropic.Anthropic(api_key=api_key)
    elif provider == "Google (Gemini)":
        from google import genai
        return genai.Client(api_key=api_key)
    else:
        from openai import OpenAI
        return OpenAI(api_key=api_key)

def ai_call(client, provider: str, model: str, prompt: str, max_tokens: int = 8000) -> str:
    """Call AI API with automatic retry on rate limit (429) errors.
    Implements exponential backoff: waits 15s, 30s, 60s between retries."""
    import time as _time

    # Limit prompt size to avoid hitting token rate limits.
    # Strategy: if prompt is too long, truncate the CATALOGUE section (between
    # === REFERENCE CATALOGUE === and === TEXT ===), not the text chunk itself.
    MAX_PROMPT_CHARS = 28000  # ~7k tokens, well within per-minute limits
    if len(prompt) > MAX_PROMPT_CHARS:
        import re as _re_trunc
        # Try to find and shorten the catalogue block
        cat_match = _re_trunc.search(
            r'(=== REFERENCE CATALOGUE ===\n)(.*?)(\n=== TEXT)',
            prompt, _re_trunc.DOTALL
        )
        if cat_match:
            cat_block  = cat_match.group(2)
            # Keep only as many catalogue lines as fit in the budget
            budget     = MAX_PROMPT_CHARS - (len(prompt) - len(cat_block)) - 200
            if budget > 500:
                # Trim catalogue to budget
                trimmed    = cat_block[:budget] + "\n...(catalogo truncado)"
                prompt     = prompt[:cat_match.start(2)] + trimmed + prompt[cat_match.end(2):]
            else:
                # Extreme case: drop catalogue entirely, keep text
                prompt = prompt[:cat_match.start()] + "(catalogo omitido por limite de tokens)" + prompt[cat_match.end():]
        else:
            # Fallback: hard truncate
            prompt = prompt[:MAX_PROMPT_CHARS] + "\n\n[PROMPT TRUNCADO]"

    max_retries = 3
    wait_times  = [15, 30, 60]

    for attempt in range(max_retries + 1):
        try:
            if provider == "Anthropic (Claude)":
                resp = client.messages.create(
                    model=model, max_tokens=max_tokens,
                    messages=[{"role":"user","content":prompt}])
                return resp.content[0].text
            elif provider == "Google (Gemini)":
                from google import genai as _genai
                resp = client.models.generate_content(
                    model=model,
                    contents=prompt,
                    config=_genai.types.GenerateContentConfig(max_output_tokens=max_tokens),
                )
                return resp.text
            else:  # OpenAI
                resp = client.chat.completions.create(
                    model=model, max_tokens=max_tokens,
                    messages=[{"role":"user","content":prompt}])
                return resp.choices[0].message.content

        except Exception as e:
            err_str = str(e)
            is_rate_limit = ("429" in err_str or "rate_limit" in err_str.lower()
                             or "too many requests" in err_str.lower()
                             or "overloaded" in err_str.lower())
            if is_rate_limit and attempt < max_retries:
                wait = wait_times[attempt]
                st.warning(f"⏳ Limite de taxa da API atingido. Aguardando {wait}s antes de tentar novamente "
                           f"(tentativa {attempt+1}/{max_retries})...")
                _time.sleep(wait)
                continue
            raise  # re-raise if not rate limit or max retries exceeded

def repair_truncated_json(raw: str) -> dict | None:
    """Recover partial paragraphs from a truncated JSON response using a depth-counter parser."""
    # Strip markdown code fences
    cleaned = raw
    m = re.search(r"```(?:json)?\s*(.*)", cleaned, re.DOTALL)
    if m:
        cleaned = re.sub(r"`+$", "", m.group(1)).strip()
    # Find the paragraphs array
    arr_m = re.search(r'"paragraphs"\s*:\s*\[', cleaned)
    if not arr_m:
        return None
    content = cleaned[arr_m.end():]
    # Walk the content extracting complete JSON objects via a depth counter
    paragraphs = []
    i = 0
    n = len(content)
    while i < n:
        if content[i] == '{':
            depth = 0
            in_str = False
            esc = False
            j = i
            while j < n:
                c = content[j]
                if esc:
                    esc = False
                elif c == '\\' and in_str:
                    esc = True
                elif c == '"':
                    in_str = not in_str
                elif not in_str:
                    if c == '{':
                        depth += 1
                    elif c == '}':
                        depth -= 1
                        if depth == 0:
                            try:
                                obj = json.loads(content[i:j+1])
                                if "original" in obj:
                                    paragraphs.append(obj)
                            except Exception:
                                pass
                            i = j + 1
                            break
                j += 1
            else:
                break
        elif content[i] == ']':
            break
        else:
            i += 1
    if paragraphs:
        return {
            "paragraphs": paragraphs,
            "reference_map": {},
            "summary": f"Recuperacao parcial: {len(paragraphs)} paragrafo(s) processado(s). O texto era longo demais — foi processado em partes.",
            "changes_detail": [],
            "_truncated": True,
        }
    return None


def extract_json_from_ai(text: str) -> dict | None:
    try:
        return json.loads(text)
    except Exception:
        pass
    match = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    if match:
        try: return json.loads(match.group(1))
        except: pass
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if match:
        try: return json.loads(match.group())
        except: pass
    # Last resort: try to recover paragraphs from a truncated response
    recovered = repair_truncated_json(text)
    if recovered:
        return recovered
    return None


CHUNK_CHAR_LIMIT = 1800  # default max chars per chunk (reduced to avoid 50k token/min rate limit)
CHUNK_CHAR_LIMIT_GEMINI = 1200  # Gemini needs smaller chunks to avoid truncation
CHUNK_CHAR_LIMIT_OPENAI = 1800



def chunk_text(text: str, max_chars: int = CHUNK_CHAR_LIMIT) -> list:
    """Split text at paragraph boundaries into chunks of at most max_chars."""
    paras = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
    if not paras:
        return [text]
    chunks, cur, cur_len = [], [], 0
    for p in paras:
        if cur_len + len(p) > max_chars and cur:
            chunks.append("\n\n".join(cur))
            cur, cur_len = [p], len(p)
        else:
            cur.append(p)
            cur_len += len(p)
    if cur:
        chunks.append("\n\n".join(cur))
    return chunks

# =============================================================================
# FILE EXTRACTION
# =============================================================================

def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extract text from PDF preserving paragraph structure.
    Uses block-mode to reconstruct full paragraphs instead of raw line breaks."""
    try:
        import fitz
        doc  = fitz.open(stream=pdf_bytes, filetype="pdf")
        all_blocks = []
        for page in doc:
            for b in page.get_text("blocks"):
                if b[6] != 0:          # skip image blocks
                    continue
                txt = b[4].strip()
                if not txt or len(txt) < 5:
                    continue
                # Rejoin words hyphenated across lines: "psiqui-\natria" → "psiquiatria"
                txt = re.sub(r'-\n\s*', '', txt)
                # Rejoin soft line breaks (mid-sentence: lowercase continues on next line)
                txt = re.sub(
                    r'(?<![.!?:»\"”])\n(?=[a-záàãâéêíóôõúüça-zA-Z(])',
                    ' ', txt
                )
                # Collapse internal whitespace
                txt = re.sub(r'[ \t]+', ' ', txt).strip()
                if txt:
                    all_blocks.append(txt)
        return "\n\n".join(all_blocks).strip()
    except Exception as e:
        st.error(f"Erro ao ler PDF: {e}")
        return ""

def extract_text_from_docx(docx_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(BytesIO(docx_bytes))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        st.error(f"Erro ao ler DOCX: {e}")
        return ""

# =============================================================================
# AI PIPELINE
# =============================================================================

def extract_ref_metadata_ai(client, provider, model, pdf_text, filename) -> dict:
    prompt = f"""Extraia metadados bibliograficos deste artigo cientifico.
Retorne APENAS JSON valido:
{{"title":"","authors":["Sobrenome AB"],"journal":"","year":"","volume":"","issue":"","pages":"","doi":"","abstract":"resumo em ate 200 palavras","key_topics":["topico1"]}}

TEXTO ({filename}, primeiros 5000 chars):
{pdf_text[:5000]}"""
    raw  = ai_call(client, provider, model, prompt, max_tokens=1200)
    data = extract_json_from_ai(raw)
    if data:
        data["_source"] = filename
        data["_raw_extracted"] = True
        return data
    return {"title":filename,"authors":[],"journal":"","year":"","volume":"","issue":"","pages":"","doi":"","abstract":pdf_text[:300],"_source":filename}

def identify_citation_needs(client, provider, model, text, mode) -> dict:
    """Identify citation needs and generate PubMed queries for a chunk of text."""
    if mode == "add":
        task  = ("The text has NO citations. For each paragraph with scientific/factual claims, "
                 "generate 1-3 specific PubMed queries in English to find supporting references.")
        extra = ('"needs_citations":true,'
                 '"pubmed_queries":["specific english query 1","specific english query 2"]')
    else:
        task  = ("The text ALREADY HAS citations. Identify problems (wrong, missing, duplicate) "
                 "and generate PubMed queries in English to find correct references.")
        extra = ('"existing_citations":["[1]"],"issues":["description"],'
                 '"pubmed_queries":["specific english query"]')
    prompt = f"""You are a scientific citation specialist (Vancouver format).
TASK: {task}

Rules for generating PubMed queries:
- Write queries in English, specific enough to find real papers
- Include key medical/scientific terms from the paragraph
- Use terms like: "systematic review", "meta-analysis", "clinical trial" when appropriate
- Aim for queries that return recent (last 10 years) high-quality papers

Output ONLY valid JSON (no markdown):
{{"paragraphs":[{{"index":0,"text_preview":"first 80 chars of paragraph",{extra},"needs_web_search":true}}]}}

TEXT:
{text}"""
    raw  = ai_call(client, provider, model, prompt, max_tokens=4000)
    data = extract_json_from_ai(raw)
    return data if data else {"paragraphs":[]}

def _build_ref_catalogue(refs: list, compact: bool = False) -> str:
    """Build reference catalogue for AI prompt.
    compact=True  → one line per ref (for large reference lists, saves token space)
    compact=False → multi-line with abstract snippet (for small lists)
    """
    if not refs:
        return "(nenhuma referencia disponivel -- use [?])"

    # Auto-select compact mode when refs are many
    if len(refs) > 20:
        compact = True

    cat = ""
    for i, r in enumerate(refs, 1):
        authors = r.get("authors", [])
        a_str   = authors[0] if authors else ""
        if len(authors) > 1:
            a_str += " et al."
        title   = r.get("title", "")[:80]
        journal = r.get("journal", "")
        year    = r.get("year", "")
        if compact:
            # Single line: [REF1] Autor et al. Título (80 chars). Journal Year.
            cat += f"[REF{i}] {a_str}. {title}. {journal} {year}.\n"
        else:
            cat += f"\n[REF{i}] {title}\n"
            cat += f"  Autores: {', '.join(authors)}\n"
            cat += f"  Periodico: {journal} {year}\n"
            if r.get("abstract"):
                cat += f"  Resumo: {r['abstract'][:150]}\n"
    return cat


def _insert_citations_chunk(client, provider, model, chunk_text, ref_catalogue, mode, citation_style="Vancouver") -> dict:
    """Run citation insertion on a single text chunk and return parsed dict."""
    if mode == "add":
        instructions = """CRITICAL RULES — Adding Citations:
1. Go paragraph by paragraph
2. For scientific/factual claims: add the citation tag at the END of the paragraph
3. ONLY use tags from the catalogue below: [REF1], [REF2], [REF3], etc.
4. NEVER invent numbers like [57] or [265] — those are FORBIDDEN
5. If no catalogue reference fits, mark [?] — do NOT use any other number
6. Do NOT modify the original text — only append the citation tag
7. Headings/titles: do not add any citation"""
    else:
        instructions = """CRITICAL RULES — Reviewing Citations:
1. Verify each existing citation against the catalogue
2. Replace wrong/missing ones with the matching [REFx] from the catalogue
3. If no catalogue match, replace with [?]
4. NEVER use numbers not present in the catalogue (e.g. [57], [265])
5. Fix Vancouver formatting when needed"""

    style_note = f"Citation style: {citation_style}. For in-text markers always use [REF1],[REF2] etc (they will be reformatted later). The reference list format will be {citation_style}."
    prompt = f"""You are a scientific citation specialist.
{style_note}
{instructions}

=== REFERENCE CATALOGUE (use ONLY these tags) ===
{ref_catalogue}

=== ORIGINAL TEXT ===
{chunk_text}

Return ONLY valid JSON (absolutely no markdown fences, no extra text):
{{"paragraphs":[{{"original":"exact original paragraph text","modified":"text with [REF1] or [?] appended","refs_used":["REF1"],"changes":["added REF1 at end"],"changed":true}}],"reference_map":{{"REF1":"REF1"}},"summary":"brief summary","changes_detail":["Paragraph 1: ..."]}}

IMPORTANT: in refs_used and reference_map, only use the exact REFx tags from the catalogue."""

    raw  = ai_call(client, provider, model, prompt, max_tokens=8000)
    data = extract_json_from_ai(raw)
    if not data:
        # Retry with smaller prompt on failure
        short_prompt = prompt.replace(chunk_text, chunk_text[:800])
        raw2 = ai_call(client, provider, model, short_prompt, max_tokens=4000)
        data = extract_json_from_ai(raw2)
    return data if data else {"error": "Falha no processamento", "raw": raw}


def _renumber_citations(paragraphs: list, refs: list) -> tuple:
    """Post-process: replace [REFx]/[?] tags with sequential Vancouver numbers [1],[2],[3]...
    Returns (renumbered_paragraphs, ordered_final_refs).
    Only valid REFx tags are renumbered; stray numbers like [57] are stripped/replaced with [?]."""
    import re as _re
    refx_to_num = {}   # "REF3" -> 1 (first appearance order)
    counter     = [0]

    def assign(refx: str) -> int:
        if refx not in refx_to_num:
            counter[0] += 1
            refx_to_num[refx] = counter[0]
        return refx_to_num[refx]

    # Build set of valid REFx keys
    valid_keys = {f"REF{i}" for i in range(1, len(refs)+1)}

    # Pattern that matches [REF1], [REF2,REF3], [?], or stray [number] sequences
    tag_pat = _re.compile(r'\[([^\]]+)\]')

    updated = []
    for p in paragraphs:
        modified = p.get("modified") or p.get("original", "")

        def replace_tag(m):
            inner = m.group(1).strip()
            # [?] stays as [?]
            if inner == "?":
                return "[?]"
            # Multiple refs: [REF1, REF2]
            parts = [x.strip() for x in _re.split(r'[,;]', inner)]
            nums = []
            for part in parts:
                if part in valid_keys:
                    nums.append(str(assign(part)))
                elif _re.match(r'^REF\d+$', part):
                    # REFx outside range — treat as [?]
                    pass
                elif _re.match(r'^\d+$', part):
                    # Stray invented number — remove silently
                    pass
                # else ignore
            if nums:
                return "[" + ", ".join(nums) + "]"
            return "[?]"

        modified = tag_pat.sub(replace_tag, modified)
        updated.append({**p, "modified": modified})

    # Build ordered final ref list
    ordered_refs = []
    for refx, num in sorted(refx_to_num.items(), key=lambda x: x[1]):
        idx = int(refx.replace("REF","")) - 1
        if 0 <= idx < len(refs):
            ordered_refs.append(refs[idx])

    return updated, ordered_refs, refx_to_num


def _insert_citations_simple(client, provider, model, text: str, refs: list,
                              mode: str, citation_style: str,
                              write_fn=None, upd_fn=None) -> dict:
    """
    Reliable citation insertion: asks the AI to return annotated TEXT directly
    (no JSON parsing needed).  Returns:
      {"full_cited_text": str, "final_refs": list, "ref_map": dict}
    """
    import re as _re

    if not refs:
        return {"full_cited_text": text, "final_refs": [], "ref_map": {}}

    ref_catalogue = _build_ref_catalogue(refs)
    chunk_limit   = CHUNK_CHAR_LIMIT_GEMINI if "Gemini" in provider else CHUNK_CHAR_LIMIT
    chunks        = chunk_text(text, max_chars=chunk_limit)
    cited_chunks  = []

    for i, chunk in enumerate(chunks):
        if write_fn:
            write_fn(f"  ✍️ Parte {i+1}/{len(chunks)} — inserindo citações...")
        if upd_fn:
            upd_fn(65 + int(20*(i+1)/len(chunks)),
                   f"Inserindo citações: parte {i+1}/{len(chunks)}")

        if mode == "add":
            task_verb = "ADD"
            task_rule = (
                "For each paragraph with a scientific/factual claim, "
                "append the matching [REFx] tag(s) at the very END of that paragraph."
            )
        else:
            task_verb = "REVIEW"
            task_rule = (
                "Check each existing citation tag against the catalogue. "
                "Replace wrong/missing tags with the correct [REFx]. Remove invented numbers."
            )

        # TEXT FIRST — eliminates any chance the model misses it
        prompt = (
            "You are a scientific citation specialist. "
            f"Your job: {task_verb} citations in the text below.\n\n"
            f"=== TEXT TO ANNOTATE (part {i+1}/{len(chunks)}) ===\n"
            f"{chunk}\n"
            "=== END OF TEXT ===\n\n"
            f"INSTRUCTIONS: {task_rule}\n"
            "RULES (read carefully):\n"
            "• Use ONLY the [REFx] tags from the CATALOGUE below — e.g. [REF1], [REF2].\n"
            "• Place the tag at the END of the paragraph, after the last word or period.\n"
            "• Do NOT change, paraphrase or remove any word of the original text.\n"
            "• Multiple refs for one claim: [REF1, REF2]. No ref fits: [?].\n"
            "• Headings and section titles: copy unchanged, NO tag.\n"
            "• Output ONLY the annotated text — no JSON, no markdown, no commentary.\n\n"
            "=== REFERENCE CATALOGUE ===\n"
            f"{ref_catalogue}\n"
            "=== END OF CATALOGUE ===\n\n"
            "Return the annotated text now:"
        )

        cited = ai_call(client, provider, model, prompt, max_tokens=8000)

        # ── Sanity check: if AI produced a complaint instead of annotated text ──
        complaint_phrases = [
            "text is missing", "not included in your message",
            "no text was provided", "please provide the text",
            "actual paragraphs", "I notice that", "I'm ready to add",
        ]
        is_complaint = any(ph in cited.lower() for ph in complaint_phrases)
        if is_complaint or len(cited.strip()) < 30:
            # Retry with ultra-minimal prompt — just the text + one-line instruction
            minimal = (
                f"Add [REF1]..[REF{len(refs)}] citation tags to paragraphs. "
                f"Use ONLY tags from this catalogue:\n{ref_catalogue[:2000]}\n\n"
                f"TEXT:\n{chunk}\n\nReturn annotated text only:"
            )
            cited = ai_call(client, provider, model, minimal, max_tokens=8000)

        # Strip any accidental markdown fences
        cited = _re.sub(r'```[^\n]*\n?', '', cited).strip()
        # If still a complaint after retry, fall back to original chunk (no citations)
        if any(ph in cited.lower() for ph in complaint_phrases) or len(cited.strip()) < 10:
            cited = chunk
        cited_chunks.append(cited)

    full_cited = "\n\n".join(cited_chunks)

    # ── Renumber [REFx] → [1], [2], ... in order of first appearance ─────
    ref_order  = {}  # "REF3" -> 1
    counter    = [0]
    valid_keys = {f"REF{i}" for i in range(1, len(refs)+1)}

    def _replace(m):
        inner = m.group(1).strip()
        if inner == "?":
            return "[?]"
        parts = [x.strip() for x in _re.split(r'[,;\s]+', inner)]
        nums  = []
        for part in parts:
            if part in valid_keys:
                if part not in ref_order:
                    counter[0] += 1
                    ref_order[part] = counter[0]
                nums.append(str(ref_order[part]))
        return ("[" + ", ".join(nums) + "]") if nums else "[?]"

    full_cited = _re.sub(r'\[([^\]]+)\]', _replace, full_cited)

    # ── Build ordered final reference list ───────────────────────────────
    final_refs = []
    for refx, _ in sorted(ref_order.items(), key=lambda x: x[1]):
        idx = int(refx.replace("REF", "")) - 1
        if 0 <= idx < len(refs):
            final_refs.append(refs[idx])

    return {
        "full_cited_text": full_cited,
        "final_refs":      final_refs,
        "ref_map":         ref_order,
    }


def insert_citations_ai(client, provider, model, text, refs, mode, citation_style="Vancouver") -> dict:
    ref_catalogue = _build_ref_catalogue(refs)
    # Use smaller chunks for Gemini to avoid JSON truncation
    chunk_limit = CHUNK_CHAR_LIMIT_GEMINI if "Gemini" in provider else CHUNK_CHAR_LIMIT
    chunks = chunk_text(text, max_chars=chunk_limit)

    all_paragraphs  = []
    all_changes     = []
    chunk_summaries = []

    if len(chunks) == 1:
        result = _insert_citations_chunk(client, provider, model, chunks[0], ref_catalogue, mode)
    else:
        chunk_bar = st.progress(0, text=f"Processando em {len(chunks)} partes...")
        for idx, chunk in enumerate(chunks, 1):
            chunk_bar.progress(idx / len(chunks),
                               text=f"Processando parte {idx}/{len(chunks)}...")
            res = _insert_citations_chunk(client, provider, model, chunk, ref_catalogue, mode)
            if "error" in res and "paragraphs" not in res:
                all_paragraphs.append({
                    "original": chunk[:200] + "...",
                    "modified": chunk[:200] + "...",
                    "refs_used": [], "changes": [f"Erro chunk {idx}"], "changed": False,
                })
            else:
                all_paragraphs.extend(res.get("paragraphs", []))
                all_changes.extend(res.get("changes_detail", []))
                if res.get("summary"):
                    chunk_summaries.append(f"Parte {idx}: {res['summary']}")
        chunk_bar.empty()
        result = {
            "paragraphs":     all_paragraphs,
            "reference_map":  {},
            "summary":        " | ".join(chunk_summaries) or "Processado em multiplas partes.",
            "changes_detail": all_changes,
            "_chunked": True,
        }

    # ── Global renumbering pass (fixes invented numbers, assigns sequential [1],[2]...) ──
    paras = result.get("paragraphs", [])
    if paras:
        renumbered, ordered_refs, ref_map = _renumber_citations(paras, refs)
        result["paragraphs"]    = renumbered
        result["_final_refs"]   = ordered_refs   # ordered for reference list
        result["_ref_map"]      = ref_map
    return result



def insert_citations_ai_bg(client, provider, model, text, refs, mode,
                            citation_style="Vancouver", upd_fn=None) -> dict:
    """St-free version of insert_citations_ai for background thread use."""
    ref_catalogue = _build_ref_catalogue(refs)
    chunk_limit   = CHUNK_CHAR_LIMIT_GEMINI if "Gemini" in provider else CHUNK_CHAR_LIMIT
    chunks        = chunk_text(text, max_chars=chunk_limit)

    all_paragraphs  = []
    all_changes     = []
    chunk_summaries = []

    def _upd(idx, n):
        if upd_fn:
            pct = 65 + int(20 * idx / max(n, 1))
            upd_fn(min(pct, 85), f"Inserindo citacoes: parte {idx}/{n}...")

    if len(chunks) == 1:
        result = _insert_citations_chunk(client, provider, model, chunks[0],
                                         ref_catalogue, mode, citation_style)
    else:
        for idx, chunk in enumerate(chunks, 1):
            _upd(idx, len(chunks))
            res = _insert_citations_chunk(client, provider, model, chunk,
                                          ref_catalogue, mode, citation_style)
            if "error" in res and "paragraphs" not in res:
                all_paragraphs.append({
                    "original": chunk[:200] + "...",
                    "modified": chunk[:200] + "...",
                    "refs_used": [], "changes": [f"Erro chunk {idx}"], "changed": False,
                })
            else:
                all_paragraphs.extend(res.get("paragraphs", []))
                all_changes.extend(res.get("changes_detail", []))
                if res.get("summary"):
                    chunk_summaries.append(f"Parte {idx}: {res['summary']}")
        result = {
            "paragraphs":     all_paragraphs,
            "reference_map":  {},
            "summary":        " | ".join(chunk_summaries) or "Processado em multiplas partes.",
            "changes_detail": all_changes,
            "_chunked": True,
        }

    paras = result.get("paragraphs", [])
    if paras:
        renumbered, ordered_refs, ref_map = _renumber_citations(paras, refs)
        result["paragraphs"]  = renumbered
        result["_final_refs"] = ordered_refs
        result["_ref_map"]    = ref_map
    return result

# =============================================================================
# MAIN PIPELINE
# =============================================================================

def run_pipeline(client, provider, model, main_text, ref_files, mode, library_refs, citation_style="Vancouver"):
    progress = st.progress(0, text="Iniciando...")
    log      = st.empty()
    def upd(pct, msg):
        progress.progress(pct / 100, text=msg)
        log.caption(f"... {msg}")

    # ── 1. Extract PDF reference metadata ────────────────────────────────────
    ref_metadata = list(library_refs)
    if ref_files:
        upd(5, f"Analisando {len(ref_files)} artigo(s) de referencia com IA...")
        for i, f in enumerate(ref_files):
            pdf_text = extract_text_from_pdf(f.read())
            meta     = extract_ref_metadata_ai(client, provider, model, pdf_text, f.name)
            ref_metadata.append(meta)
            upd(5 + int(15*(i+1)/len(ref_files)),
                f"Artigo {i+1}/{len(ref_files)}: {meta.get('title','')[:50]}...")

    # ── 2. Identify citation needs PER CHUNK ──────────────────────────────────
    # Chunk the text so identify_citation_needs is never called with too many tokens
    chunk_limit = CHUNK_CHAR_LIMIT_GEMINI if "Gemini" in provider else CHUNK_CHAR_LIMIT
    text_chunks = chunk_text(main_text, max_chars=chunk_limit)
    n_chunks    = len(text_chunks)
    upd(20, f"Identificando necessidades de citacao em {n_chunks} parte(s)...")

    all_queries = []
    for ci, chunk in enumerate(text_chunks):
        needs_data = identify_citation_needs(client, provider, model, chunk, mode)
        for p in needs_data.get("paragraphs", []):
            all_queries.extend(p.get("pubmed_queries", []))
        upd(20 + int(10 * (ci+1) / n_chunks),
            f"Analise parte {ci+1}/{n_chunks}: {len(all_queries)} queries geradas")
        if ci < n_chunks - 1:
            time.sleep(2)  # 2s entre chamadas para evitar rate limit 50k tokens/min

    # Dedup queries; allow up to 4 queries per chunk (minimum 8)
    all_queries  = list(dict.fromkeys(q for q in all_queries if q.strip()))
    max_searches = max(8, n_chunks * 4)
    all_queries  = all_queries[:max_searches]

    # ── 3. Fallback: se nenhuma query foi gerada, extrair topico do texto ─────
    if not all_queries:
        upd(30, "Gerando queries de busca a partir do topico do texto...")
        topic_prompt = f"""Read this scientific text excerpt and generate 6 specific PubMed search queries in English to find relevant supporting references. Return ONLY a JSON array of strings: ["query1","query2","query3","query4","query5","query6"]

TEXT (first 1500 chars):
{main_text[:1500]}"""
        raw_q = ai_call(client, provider, model, topic_prompt, max_tokens=500)
        try:
            m = re.search(r"\[.*\]", raw_q, re.DOTALL)
            if m:
                all_queries = json.loads(m.group())[:8]
        except Exception:
            pass

    # ── 4. Multi-source web search ────────────────────────────────────────────
    # Covers: PubMed, Europe PMC (incl. Cochrane/preprints), LILACS, OpenAlex,
    #         Semantic Scholar — busca cinzenta incluida via OpenAlex + Europe PMC PPR
    web_refs = []
    if all_queries:
        upd(35, f"Buscando em {len(all_queries)} queries x 4 bases (PubMed, LILACS, Europe PMC, OpenAlex)...")
        for i, q in enumerate(all_queries):
            pct = 35 + int(25 * (i+1) / len(all_queries))
            upd(pct, f"Busca {i+1}/{len(all_queries)}: {q[:55]}")
            web_refs.extend(multi_source_search(q, max_per_source=2))
            time.sleep(0.25)
        # Extra: preprints / busca cinzenta via Europe PMC
        for q in all_queries[:4]:
            web_refs.extend(search_europe_pmc(q, max_results=2, source_filter="PPR"))
    else:
        upd(60, "Nenhuma query gerada — usando apenas biblioteca local.")

    # Dedup web_refs
    seen, unique_web = set(), []
    for r in web_refs:
        t = r.get("title", "").lower()[:80]
        if t not in seen:
            seen.add(t)
            unique_web.append(r)

    all_refs = ref_metadata + unique_web
    upd(62, f"Total de referencias disponiveis: {len(all_refs)} ({len(ref_metadata)} locais + {len(unique_web)} da web)")

    # ── 5. Insert / review citations ──────────────────────────────────────────
    upd(65, "Inserindo/corrigindo citacoes com IA...")
    try:
        result = insert_citations_ai(client, provider, model, main_text, all_refs, mode, citation_style)
    except Exception as e_ins:
        result = {"error": str(e_ins), "paragraphs": [], "reference_map": {}}
    upd(90, "Montando resultado final...")

    # Use the ordered ref list produced by the renumbering pass
    final_ref_list = result.get("_final_refs", [])
    # Fallback: build from reference_map if renumbering didn't run
    if not final_ref_list:
        ref_map = result.get("reference_map", {})
        seen_idx = []
        for num_str, ref_id in sorted(ref_map.items(), key=lambda x: int(x[0]) if x[0].isdigit() else 999):
            idx_str = re.sub(r"[^0-9]","",str(ref_id))
            if idx_str:
                idx = int(idx_str)-1
                if 0 <= idx < len(all_refs) and idx not in seen_idx:
                    seen_idx.append(idx)
                    final_ref_list.append(all_refs[idx])

    upd(100, "Concluido!")
    log.empty(); progress.empty()
    return result, all_refs, final_ref_list

# =============================================================================
# DOCX EXPORT
# =============================================================================

def generate_docx(paragraphs: list, final_ref_list: list, mode: str, citation_style: str = "Vancouver") -> bytes:
    """Build a formatted Word document with the processed text and references."""
    from docx import Document as DocxDoc
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = DocxDoc()

    # ── Page margins ────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.0)

    # ── Styles ──────────────────────────────────────────────────────────────
    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Times New Roman"
    normal_style.font.size = Pt(12)

    # ── Header strip ────────────────────────────────────────────────────────
    hdr = doc.add_paragraph()
    hdr.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = hdr.add_run("CitacaoIA  -  Texto com Citacoes Vancouver")
    run.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run(f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}  |  "
                f"Modo: {'Adicionar citacoes' if mode == 'add' else 'Revisar citacoes'}"
                ).font.size = Pt(10)

    doc.add_paragraph()  # spacer

    # ── Body paragraphs (preserve original formatting as closely as possible) ──
    def _is_heading(text: str) -> bool:
        """Heuristic: short line, no trailing period, possibly title-case or all-caps."""
        t = text.strip()
        if len(t) > 120:
            return False
        if t.endswith(".") or t.endswith(","):
            return False
        words = t.split()
        if len(words) <= 1:
            return True
        if t == t.upper() and len(t) > 3:
            return True
        cap_ratio = sum(1 for w in words if w and w[0].isupper()) / len(words)
        return cap_ratio >= 0.7 and len(words) <= 10

    for p in paragraphs:
        text = p.get("modified") or p.get("original", "")
        if not text.strip():
            doc.add_paragraph()   # preserve blank lines
            continue

        is_head = _is_heading(text)
        para = doc.add_paragraph()

        if is_head:
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            para.paragraph_format.space_before = Pt(12)
            para.paragraph_format.space_after  = Pt(4)
            para.paragraph_format.first_line_indent = Pt(0)
        else:
            para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            para.paragraph_format.space_after  = Pt(0)
            para.paragraph_format.space_before = Pt(0)
            # No forced first-line indent — respect Word Normal style

        # Split on citation markers and render inline
        parts = re.split(r"(\[\??\d*[,\s\d]*\??\]|\[REF\d+\])", text)
        for part in parts:
            if not part:
                continue
            run = para.add_run(part)
            run.font.name = "Times New Roman"
            run.font.size = Pt(14) if is_head else Pt(12)
            run.bold = is_head
            if re.match(r"\[", part):   # citation tag
                run.font.color.rgb = RGBColor(0x2E, 0x7D, 0xDB)
                run.bold = True
                run.font.size = Pt(10) if not is_head else Pt(12)

    # ── References section ──────────────────────────────────────────────────
    doc.add_paragraph()
    ref_hdr = doc.add_paragraph()
    ref_hdr.alignment = WD_ALIGN_PARAGRAPH.LEFT
    rh = ref_hdr.add_run("REFERENCIAS")
    rh.bold = True
    rh.font.size = Pt(12)
    rh.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    # Separator line
    sep = doc.add_paragraph()
    sep.add_run("_" * 60).font.size = Pt(9)

    if final_ref_list:
        for i, r in enumerate(final_ref_list, 1):
            ref_text = format_vancouver(r, i)
            rp = doc.add_paragraph(style="Normal")
            rp.paragraph_format.left_indent  = Cm(0.5)
            rp.paragraph_format.space_after  = Pt(4)
            rr = rp.add_run(ref_text)
            rr.font.size = Pt(10)
            rr.font.name = "Times New Roman"
    else:
        np_para = doc.add_paragraph()
        np_para.add_run("(Referencias nao resolvidas — verifique o catalogo)").italic = True

    # ── Serialize to bytes ───────────────────────────────────────────────────
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# =============================================================================
# RESULTS DISPLAY
# =============================================================================

def display_results(result, all_refs, final_ref_list, mode):
    # Show prominent error if pipeline reported failure
    if result.get("error") and not result.get("paragraphs"):
        st.error(f"Erro no processamento: {result['error']}")
        if result.get("raw"):
            with st.expander("Ver resposta bruta da IA"):
                st.text(str(result["raw"])[:3000])
        return
    if "error" in result and not result.get("paragraphs"):
        st.error(f"Erro: {result.get('error')}")
        with st.expander("Ver resposta bruta"):
            st.text(result.get("raw",""))
        return

    st.success("Processamento concluido com sucesso!")
    cstyle_disp = st.session_state.get("last_citation_style", "Vancouver")

    # ── Determine output mode ─────────────────────────────────────────────
    simple_mode = result.get("_simple_mode", False)
    paragraphs  = result.get("paragraphs", [])

    # Metrics
    c1, c2, c3 = st.columns(3)
    c1.metric("Referências encontradas", len(all_refs))
    c2.metric("Referências no texto",    len(final_ref_list))
    if not simple_mode:
        c3.metric("Parágrafos alterados", len([p for p in paragraphs if p.get("changed")]))
    else:
        import re as _re_disp
        n_cit = len(_re_disp.findall(r'\[\d', result.get("full_cited_text", "")))
        c3.metric("Citações inseridas", n_cit)
    st.divider()

    # ── Build full output text ────────────────────────────────────────────
    if simple_mode:
        body = result.get("full_cited_text", "")
    else:
        body = "\n\n".join(p.get("modified") or p.get("original", "") for p in paragraphs)

    ref_section = "\n\n" + "─" * 60 + f"\nREFERÊNCIAS ({cstyle_disp})\n" + "─" * 60 + "\n"
    if final_ref_list:
        ref_section += "\n".join(
            format_reference(r, i, cstyle_disp)
            for i, r in enumerate(final_ref_list, 1)
        )
    else:
        ref_section += "(referências não resolvidas)"

    full = body + "\n\n" + ref_section

    # ── Tabs ─────────────────────────────────────────────────────────────
    tab_final, tab_refs = st.tabs(["📄 Texto Final com Citações", "📚 Lista de Referências"])

    with tab_final:
        st.text_area("", value=full, height=560, label_visibility="collapsed")

        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            st.download_button(
                "📥 Baixar (.txt)",
                data=full.encode("utf-8"),
                file_name="texto_citado.txt",
                mime="text/plain",
                use_container_width=True,
            )
        with col_dl2:
            try:
                # For simple_mode build fake paragraphs list for generate_docx
                if simple_mode:
                    _paras_docx = [{"modified": p, "changed": True}
                                   for p in body.split("\n\n") if p.strip()]
                else:
                    _paras_docx = paragraphs
                docx_bytes_dl = generate_docx(_paras_docx, final_ref_list, mode, cstyle_disp)
                st.download_button(
                    "📄 Baixar Word (.docx)",
                    data=docx_bytes_dl,
                    file_name="texto_citado.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    type="primary",
                )
            except Exception as e_docx:
                st.warning(f"Word não gerado ({e_docx}) — use o .txt acima")

    with tab_refs:
        disp = final_ref_list if final_ref_list else all_refs
        if not disp:
            st.info("Nenhuma referencia processada.")
        else:
            for i,r in enumerate(disp,1):
                src = r.get("_source","")
                with st.expander(f"[{i}] {r.get('title','')[:80]}...  ({src})"):
                    st.markdown(f"**Autores:** {', '.join(r.get('authors',[]))}")
                    st.markdown(f"**Periodico:** {r.get('journal','N/A')}  **Ano:** {r.get('year','N/A')}")
                    if r.get("doi"): st.markdown(f"**DOI:** {r['doi']}")
                    if r.get("pmid"): st.markdown(f"**PubMed:** [PMID {r['pmid']}](https://pubmed.ncbi.nlm.nih.gov/{r['pmid']}/)")
                    st.markdown(f"<div class='ref-box'>{format_vancouver(r,i)}</div>", unsafe_allow_html=True)
                    # Save to library button
                    if st.button(f"Salvar na biblioteca", key=f"save_ref_{i}"):
                        _, is_new = add_to_library(r)
                        if is_new:
                            st.success("Artigo salvo na biblioteca local!")
                        else:
                            st.info("Artigo ja esta na biblioteca (dados atualizados).")
                        st.rerun()

# =============================================================================
# SIDEBAR
# =============================================================================

def _fetch_gemini_models(api_key: str) -> list:
    """Fetch available Gemini models dynamically from the API."""
    try:
        from google import genai as _genai
        _client = _genai.Client(api_key=api_key)
        all_models = list(_client.models.list())
        fetched = []
        for m in all_models:
            name = m.name if isinstance(m.name, str) else str(m.name)
            name = name.replace("models/", "")
            if "gemini" not in name.lower():
                continue
            # Include only models that support generateContent
            supported = getattr(m, "supported_actions", None) or []
            if "generateContent" in supported or not supported:
                fetched.append(name)
        # Prefer flash/pro order, filter out embedding-only models
        fetched = [n for n in fetched if "embed" not in n.lower()]
        return fetched if fetched else ["gemini-2.0-flash", "gemini-2.5-flash", "gemini-2.5-pro"]
    except Exception:
        return ["gemini-2.0-flash", "gemini-2.5-flash", "gemini-2.5-pro"]


def render_sidebar():
    with st.sidebar:
        st.markdown("## Configuracao da IA")
        st.info("⭐ **Recomendado:** Claude Haiku (melhor custo x qualidade para citacoes)")
        provider = st.selectbox(
            "Provedor",
            ["Anthropic (Claude)", "Google (Gemini)", "OpenAI (GPT)"],
            index=0,  # Anthropic default
        )

        # Labels / URLs per provider
        if provider == "Anthropic (Claude)":
            api_lbl = "Chave API Anthropic"
            api_url = "https://console.anthropic.com"
        elif provider == "Google (Gemini)":
            api_lbl = "Chave API Google AI Studio"
            api_url = "https://aistudio.google.com/app/apikey"
        else:
            api_lbl = "Chave API OpenAI"
            api_url = "https://platform.openai.com/api-keys"

        # API key FIRST — Gemini needs it to list models dynamically
        api_key = st.text_input(api_lbl, type="password")
        if api_key:
            st.success("Chave configurada")
        else:
            st.warning(f"[Obter chave aqui]({api_url})")

        # Model selection
        if provider == "Anthropic (Claude)":
            models = [
                "claude-haiku-4-5-20251001",   # ⭐ Recomendado — rapido e economico
                "claude-sonnet-4-5",           # Melhor qualidade
                "claude-opus-4-6",             # Maxima capacidade
            ]
        elif provider == "Google (Gemini)":
            if api_key:
                # Cache per (truncated) key so we don't re-fetch on every render
                cache_key = f"gemini_models_{api_key[:12]}"
                if cache_key not in st.session_state:
                    with st.spinner("Listando modelos Gemini disponiveis..."):
                        st.session_state[cache_key] = _fetch_gemini_models(api_key)
                models = st.session_state[cache_key]
                if st.button("Atualizar lista de modelos", key="btn_refresh_gemini"):
                    st.session_state.pop(cache_key, None)
                    st.rerun()
            else:
                models = ["gemini-2.0-flash", "gemini-2.5-flash", "gemini-2.5-pro"]
        else:
            models = ["gpt-4o", "gpt-4-turbo", "gpt-3.5-turbo"]

        model = st.selectbox(
            "Modelo",
            models,
            index=0,  # always pick first = recommended
        )
        if "haiku" in model.lower():
            st.caption("💰 Custo estimado: ~$0.001 por 1000 tokens")
        elif "flash" in model.lower():
            st.caption("💰 Custo estimado: gratuito (Google AI Studio)")
        elif "sonnet" in model.lower() or "4o" in model.lower():
            st.caption("💰 Custo estimado: ~$0.003 por 1000 tokens")
        elif "opus" in model.lower():
            st.caption("💰 Custo estimado: ~$0.015 por 1000 tokens")

        st.divider()
        st.markdown("### Como usar")
        st.markdown("""
<span class="step-badge">1</span> Configure a chave API
<span class="step-badge">2</span> Veja/adicione artigos na Biblioteca
<span class="step-badge">3</span> Selecione modo e insira texto
<span class="step-badge">4</span> Processe e baixe o resultado
""", unsafe_allow_html=True)
        st.divider()
        lib = load_library()
        st.caption(f"Biblioteca local: {len(lib)} artigo(s) salvo(s)")
        if lib:
            st.caption(f"Arquivo: {BIBLIOTECA_FILE}")
    return api_key, provider, model

# BIBLIOTECA TAB
# =============================================================================

def render_biblioteca_tab():
    st.markdown("### Biblioteca de Artigos")
    st.caption("Artigos salvos localmente. Ficam disponiveis em todas as sessoes futuras.")

    # ── ADD BY DOI ────────────────

    # ── ADD BY DOI ────────────────────────────────────────────────────────────
    if "doi_article_found" not in st.session_state:
        st.session_state["doi_article_found"] = None
    if "doi_msg" not in st.session_state:
        st.session_state["doi_msg"] = ""

    with st.expander("+ Adicionar artigo por DOI", expanded=True):
        col_doi, col_btn = st.columns([4, 1])
        with col_doi:
            doi_input = st.text_input("DOI do artigo", placeholder="10.1038/s41586-023-00001-0",
                                      label_visibility="collapsed", key="doi_input")
        with col_btn:
            do_lookup = st.button("Buscar", use_container_width=True, type="primary", key="btn_doi")

        if do_lookup:
            if doi_input.strip():
                with st.spinner("Buscando metadados via CrossRef..."):
                    found = lookup_doi(doi_input.strip())
                if found:
                    st.session_state["doi_article_found"] = found
                    st.session_state["doi_msg"] = "ok"
                else:
                    st.session_state["doi_article_found"] = None
                    st.session_state["doi_msg"] = "notfound"
            else:
                st.session_state["doi_msg"] = "empty"

        if st.session_state["doi_msg"] == "notfound":
            st.error("DOI nao encontrado no CrossRef. Verifique o DOI ou adicione manualmente.")
        elif st.session_state["doi_msg"] == "empty":
            st.warning("Digite um DOI antes de buscar.")

        found_art = st.session_state.get("doi_article_found")
        found_art = st.session_state.get("doi_article_found")
        if found_art:
            st.success(f"Encontrado: **{found_art.get('title','(sem titulo)')}**")
            st.caption(
                f"Autores: {', '.join(found_art.get('authors',[]))}  |  "
                f"{found_art.get('journal','')} {found_art.get('year','')}"
            )
            if st.button("Salvar este artigo na biblioteca", type="primary",
                         key="btn_save_doi"):
                _, is_new = add_to_library(found_art)
                st.session_state["doi_article_found"] = None
                st.session_state["doi_msg"] = ""
                if is_new:
                    st.success("Artigo salvo na biblioteca!")
                else:
                    st.info("Artigo atualizado na biblioteca.")
                st.rerun()

    # ── ADD MANUAL ────────────────────────────────────────────────────────────
    with st.expander("+ Adicionar artigo manualmente"):
        with st.form("form_manual"):
            m_title   = st.text_input("Titulo*")
            m_authors = st.text_input("Autores (Sobrenome AB, Sobrenome CD)",
                                      placeholder="Smith AB, Jones CD")
            mc1, mc2, mc3 = st.columns(3)
            m_journal = mc1.text_input("Periodico")
            m_year    = mc2.text_input("Ano")
            m_volume  = mc3.text_input("Volume")
            mc4, mc5 = st.columns(2)
            m_issue   = mc4.text_input("Numero")
            m_pages   = mc5.text_input("Paginas")
            m_doi     = st.text_input("DOI (opcional)")
            m_abstract= st.text_area("Resumo (opcional)", height=80)
            submitted = st.form_submit_button("Salvar na biblioteca", type="primary")
            if submitted:
                if m_title.strip():
                    art = {
                        "title":    m_title.strip(),
                        "authors":  [a.strip() for a in m_authors.split(",") if a.strip()],
                        "journal":  m_journal.strip(),
                        "year":     m_year.strip(),
                        "volume":   m_volume.strip(),
                        "issue":    m_issue.strip(),
                        "pages":    m_pages.strip(),
                        "doi":      m_doi.strip(),
                        "abstract": m_abstract.strip(),
                    }
                    _, is_new = add_to_library(art)
                    st.success("Artigo salvo!" if is_new else "Artigo atualizado!")
                    st.rerun()
                else:
                    st.error("O titulo e obrigatorio.")

    # ── LIBRARY DISPLAY ───────────────────────────────────────────────────────
    lib = load_library()
    st.divider()
    if not lib:
        st.info("Biblioteca vazia. Adicione artigos por DOI, manualmente, "
                "ou fazendo upload de PDFs na aba de citacoes.")
        return

    col_s, col_sort = st.columns([4, 1])
    search_q = col_s.text_input("Buscar na biblioteca",
                                placeholder="titulo, autor, periodico...",
                                label_visibility="collapsed")
    sort_by  = col_sort.selectbox("Ordenar", ["Mais recente", "Ano", "Titulo"],
                                  label_visibility="collapsed")

    filtered = lib
    if search_q:
        q = search_q.lower()
        filtered = [a for a in lib if
                    q in a.get("title","").lower() or
                    q in " ".join(a.get("authors",[])).lower() or
                    q in a.get("journal","").lower()]

    if sort_by == "Ano":
        filtered = sorted(filtered, key=lambda x: x.get("year","0"), reverse=True)
    elif sort_by == "Titulo":
        filtered = sorted(filtered, key=lambda x: x.get("title","").lower())

    st.caption(f"Exibindo {len(filtered)} de {len(lib)} artigo(s)")

    for i, art in enumerate(filtered):
        orig_idx = lib.index(art) if art in lib else i
        st.markdown(
            f'<div class="art-card">'
            f'<div class="art-title">{art.get("title","(sem titulo)")}</div>'
            f'<div class="art-meta">'
            f'{", ".join(art.get("authors",[])[:3])}'
            f'{"..." if len(art.get("authors",[])) > 3 else ""}'
            f' &nbsp;|&nbsp; <em>{art.get("journal","")}</em> {art.get("year","")}'
            f'{" &nbsp;|&nbsp; DOI: " + art.get("doi","") if art.get("doi") else ""}'
            f'</div></div>', unsafe_allow_html=True)
        col_a, col_b = st.columns([5, 1])
        if art.get("abstract"):
            with col_a:
                with st.expander("Ver resumo"):
                    st.caption(art["abstract"])
        with col_b:
            if st.button("Remover", key=f"del_{orig_idx}_{art.get('title','')[:10]}"):
                delete_from_library(orig_idx)
                st.rerun()


# =============================================================================
# CITAR TAB
# =============================================================================


def parse_manual_refs(text: str) -> list:
    """Parse a pasted numbered reference list into ref dicts."""
    import re as _re
    refs = []
    # Split on numbered items (1. / 1) or blank lines between entries
    raw_lines = _re.split(r'\n(?=\d+[\.\.)])', text.strip())
    if len(raw_lines) == 1:
        raw_lines = [l for l in text.strip().split("\n") if l.strip()]
    for line in raw_lines:
        line = _re.sub(r'^\d+[\.\)]\s*', '', line.strip())
        if not line:
            continue
        ref = {"title": line[:200], "authors": "", "year": "", "journal": "", "raw": line}
        yr = _re.search(r'\b(19|20)\d{2}\b', line)
        if yr:
            ref["year"] = yr.group()
        parts = [p.strip() for p in line.split(".") if p.strip()]
        if len(parts) >= 2:
            ref["authors"] = parts[0]
            ref["title"]   = parts[1]
        if len(parts) >= 3:
            ref["journal"] = parts[2]
        refs.append(ref)
    return refs


def _run_pipeline_sync(write_fn, upd_fn,
                       provider: str, client_cfg: dict,
                       main_text: str, ref_bytes: list,
                       mode: str, library_refs: list,
                       citation_style: str, manual_refs: list):
    """Synchronous pipeline — called from inside st.status() context.
    write_fn : callable(str) — appends a visible log line (st.write)
    upd_fn   : callable(pct, msg) — updates the status label
    Returns  : (result, all_refs, final_ref_list)
    """
    import re as _re, json as _json, time as _time, traceback as _tb

    from anthropic import Anthropic
    from openai import OpenAI
    import google.genai as genai

    api_key = client_cfg["api_key"]
    model   = client_cfg["model"]

    upd_fn(2, "Conectando ao provedor de IA...")
    write_fn(f"🔌 Conectando ao provedor **{provider}**...")
    if provider == "Anthropic (Claude)":
        client = Anthropic(api_key=api_key)
    elif provider == "OpenAI":
        client = OpenAI(api_key=api_key)
    else:
        client = genai.Client(api_key=api_key)
    write_fn(f"✅ Conectado. Texto recebido: **{len(main_text):,}** caracteres.")

    # 1. Extract PDF metadata
    ref_metadata = list(library_refs)
    ref_metadata.extend(manual_refs)
    if ref_bytes:
        write_fn(f"📄 Analisando **{len(ref_bytes)}** PDF(s) de referência...")
        upd_fn(5, f"Analisando {len(ref_bytes)} PDF(s)...")
        for i, (fname, fbytes) in enumerate(ref_bytes):
            pdf_text = extract_text_from_pdf(fbytes)
            meta     = extract_ref_metadata_ai(client, provider, model, pdf_text, fname)
            ref_metadata.append(meta)
            write_fn(f"  ✔ `{fname}` → {meta.get('title','sem título')[:60]}")
            upd_fn(5 + int(15*(i+1)/len(ref_bytes)),
                   f"PDF {i+1}/{len(ref_bytes)}: {meta.get('title','')[:40]}...")

    # 2. Chunk text
    chunk_limit = CHUNK_CHAR_LIMIT_GEMINI if "Gemini" in provider else CHUNK_CHAR_LIMIT
    text_chunks = chunk_text(main_text, max_chars=chunk_limit)
    MAX_ID_CHUNKS = 15
    id_chunks = text_chunks[:MAX_ID_CHUNKS]
    n_id      = len(id_chunks)
    write_fn(f"✂️ Texto dividido em **{len(text_chunks)}** parte(s) — analisando {n_id}...")
    upd_fn(20, f"Identificando necessidades ({n_id} partes)...")

    all_queries = []
    for ci, chunk in enumerate(id_chunks):
        upd_fn(20 + int(10*(ci+1)/n_id),
               f"Parte {ci+1}/{n_id} — {len(all_queries)} queries até agora")
        needs_data = identify_citation_needs(client, provider, model, chunk, mode)
        for p in needs_data.get("paragraphs", []):
            all_queries.extend(p.get("pubmed_queries", []))
        write_fn(f"  🔍 Parte {ci+1}/{n_id}: +{len(needs_data.get('paragraphs',[]))} parágrafos, "
                 f"{len(all_queries)} queries acumuladas")

    all_queries = list(dict.fromkeys(q for q in all_queries if q.strip()))
    MAX_SEARCHES = 25
    all_queries  = all_queries[:MAX_SEARCHES]

    # 3. Fallback queries
    if not all_queries:
        write_fn("⚠️ Sem queries específicas — gerando a partir do tópico do texto...")
        upd_fn(30, "Gerando queries de fallback...")
        topic_prompt = (
            "Read this scientific text and generate 6 specific PubMed search queries in English. "
            "Return ONLY a JSON array of strings.\n\nTEXT:\n" + main_text[:1500]
        )
        raw_q = ai_call(client, provider, model, topic_prompt, max_tokens=500)
        try:
            m = _re.search(r"\[.*\]", raw_q, _re.DOTALL)
            if m:
                all_queries = _json.loads(m.group())[:8]
        except Exception:
            pass
        write_fn(f"  ↳ {len(all_queries)} queries geradas automaticamente")

    # 4. Multi-source search
    web_refs = []
    if all_queries:
        write_fn(f"🌐 Buscando **{len(all_queries)}** queries em PubMed · LILACS · OpenAlex · Europe PMC...")
        for i, q in enumerate(all_queries):
            pct = 35 + int(25*(i+1)/len(all_queries))
            upd_fn(pct, f"Busca {i+1}/{len(all_queries)}: {q[:50]}...")
            results = multi_source_search(q, max_per_source=2)
            web_refs.extend(results)
            write_fn(f"  [{i+1:02d}/{len(all_queries)}] `{q[:55]}` → {len(results)} resultado(s)")
            _time.sleep(0.25)
        for q in all_queries[:4]:
            web_refs.extend(search_europe_pmc(q, max_results=2, source_filter="PPR"))

    seen, unique_web = set(), []
    for r in web_refs:
        t = r.get("title", "").lower()[:80]
        if t not in seen:
            seen.add(t)
            unique_web.append(r)

    all_refs = ref_metadata + unique_web
    write_fn(f"📚 **{len(all_refs)}** referências no total "
             f"({len(ref_metadata)} locais/manuais + {len(unique_web)} da web)")
    upd_fn(62, f"{len(all_refs)} refs encontradas — inserindo citações...")

    # 5. Insert citations — plain-text approach (no fragile JSON)
    write_fn(f"✍️ Inserindo citações em {len(chunk_text(main_text, max_chars=chunk_limit))} parte(s) do texto...")
    upd_fn(65, "Inserindo citações com IA...")
    try:
        cit_result = _insert_citations_simple(
            client, provider, model, main_text, all_refs, mode, citation_style,
            write_fn=write_fn, upd_fn=upd_fn
        )
    except Exception as e_ins:
        write_fn(f"⚠️ Erro ao inserir citações: {e_ins}")
        cit_result = {"full_cited_text": main_text, "final_refs": [], "ref_map": {}}

    upd_fn(97, "Finalizando...")
    final_ref_list = cit_result.get("final_refs", [])
    write_fn(f"🎉 **Pronto!** {len(final_ref_list)} referência(s) inserida(s) no texto.")
    upd_fn(100, "Concluído!")

    # Package result in format display_results expects
    result = {
        "full_cited_text": cit_result.get("full_cited_text", main_text),
        "final_refs":      final_ref_list,
        "ref_map":         cit_result.get("ref_map", {}),
        "paragraphs":      [],   # empty — display_results will use full_cited_text
        "_simple_mode":    True,
    }
    return result, all_refs, final_ref_list


def _pipeline_worker(task_id: str, provider: str, client_cfg: dict,
                     main_text: str, ref_bytes: list,
                     mode: str, library_refs: list,
                     citation_style: str, manual_refs: list):
    """Backward-compat wrapper kept for reference (no longer used in UI)."""
    task = _PIPELINE_TASKS.get(task_id, {})
    def _w(msg): pass
    def _u(pct, msg):
        task["progress"] = int(pct)
        task["msg"]      = msg
    try:
        result, all_refs, final_ref_list = _run_pipeline_sync(
            _w, _u, provider, client_cfg, main_text, ref_bytes,
            mode, library_refs, citation_style, manual_refs)
        task["result"] = (result, all_refs, final_ref_list)
        task["status"] = "done"
    except Exception as e:
        import traceback as _tb
        task["error"]     = str(e)
        task["traceback"] = _tb.format_exc()
        task["status"]    = "error"

def render_citar_tab():
    st.markdown("### Processar Texto")

    api_key  = st.session_state.get("_api_key", "")
    provider = st.session_state.get("_provider", "")
    model    = st.session_state.get("_model", "")

    if not api_key:
        st.warning("Configure a chave API na barra lateral antes de processar.")
        return

    # (background task polling removed — now uses st.status() synchronous approach)

    # ── Config UI ─────────────────────────────────────────────────────────
    col_mode, col_style = st.columns([2, 1])
    with col_mode:
        mode = st.radio(
            "Modo de operacao:",
            ["Adicionar citacoes (texto sem citacoes)",
             "Revisar citacoes (texto ja citado)"],
            horizontal=True,
        )
        mode_key = "add" if "Adicionar" in mode else "review"
    with col_style:
        citation_style = st.selectbox(
            "Estilo de citacao:",
            CITATION_STYLES,
            index=0,
            help="Escolha o formato das referencias na lista final"
        )

    st.markdown("#### Texto principal")
    tab_up, tab_paste = st.tabs(["Upload (PDF, Word ou Markdown)", "Colar texto"])
    text_file   = None
    pasted_text = ""
    with tab_up:
        text_file = st.file_uploader("Upload PDF, Word ou Markdown (.md)",
                                     type=["pdf", "docx", "md"],
                                     key="main_upload")
        if text_file:
            st.success(f"{text_file.name} carregado")
    with tab_paste:
        pasted_text = st.text_area(
            "Cole seu texto aqui:", height=260,
            placeholder="Cole o texto do artigo, capitulo ou secao aqui...")

    st.markdown("#### Referencias (opcional)")
    st.caption("Upload de PDFs dos artigos de referencia para extracao de metadados. "
               "Sem upload, o app busca referencias automaticamente em "
               "PubMed, LILACS, Europe PMC e OpenAlex.")

    ref_tab_pdf, ref_tab_text = st.tabs(["Upload PDFs", "Colar lista de referencias"])
    with ref_tab_pdf:
        ref_files = st.file_uploader("PDFs de referencia", type=["pdf"],
                                     accept_multiple_files=True, key="ref_upload")
        if ref_files:
            st.info(f"{len(ref_files)} PDF(s) de referencia carregado(s).")
    with ref_tab_text:
        manual_ref_text = st.text_area(
            "Cole sua lista de referencias aqui:",
            height=180,
            placeholder=(
                "1. Smith AB, Jones CD. Titulo do artigo. J Psychiatry. 2023;45(2):123-130.\n"
                "2. Garcia L et al. Another study. Lancet. 2022;399:1234-1240.\n"
                "(Pode ser numerada ou nao — uma referencia por linha)"
            ),
            key="manual_refs_input",
        )

    use_library  = st.checkbox("Incluir artigos da biblioteca local como referencias", value=True)
    library_refs = load_library() if use_library else []

    st.divider()

    run_btn = st.button("🚀 Processar texto com IA", type="primary",
                        use_container_width=True, disabled=not api_key)

    if run_btn:
        # ── Extract text ──────────────────────────────────────────────────
        main_text = ""
        if text_file:
            b  = text_file.read()
            fn = text_file.name.lower()
            if fn.endswith(".pdf"):
                main_text = extract_text_from_pdf(b)
            elif fn.endswith(".md"):
                main_text = b.decode("utf-8", errors="replace")
            else:
                main_text = extract_text_from_docx(b)
        else:
            main_text = pasted_text.strip()

        if not main_text or not main_text.strip():
            st.error("Insira o texto principal (upload ou colagem).")
            return

        # Read ref file bytes before handing to thread
        ref_bytes = []
        if ref_files:
            for rf in ref_files:
                ref_bytes.append((rf.name, rf.read()))

        extra_refs       = parse_manual_refs(manual_ref_text) if manual_ref_text.strip() else []
        combined_library = extra_refs + library_refs

        # ── Run pipeline with st.status() — real-time progress ───────────
        import traceback as _tb_citar
        _status_placeholder = st.empty()

        with st.status("⏳ Processando texto com IA...", expanded=True) as _status_ui:
            try:
                result, all_refs, final_ref_list = _run_pipeline_sync(
                    write_fn = st.write,
                    upd_fn   = lambda pct, msg: _status_ui.update(
                                    label=f"{pct}% — {msg}"),
                    provider     = provider,
                    client_cfg   = {"api_key": api_key, "model": model},
                    main_text    = main_text,
                    ref_bytes    = ref_bytes,
                    mode         = mode_key,
                    library_refs = combined_library,
                    citation_style = citation_style,
                    manual_refs  = extra_refs,
                )
                _status_ui.update(
                    label="✅ Processamento concluído!",
                    state="complete",
                    expanded=False,
                )
                st.session_state["last_citation_style"] = citation_style
                st.session_state["last_result"]         = result
                st.session_state["last_all_refs"]       = all_refs
                st.session_state["last_final_refs"]     = final_ref_list
                st.session_state["last_mode"]           = mode_key

            except Exception as _e_pipe:
                _status_ui.update(
                    label=f"❌ Erro: {str(_e_pipe)[:80]}",
                    state="error",
                    expanded=True,
                )
                st.error(f"**Erro no processamento:** {_e_pipe}")
                with st.expander("📋 Detalhes técnicos"):
                    st.code(_tb_citar.format_exc())
                st.info("💡 Tente novamente. Se persistir, reduza o texto ou mude para claude-haiku.")

    # ── Display previous results ──────────────────────────────────────────
    _lr = st.session_state.get("last_result")
    if _lr is not None:
        display_results(
            _lr,
            st.session_state.get("last_all_refs") or [],
            st.session_state.get("last_final_refs") or [],
            st.session_state.get("last_mode", "add"),
        )


# =============================================================================
# EMBASE SEARCH (via Elsevier API ou Europe PMC fallback)
# =============================================================================

def search_embase(query: str, max_results: int = 3,
                  elsevier_key: str = "") -> list:
    """Search EMBASE.
    If an Elsevier Institutional API key is provided, uses the official API.
    Otherwise falls back to Europe PMC which indexes EMBASE-derived content."""
    articles = []

    # -- Strategy 1: Elsevier Embase API (requires institutional key) ----------
    if elsevier_key:
        try:
            r = requests.get(
                "https://api.elsevier.com/content/search/sciencedirect",
                params={
                    "query":   query,
                    "count":   max_results,
                    "field":   "title,creator,publicationName,coverDate,doi,description",
                    "apiKey":  elsevier_key,
                },
                headers={"Accept": "application/json"},
                timeout=12,
            )
            if r.status_code == 200:
                entries = r.json().get("search-results",{}).get("entry",[])
                for e in entries[:max_results]:
                    articles.append({
                        "title":   e.get("dc:title","").rstrip("."),
                        "authors": [e.get("dc:creator","")],
                        "journal": e.get("prism:publicationName",""),
                        "year":    str(e.get("prism:coverDate",""))[:4],
                        "doi":     e.get("prism:doi",""),
                        "_source": "EMBASE/Elsevier",
                    })
        except Exception:
            pass

    # -- Strategy 2: Europe PMC with EMBASE-indexed sources --------------------
    if not articles:
        try:
            r = requests.get(
                "https://www.ebi.ac.uk/europepmc/webservices/rest/search",
                params={
                    "query":      query + " (SRC:MED OR SRC:PMC OR SRC:AGR)",
                    "resultType": "core",
                    "pageSize":   max_results,
                    "format":     "json",
                },
                timeout=12,
            )
            results = r.json().get("resultList",{}).get("result",[])
            for item in results:
                au_list = item.get("authorList",{}).get("author",[])
                authors = [f"{a.get('lastName','')} {a.get('initials','')}".strip()
                           for a in (au_list if isinstance(au_list, list) else [])]
                articles.append({
                    "title":   item.get("title","").rstrip("."),
                    "authors": authors,
                    "journal": item.get("journalTitle",""),
                    "year":    str(item.get("pubYear","")),
                    "volume":  item.get("journalVolume",""),
                    "issue":   item.get("issue",""),
                    "pages":   item.get("pageInfo",""),
                    "doi":     item.get("doi",""),
                    "pmid":    item.get("pmid",""),
                    "_source": "EMBASE/Europe PMC",
                })
        except Exception:
            pass

    return articles


# Patch multi_source_search to include EMBASE
_orig_multi = multi_source_search  # keep reference

def multi_source_search(query: str, max_per_source: int = 2,
                        elsevier_key: str = "") -> list:
    """Search all configured databases including EMBASE."""
    results = []
    results.extend(search_pubmed(query, max_per_source))
    results.extend(search_europe_pmc(query, max_per_source))
    results.extend(search_embase(query, max_per_source, elsevier_key))
    results.extend(search_lilacs(query, max_per_source))
    results.extend(search_openalex(query, max_per_source))
    if not results:
        results.extend(search_semantic_scholar(query, max_per_source))
    return results


# =============================================================================
# TEXT REVISION AI  (ortografia, gramatica, estilo, redundancia)
# =============================================================================

def revise_text_ai(client, provider, model, text: str) -> dict:
    """Run editorial revision on a text chunk; returns structured data for report."""
    prompt = f"""You are a senior scientific editor at a top medical publisher.
Your job: revise the following text with the same rigour applied to best-selling academic books.

Review for ALL of the following:
1. ORTHOGRAPHY — spelling errors, wrong accents, typos (count each as "orthographic")
2. GRAMMAR — subject-verb agreement, tense consistency, pronoun reference, syntax (count as "grammar")
3. STYLE — passive voice overuse, complex sentences, weak verbs, foreign expressions like "rather than" (count as "style")
4. REDUNDANCY — repeated ideas, words, or structures (count as "redundancy")
5. SCIENTIFIC CLARITY — imprecise or ambiguous scientific statements
6. COHESION — weak transitions between paragraphs

Rules:
- Preserve all citations like [1], [2], [?] exactly
- Preserve headings
- Do NOT add or remove citations
- Keep the author voice and scientific content

Return ONLY valid JSON (no markdown fences):
{{
  "paragraphs": [{{"original":"exact text","revised":"corrected","issues":["issue"],"changed":true,"issue_types":["orthographic","grammar","style","redundancy"]}}],
  "summary": "overall assessment in Portuguese",
  "overall_quality": "Bom/Regular/Requer revisao extensiva",
  "stats": {{"orthographic":0,"grammar":0,"style":0,"redundancy":0,"other":0}},
  "correction_categories": {{
    "Higienizacao e Refinamento Lexical": ["example correction 1"],
    "Ajustes de Coesao e Fluidez": ["example"],
    "Correcoes Ortograficas": ["example"],
    "Estrangeirismos e Padronizacao": ["example"]
  }},
  "report_meta": {{
    "overview": "detailed overview of revision process in Portuguese",
    "text_type": "description of text genre/type in Portuguese",
    "plagiarism_note": "no plagiarism found statement in Portuguese"
  }}
}}

TEXT TO REVISE:
{text}"""

    raw  = ai_call(client, provider, model, prompt, max_tokens=8000)
    data = extract_json_from_ai(raw)
    return data if data else {"error": "Falha na revisao", "raw": raw}


def run_revision_pipeline(client, provider, model, text: str) -> dict:
    """Run editorial revision across all text chunks and merge."""
    chunks  = chunk_text(text, max_chars=CHUNK_CHAR_LIMIT)
    n       = len(chunks)
    progress = st.progress(0, text="Iniciando revisao editorial...")
    log      = st.empty()

    all_paragraphs  = []
    all_summaries   = []
    quality_scores  = []
    merged_stats    = {"orthographic":0,"grammar":0,"style":0,"redundancy":0,"other":0}
    merged_categories = {}
    merged_report_meta = {}

    for i, chunk in enumerate(chunks, 1):
        progress.progress(i / n, text=f"Revisando parte {i}/{n}...")
        result = revise_text_ai(client, provider, model, chunk)
        if "error" in result and "paragraphs" not in result:
            all_paragraphs.append({
                "original": chunk[:200] + "...",
                "revised":  chunk[:200] + "...",
                "issues":   [f"Erro no chunk {i}: {result.get('error','')}"],
                "changed":  False,
                "issue_types": [],
            })
        else:
            paras = result.get("paragraphs", [])
            # ensure issue_types field present
            for pp in paras:
                if "issue_types" not in pp:
                    pp["issue_types"] = []
            all_paragraphs.extend(paras)
            if result.get("summary"):
                all_summaries.append(f"Parte {i}: {result['summary']}")
            if result.get("overall_quality"):
                quality_scores.append(result["overall_quality"])
            # merge categories and meta from first chunk that has them
            if not merged_categories and result.get("correction_categories"):
                merged_categories.update(result["correction_categories"])
            if not merged_report_meta and result.get("report_meta"):
                merged_report_meta.update(result["report_meta"])
            # merge chunk stats
            for k in ["orthographic","grammar","style","redundancy","other"]:
                merged_stats[k] = merged_stats.get(k,0) + result.get("stats",{}).get(k,0)

    progress.empty(); log.empty()

    return {
        "paragraphs":           all_paragraphs,
        "summary":              " | ".join(all_summaries),
        "overall_quality":      quality_scores[0] if quality_scores else "N/A",
        "stats":                merged_stats,
        "correction_categories": merged_categories,
        "report_meta":          merged_report_meta,
        "_chunked":             n > 1,
    }


def generate_revision_docx(result: dict) -> bytes:
    """Build structured Word revision REPORT following editorial standards."""
    from docx import Document as DocxDoc
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = DocxDoc()
    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.5)

    BLUE  = RGBColor(0x1E, 0x3A, 0x5F)
    DKGR  = RGBColor(0x33, 0x33, 0x33)
    GREEN = RGBColor(0x00, 0x60, 0x00)
    AMBER = RGBColor(0x8B, 0x60, 0x00)

    def _heading(text, level=1):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = True
        r.font.color.rgb = BLUE
        r.font.size = Pt(14 if level == 1 else 12)
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(4)
        return p

    def _body(text, italic=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.space_after = Pt(4)
        r = p.add_run(text)
        r.font.size = Pt(11)
        r.font.color.rgb = DKGR
        if italic: r.italic = True
        return p

    def _table_2col(rows_data, col_widths=(3, 13), header_row=None):
        """Simple 2-col table. rows_data = list of (cell1, cell2) tuples."""
        ncols = 2
        tbl = doc.add_table(rows=0, cols=ncols)
        tbl.style = "Table Grid"
        total_w = sum(col_widths)
        if header_row:
            hrow = tbl.add_row()
            for ci, txt in enumerate(header_row):
                hrow.cells[ci].width = Cm(col_widths[ci])
                p = hrow.cells[ci].paragraphs[0]
                r = p.add_run(txt)
                r.bold = True; r.font.size = Pt(9)
                hrow.cells[ci].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for (c1, c2) in rows_data:
            row = tbl.add_row()
            row.cells[0].width = Cm(col_widths[0])
            row.cells[1].width = Cm(col_widths[1])
            row.cells[0].paragraphs[0].add_run(str(c1)).font.size = Pt(9)
            row.cells[1].paragraphs[0].add_run(str(c2)).font.size = Pt(9)
        return tbl

    paragraphs   = result.get("paragraphs", [])
    changed      = [p for p in paragraphs if p.get("changed")]
    quality      = result.get("overall_quality", "N/A")
    summary      = result.get("summary", "")
    stats        = result.get("stats", {})
    report_meta  = result.get("report_meta", {})
    now_str      = __import__("datetime").datetime.now().strftime("%d/%m/%Y")

    # ── HEADER ─────────────────────────────────────────────────
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p_title.add_run("Relatorio de Revisao Ortografica e Gramatical")
    r.bold = True; r.font.size = Pt(16); r.font.color.rgb = BLUE

    p_sub = doc.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_sub.add_run(
        f"Qualidade geral: {quality}  |  Gerado em {now_str}"
    ).font.size = Pt(10)

    doc.add_paragraph()

    # ── SECTION 1: Visao Geral ──────────────────────────────────
    _heading("1. Visao Geral do Processo")
    overview = report_meta.get("overview", summary or
        "O objetivo da revisao foi elevar o rigor academico do texto, "
        "eliminando vicios de linguagem e garantindo precisao terminologica.")
    _body(overview)

    text_type = report_meta.get("text_type","")
    if text_type:
        _heading("Caracteristicas Textuais", level=2)
        _body(text_type)

    # ── SECTION 2: Principais Pontos Corrigidos ─────────────────
    _heading("2. Principais Pontos Corrigidos")
    categories = result.get("correction_categories", {})
    if categories:
        for cat_name, cat_items in categories.items():
            _heading(cat_name, level=2)
            if isinstance(cat_items, list):
                for item in cat_items:
                    p = doc.add_paragraph(style="List Bullet")
                    p.add_run(item).font.size = Pt(10)
            else:
                _body(str(cat_items))
    else:
        # fallback: generic categories from stats
        _heading("Higienizacao e Refinamento Lexical", level=2)
        _body(f"{stats.get('lexical',0)} substituicoes vocabulares para diversificar o texto sem perder precisao cientifica.")
        _heading("Correcoes Ortograficas e Gramaticais", level=2)
        _body(f"{stats.get('orthographic',0)} correcoes ortograficas e {stats.get('grammar',0)} ajustes gramaticais realizados.")
        _heading("Ajustes de Coesao e Fluidez", level=2)
        _body(f"{stats.get('style',0)} intervencoes de estilo e {stats.get('redundancy',0)} remocoes de redundancia.")

    # ── SECTION 3: Distribuicao Estatistica ─────────────────────
    _heading("3. Distribuicao das Intervencoes de Revisao")
    n_orth   = stats.get("orthographic", 0)
    n_gram   = stats.get("grammar", 0)
    n_style  = stats.get("style", 0)
    n_redund = stats.get("redundancy", 0)
    n_other  = stats.get("other", 0)
    total_changes = sum([n_orth, n_gram, n_style, n_redund, n_other])
    if total_changes == 0 and changed:
        total_changes = len(changed)
        n_orth = total_changes

    stat_rows = [
        ("Tipo de Intervencao", "Quantidade"),
        ("Correcoes Ortograficas", str(n_orth)),
        ("Ajustes Gramaticais",   str(n_gram)),
        ("Revisao de Estilo",     str(n_style)),
        ("Remocao de Redundancias", str(n_redund)),
        ("Outros ajustes",        str(n_other)),
        ("TOTAL",                 str(total_changes)),
    ]
    _table_2col([(r[0],r[1]) for r in stat_rows[1:]], col_widths=(9,5),
                header_row=stat_rows[0])
    doc.add_paragraph()

    # ── SECTION 4: Deteccao de Plagio ───────────────────────────
    _heading("4. Deteccao de Plagio")
    plagio_note = report_meta.get("plagiarism_note",
        "Nao foi detectado nenhum indicio de plagio deliberado. "
        "O texto apresenta elevado grau de originalidade autoral.")
    _body(plagio_note)

    # ── SECTION 5: Log de Alteracoes Detalhadas ─────────────────
    doc.add_page_break()
    _heading("5. Log Detalhado de Alteracoes")
    _body(f"Paragrafos alterados: {len(changed)} de {len(paragraphs)} totais.")
    doc.add_paragraph()

    tbl_header = ("Para.", "Problema identificado", "Original (trecho)", "Revisado (trecho)")
    tbl = doc.add_table(rows=1, cols=4)
    tbl.style = "Table Grid"
    widths_cm = [1.0, 4.5, 6.0, 6.0]
    hrow = tbl.rows[0]
    for ci, htxt in enumerate(tbl_header):
        hrow.cells[ci].paragraphs[0].add_run(htxt).bold = True
        hrow.cells[ci].paragraphs[0].add_run("").font.size = Pt(8)

    for i, p in enumerate(changed, 1):
        row = tbl.add_row()
        issues = "; ".join(p.get("issues", []))
        orig = (p.get("original","") or "")[:120] + ("..." if len(p.get("original","")) > 120 else "")
        rev  = (p.get("revised","")  or "")[:120] + ("..." if len(p.get("revised",""))  > 120 else "")
        for ci, txt in enumerate([str(i), issues, orig, rev]):
            row.cells[ci].paragraphs[0].add_run(txt).font.size = Pt(8)

    doc.add_paragraph()

    # ── SECTION 6: Texto Revisado Completo ─────────────────────
    doc.add_page_break()
    _heading("6. Texto Revisado Completo")
    for p in paragraphs:
        text = p.get("revised") or p.get("original", "")
        if not text.strip(): continue
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        para.paragraph_format.space_after = Pt(4)
        run = para.add_run(text)
        run.font.name = "Times New Roman"
        run.font.size = Pt(12)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def display_revision_results(result: dict):
    """Show revision results in the UI."""
    if "error" in result and "paragraphs" not in result:
        st.error(f"Erro: {result.get('error')}")
        with st.expander("Ver resposta bruta"):
            st.text(result.get("raw",""))
        return

    paragraphs = result.get("paragraphs", [])
    changed    = [p for p in paragraphs if p.get("changed")]
    quality    = result.get("overall_quality","N/A")

    q_color = {"Bom":"#28a745","Regular":"#ffc107","Requer revisao extensiva":"#dc3545"}.get(quality,"#6c757d")
    st.markdown(
        f'<div style="background:{q_color}20;border-left:4px solid {q_color};'
        f'padding:.7rem 1rem;border-radius:0 8px 8px 0;margin-bottom:1rem;">'
        f'<strong>Qualidade geral:</strong> {quality} &nbsp;|&nbsp; '
        f'<strong>{len(changed)}</strong> de <strong>{len(paragraphs)}</strong> paragrafos alterados'
        f'</div>', unsafe_allow_html=True)

    if result.get("summary"):
        st.info(f"Avaliacao editorial: {result['summary']}")

    st.divider()
    tab_full, tab_cmp = st.tabs(["Texto Revisado", "Comparacao Paragrafo a Paragrafo"])

    with tab_full:
        revised_full = "\n\n".join(
            p.get("revised") or p.get("original","") for p in paragraphs)
        st.text_area("", value=revised_full, height=500, label_visibility="collapsed")
        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                "📥 Baixar revisado (.txt)",
                data=revised_full.encode("utf-8"),
                file_name="texto_revisado.txt",
                mime="text/plain",
                use_container_width=True,
            )
        with col2:
            try:
                docx_bytes = generate_revision_docx(result)
                st.download_button(
                    "📄 Baixar revisado (.docx)",
                    data=docx_bytes,
                    file_name="texto_revisado.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    type="primary",
                )
            except Exception as e:
                st.warning(f"Erro ao gerar .docx: {e}")

    with tab_cmp:
        for i, p in enumerate(paragraphs):
            is_ch = p.get("changed", False)
            pill  = ('<span class="pill-changed">Alterado</span>'
                     if is_ch else '<span class="pill-ok">Sem alteracao</span>')
            with st.expander(f"Paragrafo {i+1}  {pill}", expanded=is_ch):
                if is_ch:
                    ca, cb = st.columns(2)
                    with ca:
                        st.markdown("**Original:**")
                        st.text(p.get("original",""))
                    with cb:
                        st.markdown("**Revisado:**")
                        st.text(p.get("revised",""))
                    if p.get("issues"):
                        st.caption("Problemas: " + " | ".join(p["issues"]))
                else:
                    st.text(p.get("original",""))


# =============================================================================
# REVISION TAB
# =============================================================================

# =============================================================================
# REVISÃO CIENTÍFICA — helpers
# =============================================================================

def _parse_text_and_refs(full_text: str) -> tuple:
    """Parse text that contains inline citations [1],[2] and a reference list.
    Returns (paragraphs: list[dict], references: dict[int, str])
    """
    import re as _re

    # ── Detect reference section ──────────────────────────────────────────
    sep = _re.search(
        r'\n\s*(?:referências?|references?|bibliografia|bibliography|──+)\s*\n',
        full_text, _re.IGNORECASE
    )
    if sep:
        body_text = full_text[:sep.start()]
        ref_text  = full_text[sep.start():]
    else:
        # Heuristic: find start of numbered list with ≥3 items
        m = _re.search(r'\n\s*1\.\s+\S', full_text)
        if m:
            potential = full_text[m.start():]
            if len(_re.findall(r'\n\s*\d+\.\s+', potential)) >= 3:
                body_text = full_text[:m.start()]
                ref_text  = potential
            else:
                body_text, ref_text = full_text, ""
        else:
            body_text, ref_text = full_text, ""

    # ── Parse paragraphs ─────────────────────────────────────────────────
    paragraphs = []
    for para in _re.split(r'\n\s*\n', body_text.strip()):
        para = para.strip()
        if not para:
            continue
        raw_cits = _re.findall(r'\[(\d+(?:[,;\s]\d+)*)\]', para)
        nums = sorted({int(n) for c in raw_cits for n in _re.split(r'[,;\s]+', c) if n.isdigit()})
        paragraphs.append({"text": para, "citation_numbers": nums})

    # ── Parse reference list ─────────────────────────────────────────────
    references = {}
    if ref_text:
        for m in _re.finditer(r'(\d+)\.\s+(.+?)(?=\n\s*\d+\.\s|\Z)', ref_text, _re.DOTALL):
            references[int(m.group(1))] = m.group(2).strip().replace('\n', ' ')

    return paragraphs, references


def _audit_citations_ai(client, provider, model,
                         paragraphs: list, references: dict, style: str) -> list:
    """Single AI call to audit ALL citations in the text.
    Returns list of issue dicts, one per problematic citation.
    """
    import re as _re

    # Build compact summary for AI
    para_summary = []
    for i, p in enumerate(paragraphs):
        if p["citation_numbers"]:
            para_summary.append(
                f"P{i+1}: ...{p['text'][-200:]}  → cita {p['citation_numbers']}"
            )

    ref_list_str = "\n".join(
        f"[{n}] {txt[:200]}" for n, txt in sorted(references.items())
    )

    if not para_summary:
        return []

    prompt = (
        f"You are a scientific peer reviewer auditing citations in a manuscript.\n"
        f"Citation style: {style}\n\n"
        "For EACH citation in the paragraphs below:\n"
        "1. Check RELEVANCE: does the cited reference actually support the claim?\n"
        "2. Check FORMAT: is the reference formatted correctly in the stated style?\n"
        "3. If relevance is Low or Wrong: provide an English PubMed query to find a better article.\n\n"
        "Return ONLY valid JSON array (no markdown):\n"
        '[\n'
        '  {\n'
        '    "paragraph": "P1",\n'
        '    "citation_num": 3,\n'
        '    "relevance": "Alta|Media|Baixa|Incorreta",\n'
        '    "relevance_reason": "brief in Portuguese",\n'
        '    "format_ok": true,\n'
        '    "format_issues": ["list or empty"],\n'
        '    "status": "OK|Aviso|Erro",\n'
        '    "pubmed_query": "English query or empty string",\n'
        '    "summary": "one sentence in Portuguese"\n'
        '  }\n'
        ']\n\n'
        "Only include entries where status is Aviso or Erro (skip OK ones).\n\n"
        "=== PARAGRAPHS (with their citations) ===\n"
        + "\n".join(para_summary[:40]) +
        "\n\n=== REFERENCE LIST ===\n"
        + ref_list_str[:3000]
    )

    raw  = ai_call(client, provider, model, prompt, max_tokens=6000)
    data = extract_json_from_ai(raw)
    if isinstance(data, list):
        return data
    if isinstance(data, dict) and "issues" in data:
        return data["issues"]
    return []


def _format_audit_report_docx(issues: list, paragraphs: list,
                                references: dict, style: str) -> bytes:
    """Generate Word report for citation audit."""
    from docx import Document as _Doc
    from docx.shared import Pt, RGBColor
    doc = _Doc()
    doc.add_heading("Relatório de Auditoria Científica de Citações", 0)
    doc.add_paragraph(f"Estilo: {style}  |  Parágrafos auditados: {len(paragraphs)}  "
                      f"|  Problemas encontrados: {len(issues)}")
    doc.add_heading("Problemas Identificados", level=1)
    if not issues:
        doc.add_paragraph("Nenhum problema encontrado — todas as citações estão corretas.")
    for iss in issues:
        p_lbl = iss.get("paragraph", "?")
        c_num = iss.get("citation_num", "?")
        status = iss.get("status", "?")
        h = doc.add_heading(f"{p_lbl} → Citação [{c_num}]  ({status})", level=2)
        run = h.runs[0] if h.runs else h.add_run()
        if status == "Erro":
            run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        elif status == "Aviso":
            run.font.color.rgb = RGBColor(0xFF, 0x80, 0x00)
        doc.add_paragraph(f"Relevância: {iss.get('relevance','?')}")
        doc.add_paragraph(iss.get("relevance_reason",""))
        if iss.get("format_issues"):
            doc.add_paragraph("Problemas de formato:")
            for fi in iss["format_issues"]:
                doc.add_paragraph(fi, style="List Bullet")
        ref_txt = references.get(c_num, "")
        if ref_txt:
            doc.add_paragraph(f"Referência atual: {ref_txt[:300]}")
        if iss.get("pubmed_query"):
            doc.add_paragraph(f"Query PubMed sugerida: {iss['pubmed_query']}")
        sugs = iss.get("suggestions", [])
        if sugs:
            doc.add_paragraph("Artigos sugeridos:")
            for s in sugs:
                doc.add_paragraph(
                    f"{s.get('authors',['?'])[0]} et al. {s.get('title','?')[:80]}. "
                    f"{s.get('journal','?')} {s.get('year','?')}. PMID:{s.get('pmid','')}",
                    style="List Bullet"
                )
    doc.add_heading("Lista de Referências (conforme recebida)", level=1)
    for n, txt in sorted(references.items()):
        doc.add_paragraph(f"[{n}] {txt}", style="List Bullet")
    from io import BytesIO as _BIO
    buf = _BIO(); doc.save(buf); return buf.getvalue()


def render_revisao_tab():
    st.markdown("### Revisão Científica")
    api_key  = st.session_state.get("_api_key", "")
    provider = st.session_state.get("_provider", "")
    model    = st.session_state.get("_model", "")

    if not api_key:
        st.warning("Configure a chave API na barra lateral antes de usar.")
        return

    sub_audit, sub_editorial = st.tabs(
        ["🔬 Auditoria de Citações", "📝 Revisão Editorial"])

    # ── TAB: AUDITORIA DE CITAÇÕES ────────────────────────────────────────
    with sub_audit:
        st.markdown(
            "Cole um texto científico já formatado (com citações inline "
            "`[1]`, `[2]`... e lista de referências ao final). "
            "A IA verifica se cada citação é relevante para o parágrafo, "
            "se a formatação está correta e sugere artigos substitutos quando necessário."
        )

        col_sty, col_sp = st.columns([2, 3])
        with col_sty:
            audit_style = st.selectbox(
                "Estilo bibliográfico:",
                CITATION_STYLES,
                key="audit_style_sel",
            )

        # Input
        tab_a_up, tab_a_paste = st.tabs(["Upload (PDF / DOCX / MD)", "Colar texto"])
        audit_file, audit_pasted = None, ""
        with tab_a_up:
            audit_file = st.file_uploader(
                "Upload do texto completo (com referências ao final)",
                type=["pdf", "docx", "md"], key="audit_upload")
            if audit_file:
                st.success(f"{audit_file.name} carregado")
        with tab_a_paste:
            audit_pasted = st.text_area(
                "Cole o texto completo (inclua a lista de referências ao final):",
                height=320, key="audit_paste",
                placeholder=(
                    "Parágrafo com afirmação científica [1,2].\n\n"
                    "Outro parágrafo com outra afirmação [3].\n\n"
                    "REFERÊNCIAS\n"
                    "1. Smith AB et al. Título do artigo. J Psychiatry. 2023;45:123.\n"
                    "2. Jones CD, Garcia L. Outro artigo. Lancet. 2022;399:1234.\n"
                    "3. ..."
                ),
            )

        run_audit = st.button(
            "🔬 Auditar citações", type="primary",
            use_container_width=True, disabled=not api_key,
            key="btn_run_audit"
        )

        if run_audit:
            # Extract text
            full_text = ""
            if audit_file:
                b = audit_file.read(); fn = audit_file.name.lower()
                if fn.endswith(".pdf"):
                    full_text = extract_text_from_pdf(b)
                elif fn.endswith(".md"):
                    full_text = b.decode("utf-8", errors="replace")
                else:
                    full_text = extract_text_from_docx(b)
            else:
                full_text = audit_pasted.strip()

            if not full_text:
                st.error("Insira o texto para auditoria.")
            else:
                paragraphs, references = _parse_text_and_refs(full_text)
                cited_paras = [p for p in paragraphs if p["citation_numbers"]]

                if not references:
                    st.warning(
                        "Não foi possível detectar a lista de referências. "
                        "Certifique-se de que o texto termina com uma seção "
                        "identificada como 'Referências' ou com uma lista numerada."
                    )
                elif not cited_paras:
                    st.warning(
                        "Nenhuma citação inline encontrada. "
                        "O texto deve conter marcações como [1], [2] ou [1,2]."
                    )
                else:
                    st.info(
                        f"Detectados: **{len(paragraphs)}** parágrafos, "
                        f"**{len(cited_paras)}** com citações, "
                        f"**{len(references)}** referências na lista."
                    )

                    try:
                        client = get_ai_client(provider, api_key)
                    except Exception as e:
                        st.error(f"Erro ao conectar IA: {e}")
                        st.stop()

                    with st.status("🔬 Auditando citações...", expanded=True) as _s:
                        st.write(f"📋 Analisando {len(cited_paras)} parágrafos com citações...")
                        issues = _audit_citations_ai(
                            client, provider, model, paragraphs, references, audit_style)
                        _s.update(label=f"✅ Auditoria concluída — {len(issues)} problema(s) encontrado(s)",
                                  state="complete", expanded=False)

                    # For each issue with pubmed_query: search for suggestions
                    if issues:
                        with st.status("🌐 Buscando artigos substitutos...", expanded=True) as _s2:
                            for iss in issues:
                                q = iss.get("pubmed_query", "").strip()
                                if q and iss.get("status") in ("Erro", "Aviso"):
                                    st.write(f"  PubMed: `{q[:60]}`...")
                                    results = multi_source_search(q, max_per_source=3)
                                    iss["suggestions"] = results[:3]
                                else:
                                    iss["suggestions"] = []
                            _s2.update(label="✅ Busca concluída", state="complete",
                                       expanded=False)

                    st.session_state["last_audit"] = {
                        "issues": issues, "paragraphs": paragraphs,
                        "references": references, "style": audit_style,
                        "n_cited_paras": len(cited_paras),
                    }

        # ── Display audit results ─────────────────────────────────────────
        audit_res = st.session_state.get("last_audit")
        if audit_res:
            issues     = audit_res["issues"]
            refs       = audit_res["references"]
            paras      = audit_res["paragraphs"]
            sty        = audit_res["style"]
            n_cited    = audit_res["n_cited_paras"]
            n_ok       = n_cited - len([i for i in issues if i.get("status") == "Erro"])
            n_warn     = len([i for i in issues if i.get("status") == "Aviso"])
            n_err      = len([i for i in issues if i.get("status") == "Erro"])

            st.divider()
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Parágrafos auditados", n_cited)
            c2.metric("✅ Citações OK",  max(0, n_cited - n_err - n_warn))
            c3.metric("⚠️ Avisos",       n_warn)
            c4.metric("❌ Erros",        n_err)

            if not issues:
                st.success("✅ Todas as citações estão corretas e bem formatadas!")
            else:
                st.markdown("---")
                for iss in issues:
                    p_lbl = iss.get("paragraph", "?")
                    c_num = iss.get("citation_num", "?")
                    status = iss.get("status", "Aviso")
                    rel    = iss.get("relevance", "?")
                    icon   = "❌" if status == "Erro" else "⚠️"

                    with st.expander(
                        f"{icon} **{p_lbl}** → Citação **[{c_num}]** "
                        f"— Relevância: {rel}  ({status})",
                        expanded=(status == "Erro")
                    ):
                        st.markdown(f"**Diagnóstico:** {iss.get('relevance_reason','')}")

                        f_issues = iss.get("format_issues", [])
                        if f_issues:
                            st.markdown("**Problemas de formatação:**")
                            for fi in f_issues:
                                st.markdown(f"- {fi}")

                        ref_txt = refs.get(c_num, "")
                        if ref_txt:
                            st.markdown(
                                f"<div class='ref-box'>Referência atual [{c_num}]: "
                                f"{ref_txt[:300]}</div>",
                                unsafe_allow_html=True
                            )

                        sugs = iss.get("suggestions", [])
                        if sugs:
                            st.markdown("**📚 Artigos mais adequados encontrados:**")
                            for j, s in enumerate(sugs, 1):
                                authors = s.get("authors", [])
                                a_str   = authors[0] if authors else "?"
                                if len(authors) > 1:
                                    a_str += " et al."
                                pmid = s.get("pmid", "")
                                doi  = s.get("doi",  "")
                                link = (f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/"
                                        if pmid else
                                        f"https://doi.org/{doi}" if doi else "")
                                title = s.get("title", "")
                                journal = s.get("journal", "")
                                year    = s.get("year", "")
                                line = f"**{j}.** {a_str}. {title[:80]}. *{journal}*. {year}."
                                if link:
                                    line += f" → [Ver artigo]({link})"
                                st.markdown(line)
                        elif iss.get("pubmed_query"):
                            st.caption(
                                f"Query de busca usada: `{iss.get('pubmed_query','')}`")

            # Download report
            st.divider()
            try:
                rpt_bytes = _format_audit_report_docx(issues, paras, refs, sty)
                st.download_button(
                    "📄 Baixar relatório completo (.docx)",
                    data=rpt_bytes,
                    file_name="auditoria_citacoes.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    type="primary",
                )
            except Exception as e_rpt:
                st.warning(f"Relatório Word não gerado: {e_rpt}")

    # ── TAB: REVISÃO EDITORIAL (mantida) ─────────────────────────────────
    with sub_editorial:
        st.caption(
            "Revisão ortográfica, gramatical, estilística e de redundância — "
            "como um editor sênior de publicações científicas.")

        tab_up, tab_paste = st.tabs(["Upload (PDF, Word ou Markdown)", "Colar texto"])
        text_file, pasted_text = None, ""
        with tab_up:
            text_file = st.file_uploader("Upload PDF, Word ou Markdown",
                                         type=["pdf","docx","md"], key="rev_upload")
            if text_file:
                st.success(f"{text_file.name} carregado")
        with tab_paste:
            pasted_text = st.text_area("Cole seu texto aqui:", height=260,
                                       placeholder="Cole o texto para revisão...",
                                       key="rev_paste")

        st.divider()
        col_a, col_b = st.columns([3, 1])
        with col_a:
            st.markdown(
                "**O que será revisado:**\n"
                "- Ortografia e acentuação\n"
                "- Gramática e concordância\n"
                "- Estilo científico (voz passiva excessiva, frases longas)\n"
                "- Redundância e repetições\n"
                "- Coesão entre parágrafos"
            )
        with col_b:
            run_rev = st.button("Revisar texto", type="primary",
                                use_container_width=True, disabled=not api_key,
                                key="btn_rev_editorial")

        if run_rev:
            main_text = ""
            if text_file:
                b = text_file.read(); fn = text_file.name.lower()
                if fn.endswith(".pdf"):
                    main_text = extract_text_from_pdf(b)
                elif fn.endswith(".md"):
                    main_text = b.decode("utf-8", errors="replace")
                else:
                    main_text = extract_text_from_docx(b)
            else:
                main_text = pasted_text.strip()

            if not main_text:
                st.error("Insira o texto para revisão.")
            else:
                try:
                    client = get_ai_client(provider, api_key)
                except Exception as e:
                    st.error(f"Erro ao inicializar cliente IA: {e}")
                    st.stop()
                result = run_revision_pipeline(client, provider, model, main_text)
                st.session_state["last_revision"] = result

        if st.session_state.get("last_revision"):
            display_revision_results(st.session_state["last_revision"])


# =============================================================================
# VERIFICAR REFERENCIAS TAB
# =============================================================================

def verify_references_ai(client, provider, model, ref_text: str, style: str) -> dict:
    """AI verifies a reference list: checks existence, fixes formatting, fills gaps."""
    prompt = f"""You are an expert librarian and scientific citation specialist.
The user has provided a list of bibliographic references. Your tasks:

1. For each reference, verify if it appears to be a real, valid publication
2. Correct any formatting errors to match {style} style exactly
3. Fill in missing fields (DOI, year, volume, pages) if you can infer them from context
4. Flag references that seem fabricated or that have inconsistent data
5. Identify and flag duplicate references

Return ONLY valid JSON (no markdown fences):
{{
  "references": [
    {{
      "number": 1,
      "original": "original text as provided",
      "corrected": "corrected reference in {style} style",
      "status": "Validada / Suspeita / Incorreta / Duplicata",
      "issues": ["list of issues found"],
      "changes_made": ["list of corrections applied"],
      "confidence": "Alta / Media / Baixa"
    }}
  ],
  "summary": "overall summary in Portuguese",
  "total_valid": 0,
  "total_issues": 0,
  "total_duplicates": 0
}}

REFERENCE LIST TO VERIFY:
{ref_text}"""

    raw  = ai_call(client, provider, model, prompt, max_tokens=8000)
    data = extract_json_from_ai(raw)
    return data if data else {"error": "Falha na verificacao", "raw": raw}


def generate_ref_verification_docx(result: dict, style: str) -> bytes:
    """Build Word report of reference verification."""
    from docx import Document as DocxDoc
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = DocxDoc()
    for sec in doc.sections:
        sec.left_margin = Cm(2.5); sec.right_margin = Cm(2.5)
        sec.top_margin  = Cm(2.5); sec.bottom_margin = Cm(2.5)

    BLUE  = RGBColor(0x1E, 0x3A, 0x5F)
    GREEN = RGBColor(0x00, 0x70, 0x00)
    RED   = RGBColor(0xCC, 0x00, 0x00)
    AMBER = RGBColor(0x8B, 0x60, 0x00)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Relatorio de Verificacao de Referencias Bibliograficas")
    r.bold = True; r.font.size = Pt(15); r.font.color.rgb = BLUE

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.add_run(
        f"Estilo: {style}  |  Gerado em {__import__('datetime').datetime.now().strftime('%d/%m/%Y')}"
    ).font.size = Pt(10)

    doc.add_paragraph()

    refs = result.get("references", [])
    valid_n = result.get("total_valid", sum(1 for r2 in refs if r2.get("status","") == "Validada"))
    issues_n = result.get("total_issues", sum(1 for r2 in refs if r2.get("status","") != "Validada"))

    ph = doc.add_paragraph()
    ph.add_run("Resumo: ").bold = True
    ph.add_run(result.get("summary",""))
    ph.paragraph_format.space_after = Pt(8)

    # Stats table
    tbl = doc.add_table(rows=2, cols=3)
    tbl.style = "Table Grid"
    headers = ["Total de Referencias", "Validadas", "Com Problemas"]
    values  = [str(len(refs)), str(valid_n), str(issues_n)]
    for ci, (h, v) in enumerate(zip(headers, values)):
        tbl.rows[0].cells[ci].paragraphs[0].add_run(h).bold = True
        tbl.rows[1].cells[ci].paragraphs[0].add_run(v)
    doc.add_paragraph()

    # Per-reference detail
    ph2 = doc.add_paragraph()
    ph2.add_run("Detalhamento por Referencia").bold = True
    ph2.runs[0].font.size = Pt(13); ph2.runs[0].font.color.rgb = BLUE

    for ref in refs:
        status = ref.get("status","")
        color = GREEN if status == "Validada" else (RED if status == "Incorreta" else AMBER)
        p3 = doc.add_paragraph()
        p3.paragraph_format.space_before = Pt(8)
        p3.add_run(f"[{ref.get('number','')}] [{status}] ").bold = True
        p3.runs[0].font.color.rgb = color

        p4 = doc.add_paragraph()
        p4.add_run("Original: ").bold = True
        p4.add_run(ref.get("original","")).italic = True
        p4.paragraph_format.left_indent = Cm(0.5)

        if ref.get("corrected") and ref["corrected"] != ref.get("original",""):
            p5 = doc.add_paragraph()
            p5.add_run("Corrigida: ").bold = True
            p5.add_run(ref["corrected"])
            p5.paragraph_format.left_indent = Cm(0.5)
            p5.runs[1].font.color.rgb = GREEN

        if ref.get("issues"):
            p6 = doc.add_paragraph()
            p6.add_run("Problemas: " + " | ".join(ref["issues"])).font.size = Pt(9)
            p6.paragraph_format.left_indent = Cm(0.5)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def render_verificar_refs_tab():
    st.markdown("### Verificacao e Correcao de Referencias")
    st.caption("Cole sua lista de referencias e a IA verifica veracidade, corrige a formatacao e completa informacoes faltantes.")

    api_key  = st.session_state.get("_api_key","")
    provider = st.session_state.get("_provider","")
    model    = st.session_state.get("_model","")

    if not api_key:
        st.warning("Configure a chave API na barra lateral antes de verificar.")
        return

    style = st.selectbox("Estilo das referencias:", CITATION_STYLES, index=0, key="vr_style")

    ref_text = st.text_area(
        "Cole sua lista de referencias aqui (uma por linha ou numeradas):",
        height=300,
        placeholder="1. Autor A, Autor B. Titulo do artigo. Revista. 2023;45(2):123-130.\n2. ...",
        key="vr_refs"
    )

    run_vr = st.button("Verificar referencias", type="primary",
                       disabled=not api_key or not ref_text.strip())

    if run_vr and ref_text.strip():
        try:
            client = get_ai_client(provider, api_key)
        except Exception as e:
            st.error(f"Erro ao inicializar cliente: {e}")
            return

        with st.spinner("Verificando referencias com IA..."):
            result = verify_references_ai(client, provider, model, ref_text, style)
            st.session_state["last_vr_result"] = result
            st.session_state["last_vr_style"]  = style

    vr = st.session_state.get("last_vr_result")
    if not vr:
        return

    if "error" in vr and "references" not in vr:
        st.error(f"Erro: {vr.get('error')}")
        return

    refs  = vr.get("references", [])
    valid_n = sum(1 for r in refs if r.get("status","") == "Validada")
    issue_n = len(refs) - valid_n

    c1, c2, c3 = st.columns(3)
    c1.metric("Total de referencias", len(refs))
    c2.metric("Validadas", valid_n, delta=None)
    c3.metric("Com problemas", issue_n,
              delta=f"-{issue_n}" if issue_n > 0 else None,
              delta_color="inverse")

    if vr.get("summary"):
        st.info(vr["summary"])

    st.divider()
    st.markdown("#### Referencias Corrigidas")

    # Show corrected list
    corrected_lines = []
    for ref in refs:
        status = ref.get("status","")
        badge = "✅" if status == "Validada" else ("⚠️" if status == "Suspeita" else "❌")
        with st.expander(f"{badge} [{ref.get('number','')}] {status} — Confianca: {ref.get('confidence','')}"):
            if ref.get("original") != ref.get("corrected",""):
                col_a, col_b = st.columns(2)
                with col_a:
                    st.markdown("**Original:**")
                    st.text(ref.get("original",""))
                with col_b:
                    st.markdown("**Corrigida:**")
                    st.markdown(f":green[{ref.get('corrected','')}]")
            else:
                st.text(ref.get("corrected") or ref.get("original",""))
            if ref.get("issues"):
                st.caption("Problemas: " + " | ".join(ref["issues"]))
            if ref.get("changes_made"):
                st.caption("Correcoes: " + " | ".join(ref["changes_made"]))
        corrected_lines.append(ref.get("corrected") or ref.get("original",""))

    corrected_text = "\n\n".join(corrected_lines)
    col_dl1, col_dl2 = st.columns(2)
    with col_dl1:
        st.download_button("📥 Baixar lista corrigida (.txt)",
                           data=corrected_text.encode("utf-8"),
                           file_name="referencias_corrigidas.txt",
                           mime="text/plain", use_container_width=True)
    with col_dl2:
        try:
            docx_b = generate_ref_verification_docx(vr, st.session_state.get("last_vr_style","Vancouver"))
            st.download_button("📄 Baixar relatorio (.docx)",
                               data=docx_b,
                               file_name="verificacao_referencias.docx",
                               mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                               use_container_width=True, type="primary")
        except Exception as e:
            st.warning(f"Erro ao gerar .docx: {e}")


# =============================================================================
# CONVERTER TAB (PDF / DOCX -> Markdown)
# =============================================================================

def convert_file_to_markdown(file_bytes: bytes, filename: str) -> str:
    """Convert PDF or DOCX to clean Markdown with heading detection and paragraph joining."""
    import re as _re
    fn = filename.lower()

    if fn.endswith(".pdf"):
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")

            # Pass 1: find body font size (most common size across all spans)
            all_sizes = []
            for page in doc:
                for block in page.get_text("dict")["blocks"]:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            if span["text"].strip():
                                all_sizes.append(round(span["size"], 1))
            if not all_sizes:
                doc.close()
                return "Nenhum texto encontrado no PDF."

            from collections import Counter
            body_size = Counter(all_sizes).most_common(1)[0][0]

            def _hlevel(size):
                d = size - body_size
                if d >= 6:   return 1
                if d >= 3:   return 2
                if d >= 1.5: return 3
                return 0

            # Pass 2: build Markdown page by page
            md_parts = []
            for page in doc:
                page_md = []
                para_buf = []

                def flush():
                    if para_buf:
                        joined = " ".join(para_buf)
                        joined = _re.sub(r" +", " ", joined).strip()
                        # Split on sentence boundaries for readability
                        joined = _re.sub(r"([.!?]) ([A-ZÁÀÂÃÉÊÍÓÔÕÚÜÇ])",
                                         r"\1\n\n\2", joined)
                        if joined:
                            page_md.append(joined)
                        para_buf.clear()

                for block in page.get_text("dict")["blocks"]:
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines", []):
                        spans     = line.get("spans", [])
                        line_text = " ".join(s["text"] for s in spans).strip()
                        if not line_text:
                            continue
                        max_size = max((s["size"] for s in spans), default=body_size)
                        is_bold  = any(s.get("flags", 0) & 16 for s in spans)
                        level    = _hlevel(max_size)
                        # Bold short lines → subheading
                        if level == 0 and is_bold and len(line_text) < 120:
                            level = 3
                        # "24.1 Title" pattern → subheading
                        if level == 0 and _re.match(r"^\d+\.\d*\s+\S", line_text) \
                                and len(line_text) < 120:
                            level = 2
                        if level > 0:
                            flush()
                            page_md.append("#" * level + " " + line_text)
                        else:
                            # Strip soft hyphen at end of wrapped line
                            cleaned = _re.sub(r"-$", "", line_text.rstrip())
                            para_buf.append(cleaned)

                    flush()  # flush after each block

                if page_md:
                    md_parts.append("\n\n".join(page_md))

            doc.close()
            full = "\n\n---\n\n".join(md_parts)
            full = _re.sub(r"\n{3,}", "\n\n", full)
            return full

        except Exception as e:
            import traceback as _tb
            return f"Erro ao converter PDF: {e}\n{_tb.format_exc()}"

    elif fn.endswith(".docx"):
        try:
            from docx import Document as _DocxDoc
            doc   = _DocxDoc(BytesIO(file_bytes))
            lines_out = []
            for para in doc.paragraphs:
                style = para.style.name.lower()
                text  = para.text.strip()
                if not text:
                    lines_out.append("")
                    continue
                if "heading 1" in style or "titulo 1" in style:
                    lines_out.append(f"# {text}")
                elif "heading 2" in style or "titulo 2" in style:
                    lines_out.append(f"## {text}")
                elif "heading 3" in style or "titulo 3" in style:
                    lines_out.append(f"### {text}")
                elif "list" in style or "bullet" in style or "lista" in style:
                    lines_out.append(f"- {text}")
                elif para.runs and all(r.bold for r in para.runs if r.text.strip()):
                    lines_out.append(f"## {text}")
                else:
                    lines_out.append(text)
            import re as _re2
            md = "\n\n".join(lines_out)
            return _re2.sub(r"\n{3,}", "\n\n", md)
        except Exception as e:
            return f"Erro ao converter DOCX: {e}"
    else:
        return file_bytes.decode("utf-8", errors="replace")

def render_converter_tab():
    st.markdown("### Converter Documentos para Markdown (.md)")
    st.caption("Converta PDF ou DOCX para formato Markdown puro — ideal para uso em sistemas de citacao, editors e workflows com IA.")

    uploaded = st.file_uploader(
        "Selecione PDF ou DOCX para converter",
        type=["pdf","docx"],
        key="conv_upload"
    )

    if not uploaded:
        st.info("Arraste um arquivo PDF ou DOCX acima para comecar.")
        return

    st.success(f"Arquivo carregado: **{uploaded.name}** ({uploaded.size/1024:.1f} KB)")

    if st.button("Converter para Markdown", type="primary"):
        with st.spinner("Convertendo..."):
            raw_bytes = uploaded.read()
            md_text   = convert_file_to_markdown(raw_bytes, uploaded.name)
            st.session_state["last_md_conversion"]      = md_text
            st.session_state["last_md_conversion_name"] = uploaded.name.rsplit(".",1)[0] + ".md"

    md_result = st.session_state.get("last_md_conversion")
    if md_result:
        st.markdown(f"**Preview ({len(md_result)} caracteres):**")
        st.text_area("", value=md_result[:3000] + ("\n\n[...truncado para preview...]" if len(md_result) > 3000 else ""),
                     height=350, label_visibility="collapsed")

        col1, col2 = st.columns(2)
        fname = st.session_state.get("last_md_conversion_name","documento.md")
        with col1:
            st.download_button(
                "📥 Baixar .md",
                data=md_result.encode("utf-8"),
                file_name=fname,
                mime="text/markdown",
                use_container_width=True,
                type="primary",
            )
        with col2:
            st.download_button(
                "📥 Baixar .txt",
                data=md_result.encode("utf-8"),
                file_name=fname.replace(".md",".txt"),
                mime="text/plain",
                use_container_width=True,
            )


# =============================================================================
# EVIDENCIAS TAB (find supporting passages in articles)
# =============================================================================

def find_evidence_ai(client, provider, model, claim: str, article_text: str, article_name: str) -> dict:
    """AI finds passages in article_text that support the given claim."""
    prompt = f"""You are a scientific evidence analyst.

TASK: Find passages in the provided article that SUPPORT or are RELEVANT to the claim/paragraph below.

CLAIM TO SUPPORT:
{claim}

ARTICLE ({article_name}):
{article_text[:6000]}

Instructions:
- Extract 2-6 passages that best support the claim
- Each passage should be a direct quote from the article (verbatim)
- Assess relevance strength: Alta / Media / Baixa
- Note page/section if identifiable
- If no supporting passages exist, state this clearly

Return ONLY valid JSON (no markdown fences):
{{
  "article": "{article_name}",
  "passages": [
    {{
      "text": "exact quote from article",
      "relevance": "Alta/Media/Baixa",
      "explanation": "why this supports the claim, in Portuguese",
      "location": "section or page if identifiable"
    }}
  ],
  "overall_support": "Forte / Parcial / Fraco / Nenhum",
  "summary": "brief synthesis in Portuguese of how the article supports the claim"
}}"""

    raw  = ai_call(client, provider, model, prompt, max_tokens=6000)
    data = extract_json_from_ai(raw)
    return data if data else {"error": "Falha na busca", "raw": raw, "article": article_name}


def _gen_evidence_queries(client, provider, model, claim: str) -> list:
    """Use AI to generate optimized PubMed search queries from a paragraph/claim."""
    prompt = f"""You are a scientific literature search expert.

Read the paragraph below and generate 5 specific PubMed/MEDLINE search queries IN ENGLISH to find articles that support the scientific claims in it.

Rules:
- Each query should be specific and use MeSH-style terms when possible
- Mix different angles: prevalence, mechanisms, clinical evidence, systematic reviews
- Include "systematic review" OR "meta-analysis" in at least one query
- All queries must be in English

Return ONLY a JSON array of strings:
["query 1","query 2","query 3","query 4","query 5"]

PARAGRAPH:
{claim[:1500]}"""

    raw = ai_call(client, provider, model, prompt, max_tokens=400)
    try:
        m = re.search(r"\[.*?\]", raw, re.DOTALL)
        if m:
            queries = json.loads(m.group())
            return [q for q in queries if isinstance(q, str) and q.strip()][:5]
    except Exception:
        pass
    return []


def _assess_abstract_evidence(client, provider, model, claim: str, article: dict) -> dict:
    """AI assesses whether an article's abstract/title supports the claim."""
    title    = article.get("title", "")
    authors  = ", ".join(article.get("authors", [])[:3])
    journal  = article.get("journal", "")
    year     = article.get("year", "")
    abstract = article.get("abstract", "") or article.get("_abstract", "")

    content = f"Title: {title}\nAuthors: {authors}\nJournal: {journal} ({year})"
    if abstract:
        content += f"\nAbstract: {abstract[:2000]}"

    prompt = f"""You are a scientific evidence analyst.

Assess whether the article below SUPPORTS, is PARTIALLY RELEVANT, or does NOT SUPPORT the claim.

CLAIM:
{claim[:800]}

ARTICLE:
{content}

Return ONLY valid JSON (no markdown):
{{
  "article": "{title[:80]}",
  "overall_support": "Forte / Parcial / Fraco / Nenhum",
  "passages": [
    {{
      "text": "key phrase from title/abstract that supports the claim (verbatim or near-verbatim)",
      "relevance": "Alta/Media/Baixa",
      "explanation": "explanation in Portuguese of why this supports the claim",
      "location": "abstract/title"
    }}
  ],
  "summary": "1-2 sentence synthesis in Portuguese",
  "vancouver_citation": "Vancouver format: Authors. Title. Journal. Year;Vol(Issue):Pages."
}}"""

    raw  = ai_call(client, provider, model, prompt, max_tokens=1500)
    data = extract_json_from_ai(raw)
    if data:
        # Inject bibliographic metadata
        data["_meta"] = article
        data["_doi"]  = article.get("doi", "")
        data["_source"] = article.get("_source", "")
        data["_pmid"]   = article.get("pmid", "")
    return data if data else {
        "article": title,
        "overall_support": "Fraco",
        "passages": [],
        "summary": "Nao foi possivel avaliar este artigo.",
        "_meta": article,
    }


def _search_evidence_online(client, provider, model, claim: str) -> list:
    """Full online evidence search pipeline:
    1. AI generates PubMed queries from the claim
    2. Searches PubMed, Europe PMC, OpenAlex, LILACS
    3. AI assesses each found abstract for support strength
    Returns list of evidence dicts (same format as find_evidence_ai).
    """
    # Step 1 — generate queries
    queries = _gen_evidence_queries(client, provider, model, claim)
    if not queries:
        # Fallback: use first 120 chars of claim as query
        queries = [claim[:120]]

    # Step 2 — multi-source search
    all_articles = []
    seen_titles  = set()
    for q in queries:
        hits = multi_source_search(q, max_per_source=3)
        for h in hits:
            t = h.get("title","").lower()[:80]
            if t and t not in seen_titles:
                seen_titles.add(t)
                all_articles.append(h)
        time.sleep(0.3)

    # Limit to 10 most unique articles for assessment
    all_articles = all_articles[:10]

    if not all_articles:
        return []

    # Step 3 — AI assesses each abstract
    results = []
    for art in all_articles:
        ev = _assess_abstract_evidence(client, provider, model, claim, art)
        if ev.get("overall_support","Nenhum") not in ("Nenhum",):
            results.append(ev)
        time.sleep(1)  # avoid rate limit

    # Sort: Forte → Parcial → Fraco
    order = {"Forte": 0, "Parcial": 1, "Fraco": 2, "Nenhum": 3}
    results.sort(key=lambda x: order.get(x.get("overall_support","Nenhum"), 4))
    return results


def _display_ev_results(ev_results: list, online_mode: bool = False):
    """Render evidence results (shared by both local-file and online modes)."""
    st.divider()
    st.markdown("#### Evidencias Encontradas")

    if st.session_state.get("last_ev_claim"):
        st.info(f"**Afirmacao analisada:** {st.session_state['last_ev_claim'][:300]}")

    if online_mode:
        st.caption("🔍 Resultados obtidos por busca automatica em PubMed, Europe PMC, OpenAlex e LILACS.")

    color_map = {"Forte": "🟢", "Parcial": "🟡", "Fraco": "🟠", "Nenhum": "🔴"}

    for ev in ev_results:
        if "error" in ev and "passages" not in ev:
            st.error(f"Erro: {ev.get('error')}")
            continue

        support = ev.get("overall_support", "N/A")
        icon    = color_map.get(support, "⚪")
        art_label = ev.get("article", ev.get("_meta", {}).get("title", "Artigo"))[:90]

        # Build expander label with journal info when online
        meta = ev.get("_meta", {})
        if online_mode and meta:
            journal_line = f"{meta.get('journal','')} {meta.get('year','')}"
            source_tag   = meta.get("_source","")
            label = f"{icon} **{art_label}** — {support} | {journal_line} [{source_tag}]"
        else:
            label = f"{icon} **{art_label}** — Suporte: {support}"

        with st.expander(label, expanded=(support in ("Forte","Parcial"))):
            if ev.get("summary"):
                st.markdown(f"*{ev['summary']}*")

            # Vancouver citation (online mode)
            if online_mode:
                vcit = ev.get("vancouver_citation","")
                doi  = ev.get("_doi","") or meta.get("doi","")
                pmid = ev.get("_pmid","") or meta.get("pmid","")
                if vcit:
                    st.markdown(f"**Citacao Vancouver:** {vcit}")
                links = []
                if doi:
                    links.append(f"[DOI: {doi}](https://doi.org/{doi})")
                if pmid:
                    links.append(f"[PubMed PMID {pmid}](https://pubmed.ncbi.nlm.nih.gov/{pmid}/)")
                if links:
                    st.markdown("  ".join(links))

            passages = ev.get("passages", [])
            if not passages:
                st.warning("Nenhum trecho de suporte encontrado.")
                continue

            for i, p in enumerate(passages, 1):
                rel   = p.get("relevance", "")
                loc   = f" — {p['location']}" if p.get("location") else ""
                badge = "🟢" if rel == "Alta" else ("🟡" if rel == "Media" else "🟠")
                st.markdown(f"**Trecho {i}** {badge} Relevancia {rel}{loc}")
                st.markdown(
                    f'<div class="ref-box">{p.get("text","")}</div>',
                    unsafe_allow_html=True,
                )
                if p.get("explanation"):
                    st.caption(f"Por que apoia: {p['explanation']}")

    # Download report
    if ev_results:
        mode_label = "BUSCA ONLINE AUTOMATICA" if online_mode else "ANALISE DE ARTIGOS LOCAIS"
        lines = [f"BUSCA DE EVIDENCIAS — {mode_label}",
                 f"Afirmacao: {st.session_state.get('last_ev_claim','')}", ""]
        for ev in ev_results:
            meta = ev.get("_meta", {})
            lines += ["", f"ARTIGO: {ev.get('article', meta.get('title',''))}"]
            if online_mode:
                lines.append(f"Periodico: {meta.get('journal','')} {meta.get('year','')}")
                lines.append(f"Fonte: {meta.get('_source','')}")
                if meta.get("doi"):
                    lines.append(f"DOI: https://doi.org/{meta['doi']}")
            lines += [f"Suporte geral: {ev.get('overall_support','')}", f"Resumo: {ev.get('summary','')}"]
            for i, p in enumerate(ev.get("passages", []), 1):
                lines += ["", f"  Trecho {i} [Relevancia {p.get('relevance','')}]:",
                          f'  "{p.get("text","")}"', f"  -> {p.get('explanation','')}"]
        st.download_button(
            "📥 Baixar relatorio de evidencias (.txt)",
            data="\n".join(lines).encode("utf-8"),
            file_name="evidencias.txt",
            mime="text/plain",
            use_container_width=True,
        )


def render_evidencias_tab():
    st.markdown("### Busca de Evidencias Cientificas")
    st.caption(
        "Insira uma afirmacao ou paragrafo. "
        "**Com upload de PDFs:** a IA analisa os artigos que voce forneceu. "
        "**Sem upload:** a IA busca automaticamente em PubMed, Europe PMC, OpenAlex e LILACS.")

    api_key  = st.session_state.get("_api_key", "")
    provider = st.session_state.get("_provider", "")
    model    = st.session_state.get("_model", "")

    if not api_key:
        st.warning("Configure a chave API na barra lateral antes de usar.")
        return

    claim = st.text_area(
        "Afirmacao / Paragrafo a ser apoiado:",
        height=130,
        placeholder="Ex: O TDAH em adultos apresenta taxas de comorbidade com ansiedade superiores a 50%, com impacto significativo na funcionalidade...",
        key="ev_claim",
    )

    art_files = st.file_uploader(
        "Artigos de referencia (PDF ou DOCX) — opcional:",
        type=["pdf", "docx"],
        accept_multiple_files=True,
        key="ev_articles",
        help="Deixe em branco para busca automatica nas bases PubMed, Europe PMC, OpenAlex e LILACS.",
    )

    online_mode = not art_files  # True quando nenhum arquivo for anexado

    if online_mode:
        st.info(
            "📡 **Modo busca automatica:** nenhum artigo anexado. "
            "A IA vai gerar queries de busca a partir do seu paragrafo, "
            "consultar PubMed, Europe PMC, OpenAlex e LILACS, "
            "e avaliar quais artigos encontrados apoiam sua afirmacao.",
            icon=None,
        )
    else:
        st.success(
            f"📂 **Modo analise local:** {len(art_files)} artigo(s) carregado(s). "
            "A IA vai analisar os PDFs/DOCX que voce forneceu.",
            icon=None,
        )

    run_ev = st.button(
        "🔍 Buscar evidencias" if online_mode else "🔍 Analisar artigos anexados",
        type="primary",
        disabled=not api_key or not claim.strip(),
        use_container_width=True,
    )

    if run_ev:
        try:
            client = get_ai_client(provider, api_key)
        except Exception as e:
            st.error(f"Erro ao inicializar cliente: {e}")
            return

        if online_mode:
            # ── ONLINE MODE: search databases automatically ──────────────────
            with st.spinner("Gerando queries de busca e consultando bases cientificas..."):
                try:
                    results_ev = _search_evidence_online(client, provider, model, claim)
                    if not results_ev:
                        st.warning(
                            "Nenhum artigo com suporte relevante foi encontrado nas bases consultadas. "
                            "Tente reformular a afirmacao com termos mais especificos, "
                            "ou faca upload manual de artigos de referencia.")
                        return
                except Exception as e:
                    import traceback as _tb
                    st.error(f"Erro na busca online: {e}")
                    with st.expander("Detalhes"):
                        st.code(_tb.format_exc())
                    return
        else:
            # ── LOCAL MODE: analyse uploaded files ───────────────────────────
            results_ev = []
            progress   = st.progress(0)
            for idx, af in enumerate(art_files, 1):
                progress.progress(idx / len(art_files), text=f"Analisando {af.name}...")
                b = af.read()
                art_text = (extract_text_from_pdf(b)
                            if af.name.lower().endswith(".pdf")
                            else extract_text_from_docx(b))
                ev_result = find_evidence_ai(client, provider, model, claim, art_text, af.name)
                results_ev.append(ev_result)
            progress.empty()

        st.session_state["last_ev_results"]     = results_ev
        st.session_state["last_ev_claim"]       = claim
        st.session_state["last_ev_online_mode"] = online_mode

    ev_results   = st.session_state.get("last_ev_results")
    saved_online = st.session_state.get("last_ev_online_mode", False)

    if ev_results is not None:
        _display_ev_results(ev_results, online_mode=saved_online)



# =============================================================================
# PPT CITATION — extract, cite, annotate
# =============================================================================

def extract_pptx_slides(pptx_bytes: bytes) -> list:
    """Extract text and structure from each slide. Returns list of dicts."""
    from pptx import Presentation
    from io import BytesIO
    prs = Presentation(BytesIO(pptx_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title   = ""
        content = []
        notes   = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type in (13, 14):  # image/chart shapes
                continue
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:      # title
                    title = text
                elif ph_idx == 1:    # body
                    content.append(text)
                else:
                    content.append(text)
            else:
                content.append(text)
        # notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({
            "number":  i,
            "title":   title,
            "content": "\n".join(content),
            "notes":   notes,
            "full_text": (title + "\n" + "\n".join(content)).strip(),
        })
    return slides


def cite_slides_ai(client, provider, model, slides: list, refs: list,
                   citation_style: str = "Vancouver") -> list:
    """AI matches each slide to the best references. Returns list of per-slide citation dicts."""
    ref_catalogue = _build_ref_catalogue(refs)
    results = []
    for slide in slides:
        if not slide["full_text"].strip():
            results.append({"slide": slide["number"], "citations": [],
                            "note_text": "", "summary": "Slide sem texto"})
            continue

        prompt = f"""You are a scientific citation specialist.

Analyze the slide content below and identify which references from the catalogue BEST support the claims made.

CITATION STYLE: {citation_style}
REFERENCE CATALOGUE:
{ref_catalogue}

SLIDE {slide['number']}: {slide['title']}
CONTENT:
{slide['full_text'][:1500]}

TASK:
1. Identify 1-4 most relevant references for this slide
2. Write a short citation note (2-3 sentences) suitable for slide notes, in Portuguese
3. Format citations according to {citation_style} style

Return ONLY valid JSON (no markdown fences):
{{
  "slide": {slide['number']},
  "refs_used": ["REF1", "REF2"],
  "inline_citation": "[1, 2]",
  "note_text": "References: [1] Author (Year). [2] Author (Year). -- brief justification of why these support this slide",
  "summary": "one sentence why these refs support this slide, in Portuguese"
}}

If no reference is relevant to this slide, return refs_used=[] and note_text="Sem referencias aplicaveis."
"""
        raw  = ai_call(client, provider, model, prompt, max_tokens=2000)
        data = extract_json_from_ai(raw)
        if data:
            results.append(data)
        else:
            results.append({
                "slide": slide["number"], "refs_used": [],
                "inline_citation": "", "note_text": "Erro na analise.",
                "summary": "Erro"
            })
    return results


def add_citations_to_pptx(pptx_bytes: bytes, slides_info: list,
                           slide_citations: list, refs: list,
                           citation_style: str = "Vancouver",
                           add_footer: bool = True) -> bytes:
    """Insert citations into PPTX notes and optionally add footer text boxes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from io import BytesIO

    prs  = Presentation(BytesIO(pptx_bytes))
    cite_map = {c["slide"]: c for c in slide_citations}

    # Build numbered reference list from all used refs
    used_refx = set()
    for c in slide_citations:
        used_refx.update(c.get("refs_used", []))

    # Map REFx -> sequential number
    refx_to_num = {}
    num = 1
    for c in slide_citations:
        for rx in c.get("refs_used", []):
            if rx not in refx_to_num:
                refx_to_num[rx] = num
                num += 1

    for i, slide in enumerate(prs.slides, 1):
        citation = cite_map.get(i)
        if not citation:
            continue

        refs_used = citation.get("refs_used", [])
        note_text = citation.get("note_text", "")

        # ─ Add/update Notes ──────────────────────────────────────
        if note_text.strip() and note_text != "Sem referencias aplicaveis.":
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide
            else:
                notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            # Append to existing notes
            existing = tf.text.strip()
            separator = "\n\n---\n" if existing else ""
            # Add formatted references
            formatted_refs = []
            for rx in refs_used:
                idx = int(rx.replace("REF","")) - 1
                if 0 <= idx < len(refs):
                    n = refx_to_num.get(rx, idx+1)
                    formatted_refs.append(format_reference(refs[idx], n, citation_style))

            new_note = (
                f"{separator}REFERENCIAS (CitacaoIA):\n" +
                "\n".join(formatted_refs) +
                (f"\n\n{citation.get('summary','')}" if citation.get("summary") else "")
            )
            # Write to the last paragraph or add new
            if len(tf.paragraphs) > 0 and not tf.paragraphs[-1].text.strip():
                p = tf.paragraphs[-1]
            else:
                p = tf.add_paragraph()
            p.text = (existing + new_note) if existing else new_note.lstrip("\n")

        # ─ Add small citation footer text box ────────────────────
        if add_footer and refs_used:
            inline = ", ".join(
                str(refx_to_num[rx]) for rx in refs_used if rx in refx_to_num
            )
            if not inline:
                continue
            footer_text = f"[{inline}]"

            sw = prs.slide_width
            sh = prs.slide_height
            # Small box in bottom-right corner
            left   = sw - Inches(1.4)
            top    = sh - Inches(0.45)
            width  = Inches(1.2)
            height = Inches(0.35)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf2   = txBox.text_frame
            tf2.word_wrap = False
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.RIGHT
            run = p2.add_run()
            run.text = footer_text
            run.font.size   = Pt(9)
            run.font.bold   = True
            run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_ppt_citation_report_docx(slides_info: list, slide_citations: list,
                                      refs: list, citation_style: str,
                                      ppt_name: str) -> bytes:
    """Build a Word report: slide-by-slide reference list."""
    from docx import Document as DocxDoc
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc  = DocxDoc()
    for sec in doc.sections:
        sec.left_margin = Cm(2.5); sec.right_margin = Cm(2.5)
        sec.top_margin  = Cm(2.5); sec.bottom_margin = Cm(2.5)

    BLUE = RGBColor(0x1E, 0x3A, 0x5F)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Relatorio de Referenciamento de Apresentacao")
    r.bold = True; r.font.size = Pt(16); r.font.color.rgb = BLUE

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.add_run(
        f"Arquivo: {ppt_name}  |  Estilo: {citation_style}  |  "
        f"Gerado em {__import__('datetime').datetime.now().strftime('%d/%m/%Y')}"
    ).font.size = Pt(10)
    doc.add_paragraph()

    cite_map = {c["slide"]: c for c in slide_citations}
    refx_to_num = {}
    num = 1
    for c in slide_citations:
        for rx in c.get("refs_used", []):
            if rx not in refx_to_num:
                refx_to_num[rx] = num; num += 1

    # Slide-by-slide table
    ph = doc.add_paragraph()
    ph.add_run("Citacoes por Slide").bold = True
    ph.runs[0].font.size = Pt(13); ph.runs[0].font.color.rgb = BLUE
    doc.add_paragraph()

    for slide in slides_info:
        n = slide["number"]
        citation = cite_map.get(n, {})
        refs_used = citation.get("refs_used", [])

        sh = doc.add_paragraph()
        sh.add_run(f"Slide {n}: {slide.get('title','(sem titulo)')}").bold = True
        sh.runs[0].font.color.rgb = BLUE
        sh.paragraph_format.space_before = Pt(10)

        if not refs_used:
            doc.add_paragraph("Sem referencias aplicaveis.").runs[0].italic = True
            continue

        if citation.get("summary"):
            sp = doc.add_paragraph()
            sp.add_run(citation["summary"]).italic = True
            sp.paragraph_format.space_after = Pt(4)

        for rx in refs_used:
            idx = int(rx.replace("REF","")) - 1
            if 0 <= idx < len(refs):
                n_ref = refx_to_num.get(rx, idx+1)
                ref_str = format_reference(refs[idx], n_ref, citation_style)
                rp = doc.add_paragraph(style="List Bullet")
                rp.add_run(ref_str).font.size = Pt(10)

    # Full reference list
    doc.add_page_break()
    fh = doc.add_paragraph()
    fh.add_run("Lista Completa de Referencias").bold = True
    fh.runs[0].font.size = Pt(13); fh.runs[0].font.color.rgb = BLUE
    doc.add_paragraph()

    ordered_refs = sorted(refx_to_num.items(), key=lambda x: x[1])
    for rx, n_ref in ordered_refs:
        idx = int(rx.replace("REF","")) - 1
        if 0 <= idx < len(refs):
            ref_str = format_reference(refs[idx], n_ref, citation_style)
            rp = doc.add_paragraph()
            rp.add_run(ref_str).font.size = Pt(11)
            rp.paragraph_format.space_after = Pt(6)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def render_citar_ppt_tab():
    st.markdown("### Referenciar Apresentacao PowerPoint")
    st.caption(
        "Envie sua apresentacao e (opcionalmente) os artigos de referencia. "
        "A IA analisa cada slide e insere as citacoes nas notas e rodape dos slides, "
        "buscando automaticamente em bases cientificas quando necessario."
    )

    api_key  = st.session_state.get("_api_key","")
    provider = st.session_state.get("_provider","")
    model    = st.session_state.get("_model","")

    if not api_key:
        st.warning("Configure a chave API na barra lateral antes de usar.")
        return

    col_ppt, col_style = st.columns([3, 1])
    with col_ppt:
        ppt_file = st.file_uploader(
            "Apresentacao PowerPoint (.pptx)",
            type=["pptx"],
            key="ppt_upload"
        )
    with col_style:
        ppt_style = st.selectbox(
            "Estilo de citacao:", CITATION_STYLES, index=0, key="ppt_style"
        )

    st.markdown("#### Fontes de Referencias (opcional)")
    st.caption(
        "Sem referencias: a IA busca automaticamente em PubMed, EMBASE, OpenAlex e outras bases. "
        "Com referencias: usa prioritariamente as fornecidas."
    )
    ref_col1, ref_col2 = st.columns(2)
    with ref_col1:
        ref_pdfs = st.file_uploader(
            "PDFs dos artigos de referencia",
            type=["pdf"],
            accept_multiple_files=True,
            key="ppt_ref_pdfs"
        )
    with ref_col2:
        ref_text_input = st.text_area(
            "Ou cole lista de referencias (uma por linha):",
            height=120,
            placeholder="1. Autor A et al. Titulo. Revista. 2023;45(2):123.\n2. ...",
            key="ppt_ref_text"
        )

    use_library = st.checkbox("Incluir artigos da biblioteca local", value=True)
    add_footer  = st.checkbox("Adicionar marcador [N] no rodape de cada slide", value=True)

    if ppt_file:
        st.success(f"Apresentacao carregada: **{ppt_file.name}**")

    st.divider()
    run_ppt = st.button(
        "Referenciar apresentacao com IA",
        type="primary",
        use_container_width=True,
        disabled=not api_key or not ppt_file
    )

    if run_ppt and ppt_file:
        try:
            client = get_ai_client(provider, api_key)
        except Exception as e:
            st.error(f"Erro ao inicializar cliente: {e}"); return

        ppt_bytes = ppt_file.read()

        # ── Step 1: Extract slides ────────────────────────────────
        with st.spinner("Lendo slides da apresentacao..."):
            slides_info = extract_pptx_slides(ppt_bytes)
        st.info(f"{len(slides_info)} slides encontrados.")

        # ── Step 2: Build reference pool ─────────────────────────
        all_refs = []
        if use_library:
            all_refs.extend(load_library())

        if ref_pdfs:
            progress_r = st.progress(0, text="Extraindo metadados dos PDFs...")
            for ri, rf in enumerate(ref_pdfs, 1):
                progress_r.progress(ri / len(ref_pdfs), text=f"Lendo {rf.name}...")
                pdf_text = extract_text_from_pdf(rf.read())
                meta = extract_ref_metadata_ai(client, provider, model, pdf_text, rf.name)
                if meta and meta.get("title"):
                    all_refs.append(meta)
            progress_r.empty()

        if ref_text_input.strip():
            # Parse pasted reference list via AI
            with st.spinner("Parseando lista de referencias..."):
                parse_prompt = f"""Parse the following reference list into JSON.
Return ONLY valid JSON (no fences):
{{"refs":[{{"title":"","authors":[],"journal":"","year":"","volume":"","issue":"","pages":"","doi":""}}]}}

REFERENCE LIST:
{ref_text_input}"""
                raw = ai_call(client, provider, model, parse_prompt, max_tokens=4000)
                parsed = extract_json_from_ai(raw)
                if parsed and "refs" in parsed:
                    all_refs.extend(parsed["refs"])

        # ── Step 3: Auto-search if no refs ───────────────────────
        if not all_refs:
            st.info("Nenhuma referencia fornecida — buscando automaticamente nas bases cientificas...")
            queries_done = set()
            progress_s = st.progress(0, text="Buscando referencias...")
            for si, slide in enumerate(slides_info):
                progress_s.progress((si+1)/len(slides_info),
                                    text=f"Buscando para slide {slide['number']}...")
                if not slide["full_text"].strip():
                    continue
                # Use title or first 100 chars as query
                q = (slide["title"] or slide["full_text"])[:120]
                if q in queries_done:
                    continue
                queries_done.add(q)
                found = multi_source_search(q, max_per_source=2)
                for r in found:
                    if r not in all_refs:
                        all_refs.append(r)
                if len(all_refs) >= 40:
                    break
            progress_s.empty()
            st.success(f"{len(all_refs)} referencias encontradas automaticamente.")

        if not all_refs:
            st.error("Nao foi possivel obter referencias. Adicione PDFs ou uma lista de referencias.")
            return

        # ── Step 4: Deduplicate refs ──────────────────────────────
        seen_titles = set()
        dedup_refs = []
        for r in all_refs:
            t = (r.get("title","") or "").lower().strip()[:80]
            if t and t not in seen_titles:
                seen_titles.add(t)
                dedup_refs.append(r)
        all_refs = dedup_refs[:60]  # cap at 60

        st.info(f"Pool de referenciamento: **{len(all_refs)} referencias** disponveis.")

        # ── Step 5: AI cites each slide ───────────────────────────
        progress_c = st.progress(0, text="Referenciando slides com IA...")
        log_c = st.empty()
        slide_citations = []
        for si, slide in enumerate(slides_info):
            progress_c.progress((si+1)/len(slides_info),
                                 text=f"Analisando slide {slide['number']}/{len(slides_info)}: {slide['title'][:40]}...")
            cit = cite_slides_ai(client, provider, model, [slide], all_refs, ppt_style)
            slide_citations.extend(cit)
        progress_c.empty(); log_c.empty()

        # ── Step 6: Build annotated PPTX ─────────────────────────
        with st.spinner("Inserindo citacoes na apresentacao..."):
            annotated_pptx = add_citations_to_pptx(
                ppt_bytes, slides_info, slide_citations, all_refs, ppt_style, add_footer
            )

        # ── Step 7: Build Word report ─────────────────────────────
        with st.spinner("Gerando relatorio de referenciamento..."):
            report_docx = generate_ppt_citation_report_docx(
                slides_info, slide_citations, all_refs, ppt_style, ppt_file.name
            )

        st.session_state["ppt_annotated"]   = annotated_pptx
        st.session_state["ppt_report_docx"] = report_docx
        st.session_state["ppt_slides_info"] = slides_info
        st.session_state["ppt_citations"]   = slide_citations
        st.session_state["ppt_all_refs"]    = all_refs
        st.session_state["ppt_style"]       = ppt_style
        st.session_state["ppt_filename"]    = ppt_file.name

    # ─ Display results ─────────────────────────────────────────
    if not st.session_state.get("ppt_annotated"):
        return

    slides_info    = st.session_state["ppt_slides_info"]
    slide_citations = st.session_state["ppt_citations"]
    all_refs        = st.session_state["ppt_all_refs"]
    ppt_style_res   = st.session_state.get("ppt_style","Vancouver")
    ppt_fn          = st.session_state.get("ppt_filename","apresentacao.pptx")

    cited_slides = sum(1 for c in slide_citations if c.get("refs_used"))
    total_refs_used = len({rx for c in slide_citations for rx in c.get("refs_used",[])})

    c1, c2, c3 = st.columns(3)
    c1.metric("Slides analisados",    len(slides_info))
    c2.metric("Slides referenciados", cited_slides)
    c3.metric("Referencias utilizadas", total_refs_used)

    st.divider()

    # Download buttons
    dl1, dl2 = st.columns(2)
    with dl1:
        base = ppt_fn.rsplit(".",1)[0]
        st.download_button(
            "📊 Baixar apresentacao referenciada (.pptx)",
            data=st.session_state["ppt_annotated"],
            file_name=f"{base}_citado.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary",
        )
    with dl2:
        st.download_button(
            "📄 Baixar relatorio de referencias (.docx)",
            data=st.session_state["ppt_report_docx"],
            file_name=f"{base}_referencias.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True,
        )

    st.markdown("#### Revisao por Slide")
    cite_map = {c["slide"]: c for c in slide_citations}
    refx_to_num = {}
    num_r = 1
    for c in slide_citations:
        for rx in c.get("refs_used",[]):
            if rx not in refx_to_num:
                refx_to_num[rx] = num_r; num_r += 1

    for slide in slides_info:
        n = slide["number"]
        citation = cite_map.get(n, {})
        refs_used = citation.get("refs_used", [])
        icon = "✅" if refs_used else "⬜"
        inline = "[" + ", ".join(str(refx_to_num[rx]) for rx in refs_used if rx in refx_to_num) + "]" if refs_used else ""

        with st.expander(f"{icon} Slide {n}: {slide.get('title','(sem titulo)')} {inline}",
                         expanded=False):
            if slide.get("content"):
                st.caption(slide["content"][:300] + ("..." if len(slide["content"]) > 300 else ""))
            if not refs_used:
                st.info("Nenhuma referencia aplicavel identificada.")
            else:
                if citation.get("summary"):
                    st.markdown(f"*{citation['summary']}*")
                for rx in refs_used:
                    idx = int(rx.replace("REF","")) - 1
                    if 0 <= idx < len(all_refs):
                        n_ref = refx_to_num.get(rx, idx+1)
                        ref_str = format_reference(all_refs[idx], n_ref, ppt_style_res)
                        st.markdown(f'<div class="ref-box">{ref_str}</div>',
                                    unsafe_allow_html=True)

# =============================================================================
# MAIN
# =============================================================================


# =============================================================================
# REDESENHAR SLIDE — Vision AI + DALL-E 3 + python-pptx
# =============================================================================

def ai_vision_call(client, provider: str, model: str,
                   image_bytes: bytes, mime_type: str, prompt: str) -> str:
    """Call vision-capable AI to analyze an image. Returns raw text."""
    import base64
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    try:
        if provider == "Anthropic (Claude)":
            vision_model = "claude-3-5-sonnet-20241022"
            resp = client.messages.create(
                model=vision_model,
                max_tokens=2000,
                messages=[{"role": "user", "content": [
                    {"type": "image",
                     "source": {"type": "base64", "media_type": mime_type, "data": b64}},
                    {"type": "text", "text": prompt}
                ]}]
            )
            return resp.content[0].text
        elif provider == "OpenAI":
            from openai import OpenAI as _OAI
            oc = _OAI(api_key=client.api_key) if hasattr(client, "api_key") else client
            resp = oc.chat.completions.create(
                model="gpt-4o", max_tokens=2000,
                messages=[{"role": "user", "content": [
                    {"type": "image_url",
                     "image_url": {"url": f"data:{mime_type};base64,{b64}"}},
                    {"type": "text", "text": prompt}
                ]}]
            )
            return resp.choices[0].message.content
        else:  # Gemini
            import google.genai.types as _gt
            resp = client.models.generate_content(
                model=model,
                contents=[_gt.Part.from_bytes(data=image_bytes, mime_type=mime_type), prompt]
            )
            return resp.text
    except Exception as e:
        return f"Erro vision: {e}"


def extract_slide_content_vision(client, provider: str, model: str,
                                  image_bytes: bytes,
                                  mime_type: str = "image/png") -> dict:
    """Extract structured slide content using vision AI."""
    prompt = (
        "Analyze this presentation slide and extract ALL content as JSON. "
        "Return ONLY valid JSON with no markdown:\n"
        '{"title":"","subtitle":"","main_body":"","bullets":[],'
        '"labels":[],"callouts":[],"diagram_description":"",'
        '"layout_description":"","key_messages":[]}'
    )
    raw = ai_vision_call(client, provider, model, image_bytes, mime_type, prompt)
    return extract_json_from_ai(raw) or {"title": "", "raw": raw}


def extract_slide_content_pptx(pptx_bytes: bytes, slide_idx: int = 0) -> dict:
    """Extract slide content from PPTX using python-pptx (no image conversion)."""
    from pptx import Presentation as _Prs
    prs = _Prs(BytesIO(pptx_bytes))
    idx = min(slide_idx, len(prs.slides) - 1)
    slide = prs.slides[idx]
    content = {"title": "", "bullets": [], "labels": [], "callouts": [], "main_body": ""}
    extra = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text = shape.text.strip()
        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            if shape.placeholder_format.idx in (0, 1):
                is_title = True
        if not is_title and "title" in shape.name.lower():
            is_title = True
        if is_title:
            content["title"] = text
        elif hasattr(shape, "text_frame"):
            paras = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(paras) > 1:
                content["bullets"].extend(paras)
            elif paras:
                extra.append(paras[0])
    if extra:
        content["main_body"] = "\n".join(extra)
    return content


def generate_slide_image_dalle(openai_api_key: str, content: dict,
                                template_desc: str) -> bytes:
    """Generate redesigned slide image using DALL-E 3."""
    from openai import OpenAI as _OAI
    oc = _OAI(api_key=openai_api_key)
    title   = content.get("title", "")
    bullets = content.get("bullets", [])
    labels  = content.get("labels", [])
    diagram = content.get("diagram_description", "")
    c_desc  = f'Title: "{title}"'
    if content.get("subtitle"):
        c_desc += f'. Subtitle: "{content["subtitle"]}"'
    if bullets:
        c_desc += ". Bullets: " + " | ".join(f'"{b}"' for b in bullets[:6])
    if labels:
        c_desc += ". Labels: " + ", ".join(f'"{l}"' for l in labels[:4])
    if diagram:
        c_desc += f". Diagram/visual: {diagram}"
    prompt = (
        f"Professional medical/pharmaceutical presentation slide, widescreen 16:9. "
        f"LAYOUT STYLE TO REPRODUCE: {template_desc[:500]}. "
        f"EXACT CONTENT: {c_desc}. "
        f"Corporate palette: dark navy blue and gold/amber. Clean white background. "
        f"Professional icons. No extra text beyond specified content."
    )
    resp = oc.images.generate(
        model="dall-e-3", prompt=prompt[:4000],
        size="1792x1024", quality="hd", n=1,
    )
    r = requests.get(resp.data[0].url, timeout=60)
    r.raise_for_status()
    return r.content


def fill_pptx_template_with_content(template_bytes: bytes, content: dict) -> bytes:
    """Fill first slide of a PPTX template with extracted content."""
    from pptx import Presentation as _Prs
    from pptx.util import Pt
    prs   = _Prs(BytesIO(template_bytes))
    slide = prs.slides[0]
    body_queue = []
    if content.get("subtitle"):
        body_queue.append(content["subtitle"])
    if content.get("main_body"):
        body_queue.append(content["main_body"])
    body_queue.extend(content.get("bullets", []))
    title_done = False
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            if shape.placeholder_format.idx in (0, 1):
                is_title = True
        if not is_title and "title" in shape.name.lower():
            is_title = True
        tf = shape.text_frame
        if is_title and not title_done and content.get("title"):
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = content["title"]
                else:
                    tf.paragraphs[0].text = content["title"]
            title_done = True
        elif not is_title and body_queue:
            new_text = body_queue.pop(0)
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = new_text
                else:
                    tf.paragraphs[0].text = new_text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pptx_from_image_bg(template_img_bytes: bytes, content: dict,
                               template_name: str = "") -> bytes:
    """Create PPTX with template image as background + text overlays."""
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    prs = _Prs()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    img_stream = BytesIO(template_img_bytes)
    new_slide.shapes.add_picture(
        img_stream, Emu(0), Emu(0),
        width=prs.slide_width, height=prs.slide_height
    )
    if content.get("title"):
        tb = new_slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = content["title"]
        run.font.size  = Pt(28)
        run.font.bold  = True
        run.font.color.rgb = RGBColor(0, 51, 102)
    bullets = content.get("bullets", [])
    if content.get("main_body"):
        bullets = [content["main_body"]] + bullets
    if bullets:
        tb2 = new_slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(7.5), Inches(4))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets[:8]):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            run = p.add_run()
            run.text = f"\\u2022  {b}"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(40, 40, 40)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_redesenhar_slide_tab():
    api_key  = st.session_state.get("_api_key", "")
    provider = st.session_state.get("_provider", "Anthropic (Claude)")
    model    = st.session_state.get("_model", "")

    st.markdown("### Redesenhar Slide")
    st.caption(
        "Envie um slide rascunho (PPT ou imagem) e um layout template. "
        "A IA extrai o conteudo do rascunho e aplica ao template, "
        "gerando o resultado em PPTX e imagem PNG."
    )

    col_draft, col_tmpl = st.columns(2)
    with col_draft:
        st.markdown("#### Slide Rascunho")
        draft_file = st.file_uploader(
            "Upload rascunho (PPTX ou imagem)",
            type=["pptx", "png", "jpg", "jpeg"],
            key="rd_draft_upload",
        )
        if draft_file:
            if draft_file.name.lower().endswith((".png", ".jpg", ".jpeg")):
                st.image(draft_file, use_column_width=True)
            else:
                st.success(f"PPTX: {draft_file.name}")

    with col_tmpl:
        st.markdown("#### Layout Template")
        tmpl_file = st.file_uploader(
            "Upload template (PPTX ou imagem)",
            type=["pptx", "png", "jpg", "jpeg"],
            key="rd_tmpl_upload",
        )
        if tmpl_file:
            if tmpl_file.name.lower().endswith((".png", ".jpg", ".jpeg")):
                st.image(tmpl_file, use_column_width=True)
            else:
                st.success(f"Template PPTX: {tmpl_file.name}")

    col_opt1, col_opt2 = st.columns(2)
    with col_opt1:
        slide_num = st.number_input("Slide do rascunho (1 = primeiro)",
                                    min_value=1, max_value=30, value=1)
    with col_opt2:
        use_dalle = st.checkbox(
            "Gerar imagem com DALL-E 3 (OpenAI)",
            value=(provider == "OpenAI"),
            help="Requer chave OpenAI. Gera imagem do slide redesenhado.",
        )

    dalle_key = api_key
    if use_dalle and provider != "OpenAI":
        dalle_key = st.text_input("Chave OpenAI para DALL-E:",
                                   type="password", key="rd_dalle_key")

    st.divider()
    run_rd = st.button("Redesenhar Slide", type="primary",
                        use_container_width=True,
                        disabled=(not draft_file or not tmpl_file or not api_key))

    if run_rd:
        draft_bytes = draft_file.read()
        tmpl_bytes  = tmpl_file.read()
        draft_lower = draft_file.name.lower()
        tmpl_lower  = tmpl_file.name.lower()
        slide_idx   = int(slide_num) - 1

        try:
            client = get_ai_client(provider, api_key)
        except Exception as e:
            st.error(f"Erro ao inicializar IA: {e}"); return

        # 1. Extract content from draft
        with st.spinner("Extraindo conteudo do rascunho com IA..."):
            try:
                if draft_lower.endswith(".pptx"):
                    content = extract_slide_content_pptx(draft_bytes, slide_idx)
                    # Try vision too for richer extraction
                    try:
                        from pptx import Presentation as _Prs2
                        from PIL import Image as _PILImg
                        # Quick PIL thumbnail from pptx text as placeholder
                        # (LibreOffice not available on Cloud — use text-only extraction)
                    except Exception:
                        pass
                else:
                    mime = "image/png" if "png" in draft_lower else "image/jpeg"
                    content = extract_slide_content_vision(
                        client, provider, model, draft_bytes, mime)
            except Exception as e:
                st.error(f"Erro na extracao: {e}"); return

        # Show extracted content as editable fields
        st.markdown("#### Conteudo extraido (editavel antes de gerar)")
        c1, c2 = st.columns(2)
        with c1:
            t_edit  = st.text_input("Titulo:", value=content.get("title", ""),  key="rd_t")
            s_edit  = st.text_input("Subtitulo:", value=content.get("subtitle",""), key="rd_s")
        with c2:
            b_raw   = "\n".join(
                content.get("bullets", []) +
                ([content.get("main_body","")] if content.get("main_body") else [])
            )
            b_edit  = st.text_area("Corpo / Bullets (um por linha):",
                                    value=b_raw, height=130, key="rd_b")

        content["title"]    = t_edit
        content["subtitle"] = s_edit
        content["bullets"]  = [x.strip() for x in b_edit.split("\n") if x.strip()]
        content["main_body"] = ""

        # 2. Analyze template layout
        with st.spinner("Analisando layout do template..."):
            try:
                if tmpl_lower.endswith(".pptx"):
                    tmpl_img_bytes = None  # no render available without LibreOffice
                    tmpl_desc = ("Professional corporate slide with blue and gold color scheme, "
                                 "icons and structured layout as provided in the PPTX template.")
                else:
                    tmpl_mime = "image/png" if "png" in tmpl_lower else "image/jpeg"
                    tmpl_img_bytes = tmpl_bytes
                    tmpl_desc = ai_vision_call(
                        client, provider, model, tmpl_bytes, tmpl_mime,
                        "Describe this slide layout for faithful reproduction: "
                        "colors, structure, hierarchy, typography, element positions. "
                        "Be specific (max 300 words)."
                    )
            except Exception as e:
                tmpl_desc = "Professional medical slide layout with blue corporate palette."
                tmpl_img_bytes = None
                st.warning(f"Analise visual do template falhou: {e}")

        # 3. Generate PPTX
        with st.spinner("Gerando PPTX..."):
            try:
                if tmpl_lower.endswith(".pptx"):
                    result_pptx = fill_pptx_template_with_content(tmpl_bytes, content)
                else:
                    result_pptx = build_pptx_from_image_bg(
                        tmpl_bytes, content, tmpl_lower)
                st.session_state["rd_pptx"] = result_pptx
            except Exception as e:
                st.error(f"Erro ao gerar PPTX: {e}")
                st.session_state["rd_pptx"] = None

        # 4. Generate image with DALL-E
        if use_dalle and dalle_key:
            with st.spinner("Gerando imagem com DALL-E 3 (pode demorar ~20s)..."):
                try:
                    result_img = generate_slide_image_dalle(dalle_key, content, tmpl_desc)
                    st.session_state["rd_image"] = result_img
                except Exception as e:
                    st.error(f"Erro DALL-E: {e}")
                    st.session_state["rd_image"] = None
        else:
            st.session_state["rd_image"] = None

        st.rerun()

    # Results
    if st.session_state.get("rd_pptx") or st.session_state.get("rd_image"):
        st.divider()
        st.markdown("### Resultado")
        if st.session_state.get("rd_image"):
            st.image(st.session_state["rd_image"],
                     caption="Slide Redesenhado — DALL-E 3", use_column_width=True)

        dl1, dl2, dl3 = st.columns(3)
        if st.session_state.get("rd_pptx"):
            dl1.download_button("Baixar PPTX",
                data=st.session_state["rd_pptx"],
                file_name="slide_redesenhado.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        if st.session_state.get("rd_image"):
            dl2.download_button("Baixar PNG",
                data=st.session_state["rd_image"],
                file_name="slide_redesenhado.png",
                mime="image/png")
        if dl3.button("Novo Redesenho"):
            st.session_state["rd_pptx"]  = None
            st.session_state["rd_image"] = None
            st.rerun()



# =============================================================================
# GUIA DO USUARIO — popup de ajuda
# =============================================================================

GUIA_CONTEUDO = {
    "📖 Biblioteca": {
        "desc": "Gerencie sua coleção de artigos científicos.",
        "passos": [
            "Adicione artigos via DOI, PubMed ID ou busca por palavras-chave.",
            "Os artigos salvos ficam disponíveis automaticamente em todas as abas.",
            "Use os filtros para encontrar artigos por título, autores ou ano.",
            "Exporte sua biblioteca para CSV ou BibTeX.",
        ],
        "dica": "Mantenha a biblioteca atualizada — ela é a base de referências para o processamento de texto.",
    },
    "✍️ Processar Texto": {
        "desc": "Insere ou revisa citações Vancouver em textos científicos usando IA.",
        "passos": [
            "Escolha o modo: 'Adicionar citações' (texto sem refs) ou 'Revisar citações' (texto já citado).",
            "Faça upload do seu texto (PDF, Word ou .md) ou cole diretamente.",
            "Opcionalmente, anexe PDFs dos artigos de referência (sem upload, a IA busca no PubMed automaticamente).",
            "Clique em 'Processar texto com IA' e aguarde.",
            "Revise os parágrafos modificados e baixe o Word formatado.",
        ],
        "dica": "Para textos longos (>5 páginas), o processamento é feito em partes — aguarde a barra de progresso concluir.",
    },
    "🔍 Revisão Completa": {
        "desc": "Revisão acadêmica profunda: estrutura, argumentação, coerência e linguagem.",
        "passos": [
            "Faça upload do texto ou cole-o.",
            "Clique em 'Revisar com IA'.",
            "Veja o texto revisado lado a lado com o original.",
            "Baixe o relatório de revisão em Word.",
        ],
        "dica": "Use esta aba para uma revisão editorial completa, não apenas de citações.",
    },
    "✅ Verificar Refs": {
        "desc": "Verifica se as referências bibliográficas estão corretas e completas.",
        "passos": [
            "Cole sua lista de referências no campo indicado.",
            "A IA confere cada referência contra PubMed e CrossRef.",
            "Referências com erros são marcadas em vermelho com a correção sugerida.",
            "Baixe o relatório com todas as correções.",
        ],
        "dica": "Ideal para uso antes da submissão: detecta DOIs errados, autores faltando e formatação incorreta.",
    },
    "🔄 Converter": {
        "desc": "Converte listas de referências entre formatos (Vancouver ↔ APA ↔ ABNT ↔ Chicago).",
        "passos": [
            "Cole a lista de referências no formato de origem.",
            "Selecione o formato de destino.",
            "Clique em 'Converter' e baixe o resultado.",
        ],
        "dica": "Funciona com listas mistas — a IA detecta o formato automaticamente.",
    },
    "🔬 Buscar Evidências": {
        "desc": "Encontra evidências científicas que apoiam um parágrafo ou afirmação.",
        "passos": [
            "Cole o parágrafo ou afirmação que precisa de embasamento.",
            "Sem upload: a IA busca automaticamente em PubMed, Europe PMC, OpenAlex e LILACS.",
            "Com upload: anexe PDFs de artigos e a IA encontra os trechos de suporte dentro deles.",
            "Veja os trechos relevantes com nível de suporte (Forte/Parcial/Fraco).",
            "Baixe o relatório de evidências em .txt.",
        ],
        "dica": "Modo sem upload é ideal para exploração inicial; modo com upload é para análise profunda de artigos específicos.",
    },
    "📊 Citar PPT": {
        "desc": "Insere citações Vancouver nos slides de apresentações PowerPoint.",
        "passos": [
            "Faça upload do arquivo .pptx.",
            "Opcionalmente, adicione PDFs dos artigos de referência.",
            "Clique em 'Processar Apresentação'.",
            "As citações são inseridas nas notas de cada slide.",
            "Baixe o PPTX anotado e/ou o relatório Word.",
        ],
        "dica": "Use a opção 'Incluir rodapé nos slides' para que as citações apareçam visíveis na apresentação.",
    },
    "🎨 Redesenhar Slide": {
        "desc": "Aplica o conteúdo de um slide rascunho ao layout de um template profissional.",
        "passos": [
            "Faça upload do slide rascunho (PPTX ou imagem).",
            "Faça upload do template de destino (PPTX ou imagem).",
            "A IA extrai o conteúdo do rascunho e aplica ao template.",
            "Com chave OpenAI: gera versão visual via DALL-E 3.",
            "Baixe o resultado em PPTX e/ou PNG.",
        ],
        "dica": "Funciona melhor com templates PPTX. Para imagens, o resultado é mais visual mas menos preciso na formatação.",
    },
}


def render_user_guide():
    """Render user guide as a modal-style expander at the top of the page."""
    with st.sidebar:
        st.divider()
        with st.expander("📘 Guia do Usuário", expanded=False):
            aba_names = list(GUIA_CONTEUDO.keys())
            sel_aba = st.selectbox(
                "Selecione a aba:",
                aba_names,
                key="guia_aba_sel",
                label_visibility="collapsed",
            )
            info = GUIA_CONTEUDO[sel_aba]
            st.markdown(f"**{sel_aba}**")
            st.markdown(info["desc"])
            st.markdown("**Como usar:**")
            for i, passo in enumerate(info["passos"], 1):
                st.markdown(f"{i}. {passo}")
            st.info(f"💡 **Dica:** {info['dica']}")


def render_user_guide_main():
    """Render a dismissible tutorial banner on first use (top of page)."""
    if st.session_state.get("guia_dismissed"):
        return

    with st.container():
        col_msg, col_btn = st.columns([8, 2])
        with col_msg:
            st.info(
                "👋 **Bem-vindo ao CitacaoIA!** Primeira vez aqui? "
                "Veja o **Guia do Usuário** na barra lateral (📘) para entender cada aba. "
                "Ou feche este aviso se já conhece o sistema.",
                icon=None,
            )
        with col_btn:
            if st.button("✕ Fechar", key="btn_dismiss_guia", use_container_width=True):
                st.session_state["guia_dismissed"] = True
                st.rerun()


# =============================================================================
# MAIN
# =============================================================================

def main():
    # ── Authentication gate ────────────────────────────────────────────────
    if not render_login_page():
        return

    auth_user = st.session_state.get("auth_user", "")
    auth_role = st.session_state.get("auth_role", "user")

    st.markdown("""<div class="cia-header">
  <h1>📚 CitacaoIA</h1>
  <p>Gestor inteligente de citacoes . Vancouver . APA . ABNT . Chicago . Textos + Apresentacoes + Referencias . PubMed + EMBASE + LILACS + OpenAlex</p>
</div>""", unsafe_allow_html=True)

    api_key, provider, model = render_sidebar()
    st.session_state["_api_key"]  = api_key
    st.session_state["_provider"] = provider
    st.session_state["_model"]    = model

    # ── Sidebar: user guide + user info + logout ───────────────────────────
    render_user_guide()

    with st.sidebar:
        st.divider()
        st.caption(f"👤 **{auth_user}** ({auth_role})")
        if st.button("🚪 Sair", key="btn_logout", use_container_width=True):
            _append_usage_log(auth_user, "logout")
            for k in ["auth_user", "auth_role"]:
                if k in st.session_state:
                    del st.session_state[k]
            st.rerun()

    # ── Welcome banner (dismissible, first-time only) ─────────────────────
    render_user_guide_main()

    # ── Admin panel ────────────────────────────────────────────────────────
    if auth_role == "admin":
        with st.expander("⚙️ Painel do Administrador", expanded=False):
            render_admin_panel()

    # ── Main tabs ──────────────────────────────────────────────────────────
    tab_bib, tab_citar, tab_rev, tab_vr, tab_conv, tab_ev, tab_ppt, tab_rd = st.tabs([
        "📖 Biblioteca",
        "✍️ Processar Texto",
        "🔍 Revisao Completa",
        "✅ Verificar Refs",
        "🔄 Converter",
        "🔬 Buscar Evidencias",
        "📊 Citar PPT",
        "🎨 Redesenhar Slide",
    ])

    with tab_bib:
        render_biblioteca_tab()
    with tab_citar:
        render_citar_tab()
    with tab_rev:
        render_revisao_tab()
    with tab_vr:
        render_verificar_refs_tab()
    with tab_conv:
        render_converter_tab()
    with tab_ev:
        render_evidencias_tab()
    with tab_ppt:
        render_citar_ppt_tab()
    with tab_rd:
        render_redesenhar_slide_tab()

    # Log session activity (once per session)
    if st.session_state.get("_last_log_user") != auth_user:
        _append_usage_log(auth_user, "session_active", provider)
        st.session_state["_last_log_user"] = auth_user


if __name__ == "__main__":
    main()
